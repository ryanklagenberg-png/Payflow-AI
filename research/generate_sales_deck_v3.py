from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_LABEL_POSITION
from pptx.chart.data import CategoryChartData

NAVY = RGBColor(0x1a, 0x36, 0x5d)
BLUE = RGBColor(0x2b, 0x6c, 0xb0)
LIGHT_BLUE = RGBColor(0xeb, 0xf4, 0xff)
GREEN = RGBColor(0x38, 0xa1, 0x69)
LIGHT_GREEN = RGBColor(0xc6, 0xf6, 0xd5)
DARK_GREEN = RGBColor(0x27, 0x67, 0x49)
RED = RGBColor(0xe5, 0x3e, 0x3e)
LIGHT_RED = RGBColor(0xfe, 0xd7, 0xd7)
WHITE = RGBColor(0xff, 0xff, 0xff)
DARK = RGBColor(0x1a, 0x20, 0x2c)
GRAY = RGBColor(0x4a, 0x55, 0x68)
LIGHT_GRAY = RGBColor(0xf7, 0xfa, 0xfc)
MEDIUM_GRAY = RGBColor(0xe2, 0xe8, 0xf0)
ORANGE = RGBColor(0xdd, 0x6b, 0x20)
GOLD = RGBColor(0xd6, 0x9e, 0x2e)
PURPLE = RGBColor(0x6b, 0x46, 0xc1)
LIGHT_PURPLE = RGBColor(0xe9, 0xd8, 0xfd)
TEAL = RGBColor(0x31, 0x90, 0x95)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

def add_bg(s, c):
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = c

def shape(s, l, t, w, h, fc, bc=None):
    sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    sh.fill.solid(); sh.fill.fore_color.rgb = fc
    if bc: sh.line.color.rgb = bc; sh.line.width = Pt(1)
    else: sh.line.fill.background()
    return sh

def rect(s, l, t, w, h, fc):
    sh = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    sh.fill.solid(); sh.fill.fore_color.rgb = fc; sh.line.fill.background()
    return sh

def txt(s, l, t, w, h, tx, sz=18, c=DARK, b=False, a=PP_ALIGN.LEFT):
    bx = s.shapes.add_textbox(l, t, w, h)
    tf = bx.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = tx; p.font.size = Pt(sz)
    p.font.color.rgb = c; p.font.bold = b; p.font.name = "Calibri"; p.alignment = a
    return bx

def bul(s, l, t, w, h, items, sz=15, c=DARK, sp=8):
    bx = s.shapes.add_textbox(l, t, w, h)
    tf = bx.text_frame; tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item; p.font.size = Pt(sz); p.font.color.rgb = c
        p.font.name = "Calibri"; p.space_after = Pt(sp)

def multiline(s, l, t, w, h, lines, sz=13, c=GRAY, a=PP_ALIGN.LEFT, sp=2):
    bx = s.shapes.add_textbox(l, t, w, h)
    tf = bx.text_frame; tf.word_wrap = True
    for ln in lines:
        p = tf.add_paragraph(); p.text = ln; p.font.size = Pt(sz)
        p.font.color.rgb = c; p.font.name = "Calibri"; p.alignment = a; p.space_after = Pt(sp)

def stat(s, l, t, w, h, label, val, sub="", vc=NAVY):
    shape(s, l, t, w, h, WHITE, MEDIUM_GRAY)
    txt(s, l+Inches(0.15), t+Inches(0.1), w-Inches(0.3), Inches(0.25), label, sz=10, c=GRAY, b=True)
    txt(s, l+Inches(0.15), t+Inches(0.35), w-Inches(0.3), Inches(0.5), val, sz=30, c=vc, b=True, a=PP_ALIGN.CENTER)
    if sub: txt(s, l+Inches(0.15), t+Inches(0.85), w-Inches(0.3), Inches(0.25), sub, sz=9, c=GRAY, a=PP_ALIGN.CENTER)

def hdr(s, title, sub=""):
    rect(s, Inches(0), Inches(0), Inches(13.333), Inches(0.08), NAVY)
    txt(s, Inches(0.8), Inches(0.3), Inches(11), Inches(0.6), title, sz=32, c=NAVY, b=True)
    if sub: txt(s, Inches(0.8), Inches(0.9), Inches(11), Inches(0.4), sub, sz=17, c=GRAY)

def card(s, l, t, w, h, title, lines, color, title_sz=17):
    shape(s, l, t, w, h, WHITE, MEDIUM_GRAY)
    rect(s, l+Inches(0.05), t+Inches(0.05), w-Inches(0.1), Inches(0.08), color)
    txt(s, l+Inches(0.2), t+Inches(0.2), w-Inches(0.4), Inches(0.35), title, sz=title_sz, c=NAVY, b=True)
    multiline(s, l+Inches(0.2), t+Inches(0.6), w-Inches(0.4), h-Inches(0.7), lines, c=GRAY)

def circle_step(s, l, t, num, color):
    c = s.shapes.add_shape(MSO_SHAPE.OVAL, l, t, Inches(0.8), Inches(0.8))
    c.fill.solid(); c.fill.fore_color.rgb = color; c.line.fill.background()
    tf = c.text_frame; p = tf.paragraphs[0]; p.text = str(num)
    p.font.size = Pt(24); p.font.color.rgb = WHITE; p.font.bold = True; p.alignment = PP_ALIGN.CENTER

def arrow(s, l, t):
    a = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, l, t, Inches(0.4), Inches(0.35))
    a.fill.solid(); a.fill.fore_color.rgb = MEDIUM_GRAY; a.line.fill.background()


# ============================================================
# SLIDE 1: TITLE
# ============================================================
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(sl, NAVY)
rect(sl, Inches(0), Inches(3.6), Inches(13.333), Inches(0.06), BLUE)
txt(sl, Inches(1), Inches(1.0), Inches(11), Inches(1), "PayFlow AI", sz=58, c=WHITE, b=True)
txt(sl, Inches(1), Inches(2.1), Inches(11), Inches(0.8),
    "AI-Powered Accounts Payable Automation for Construction", sz=26, c=RGBColor(0xa0,0xce,0xf0))
txt(sl, Inches(1), Inches(4.0), Inches(11), Inches(0.5),
    "Stop spending $576,000/year processing invoices manually.", sz=20, c=RGBColor(0xcb,0xd5,0xe0))
txt(sl, Inches(1), Inches(4.5), Inches(11), Inches(0.5),
    "Start capturing $100,000+ in early-pay discounts you're leaving on the table.", sz=20, c=RGBColor(0xcb,0xd5,0xe0))
txt(sl, Inches(1), Inches(6.0), Inches(11), Inches(0.4),
    "Confidential  |  2026", sz=14, c=RGBColor(0x71,0x80,0x96))

# ============================================================
# SLIDE 2: THE PROBLEM
# ============================================================
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(sl, WHITE)
hdr(sl, "Your AP Department By The Numbers", "Based on 3,000 invoices per month")

stat(sl, Inches(0.8), Inches(1.5), Inches(2.5), Inches(1.2), "INVOICES / MONTH", "3,000", "", NAVY)
stat(sl, Inches(3.6), Inches(1.5), Inches(2.5), Inches(1.2), "COST PER INVOICE", "$16-$25", "Construction average", RED)
stat(sl, Inches(6.4), Inches(1.5), Inches(2.5), Inches(1.2), "MONTHLY AP COST", "$48,000", "", RED)
stat(sl, Inches(9.2), Inches(1.5), Inches(2.5), Inches(1.2), "ANNUAL AP COST", "$576,000", "", RED)

# Pie chart
cd = CategoryChartData()
cd.categories = ["AP Staff Labor\n$360K", "Late Payments &\nMissed Discounts\n$90K", "Error Correction\n$72K", "Overhead\n$36K", "Vendor\nFollow-ups\n$18K"]
cd.add_series("Cost", (360, 90, 72, 36, 18))
ch = sl.shapes.add_chart(XL_CHART_TYPE.PIE, Inches(0.8), Inches(3.0), Inches(5.5), Inches(4.0), cd).chart
ch.has_legend = True; ch.legend.position = XL_LEGEND_POSITION.BOTTOM
ch.legend.font.size = Pt(10); ch.legend.font.color.rgb = GRAY
plot = ch.plots[0]; plot.has_data_labels = True
plot.data_labels.font.size = Pt(11); plot.data_labels.font.bold = True
plot.data_labels.font.color.rgb = WHITE; plot.data_labels.number_format = '0"%"'
for i, c in enumerate([NAVY, ORANGE, RED, BLUE, GRAY]):
    ch.series[0].points[i].format.fill.solid()
    ch.series[0].points[i].format.fill.fore_color.rgb = c

txt(sl, Inches(6.8), Inches(3.0), Inches(5.5), Inches(0.4), "Where Your $576,000 Goes:", sz=20, c=NAVY, b=True)
bul(sl, Inches(7.0), Inches(3.5), Inches(5.5), Inches(2.5), [
    "62% ($360K) - Staff manually typing invoice data into Sage",
    "16% ($90K) - Late payment penalties & missed early-pay discounts",
    "13% ($72K) - Finding and fixing data entry errors",
    "6% ($36K) - Office space, software, printing, filing",
    "3% ($18K) - Calling vendors about payment status",
], sz=14, c=GRAY)
txt(sl, Inches(6.8), Inches(6.0), Inches(5.5), Inches(0.5),
    "And you're leaving $50-100K in early-pay discounts uncaptured.", sz=14, c=RED, b=True)

# ============================================================
# SLIDE 3: SOLUTION OVERVIEW
# ============================================================
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(sl, WHITE)
hdr(sl, "PayFlow AI: Not Just AP Automation", "It's a complete financial intelligence platform for construction")

features = [
    ("AI Invoice Extraction", ["Reads any format: PDF, scan, photo, email", "99% accuracy, no templates", "5 seconds per invoice"], BLUE),
    ("5D Job Costing", ["Predicts GL, job, phase, cost code, type", "Learns your patterns over time", "Construction-native coding"], BLUE),
    ("Three-Way PO Matching", ["Matches against Procore POs", "Flags price & quantity variances", "Tracks change order impact"], BLUE),
    ("AIA & Lien Waivers", ["Reads G702/G703 natively", "Classifies 4 waiver types", "Blocks payment if missing"], GREEN),
    ("Fraud & Anomaly Detection", ["Duplicate detection", "Overbilling beyond contract", "Price spike alerts"], ORANGE),
    ("Cash Flow Forecasting", ["Predicts AP obligations 30-60 days out", "Models payment scenarios", "Helps plan draws and borrowing"], PURPLE),
    ("Early-Pay Discount Capture", ["Flags every 2/10 Net 30 opportunity", "Prioritizes in approval queue", "Target: $50-100K/year captured"], TEAL),
    ("Job Profitability Alerts", ["Real-time cost vs budget by job", "Alerts when costs exceed estimate", "Catches bleeding jobs early"], ORANGE),
    ("Vendor Spend Analytics", ["Top vendors by spend with trends", "Price trend analysis across trades", "Consolidation opportunities"], NAVY),
]

for i, (title, lines, color) in enumerate(features):
    col = i % 3; row = i // 3
    l = Inches(0.8) + Inches(col * 4.05)
    t = Inches(1.5) + Inches(row * 1.85)
    card(sl, l, t, Inches(3.7), Inches(1.6), title, lines, color, title_sz=15)

# ============================================================
# SLIDE 4: HOW IT WORKS
# ============================================================
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(sl, WHITE)
hdr(sl, "How It Works: 5 Seconds, Not 5 Days")

steps = [
    ("1", "Invoice\nArrives", "Email, upload,\nscan, portal", BLUE),
    ("2", "AI Reads\nDocument", "Extracts every\nfield in seconds", BLUE),
    ("3", "AI Codes\nto Job", "GL, job, phase,\ncost code, type", BLUE),
    ("4", "AI Validates\n& Matches", "PO match, fraud,\ncompliance check", GREEN),
    ("5", "Route or\nAuto-Approve", "85%+ straight\nthrough", GREEN),
    ("6", "Posts to\nSage", "Syncs to your\nERP directly", NAVY),
]
for i, (num, title, desc, color) in enumerate(steps):
    l = Inches(0.4) + Inches(i * 2.15)
    circle_step(sl, l+Inches(0.5), Inches(1.3), num, color)
    if i < 5: arrow(sl, l+Inches(1.45), Inches(1.5))
    txt(sl, l, Inches(2.3), Inches(2.0), Inches(0.6), title, sz=13, c=NAVY, b=True, a=PP_ALIGN.CENTER)
    txt(sl, l, Inches(2.9), Inches(2.0), Inches(0.5), desc, sz=11, c=GRAY, a=PP_ALIGN.CENTER)

# Before/After
rect(sl, Inches(0.8), Inches(3.7), Inches(5.6), Inches(3.3), LIGHT_RED)
txt(sl, Inches(1.0), Inches(3.8), Inches(5.2), Inches(0.4), "TODAY (Manual)", sz=18, c=RED, b=True, a=PP_ALIGN.CENTER)
bul(sl, Inches(1.2), Inches(4.3), Inches(5), Inches(2.5), [
    "AP clerk opens email, downloads invoice",
    "Manually types vendor, amount, date into Sage",
    "Looks up job number and cost code",
    "Emails PM for approval, waits days",
    "PM eventually approves, AP posts to Sage",
    "Misses early-pay discount (invoice sat too long)",
    "Total: 20-30 min/invoice, 14+ days to pay",
], sz=12, c=RGBColor(0xc5,0x30,0x30))

rect(sl, Inches(6.9), Inches(3.7), Inches(5.6), Inches(3.3), LIGHT_GREEN)
txt(sl, Inches(7.1), Inches(3.8), Inches(5.2), Inches(0.4), "WITH PAYFLOW AI", sz=18, c=GREEN, b=True, a=PP_ALIGN.CENTER)
bul(sl, Inches(7.3), Inches(4.3), Inches(5), Inches(2.5), [
    "Invoice arrives via email or upload",
    "AI reads and extracts all fields in 5 seconds",
    "AI codes to correct GL, job, phase, cost code",
    "AI matches against PO, checks compliance",
    "Auto-approves or routes to PM instantly",
    "Flags early-pay discount, prioritizes approval",
    "Total: 5 sec/invoice, same-day processing",
], sz=12, c=DARK_GREEN)

# ============================================================
# SLIDE 5: ROI - COST SAVINGS
# ============================================================
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(sl, WHITE)
hdr(sl, "The ROI: $462,000 in Direct AP Savings", "Based on 3,000 invoices per month")

cd3 = CategoryChartData()
cd3.categories = ["Annual AP Cost"]
cd3.add_series("Current (Manual)", (576000,))
cd3.add_series("With PayFlow AI", (114000,))
cd3.add_series("Direct Savings", (462000,))
ch3 = sl.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(0.8), Inches(1.5), Inches(5.5), Inches(3.2), cd3).chart
ch3.has_legend = True; ch3.legend.position = XL_LEGEND_POSITION.BOTTOM; ch3.legend.font.size = Pt(11)
ch3.series[0].format.fill.solid(); ch3.series[0].format.fill.fore_color.rgb = RED
ch3.series[1].format.fill.solid(); ch3.series[1].format.fill.fore_color.rgb = BLUE
ch3.series[2].format.fill.solid(); ch3.series[2].format.fill.fore_color.rgb = GREEN
p3 = ch3.plots[0]; p3.has_data_labels = True; p3.data_labels.font.size = Pt(14)
p3.data_labels.font.bold = True; p3.data_labels.number_format = '$#,##0'; p3.gap_width = 60
ch3.value_axis.visible = False

# Savings table
data = [
    ["Category", "Current", "With AI", "Savings"],
    ["AP Staff Labor", "$360,000", "$84,000", "$276,000"],
    ["Error Correction", "$72,000", "$7,200", "$64,800"],
    ["Late Payments / Discounts", "$90,000", "$9,000", "$81,000"],
    ["Overhead", "$36,000", "$12,000", "$24,000"],
    ["Vendor Follow-ups", "$18,000", "$1,800", "$16,200"],
    ["TOTAL", "$576,000", "$114,000", "$462,000"],
]
rows, cols = len(data), len(data[0])
tbl = sl.shapes.add_table(rows, cols, Inches(6.8), Inches(1.5), Inches(5.7), Inches(2.8)).table
tbl.columns[0].width = Inches(2.0); tbl.columns[1].width = Inches(1.1)
tbl.columns[2].width = Inches(1.1); tbl.columns[3].width = Inches(1.5)
for r in range(rows):
    for c in range(cols):
        cell = tbl.cell(r, c); cell.text = data[r][c]
        p = cell.text_frame.paragraphs[0]; p.font.size = Pt(12); p.font.name = "Calibri"
        if r == 0:
            cell.fill.solid(); cell.fill.fore_color.rgb = NAVY; p.font.color.rgb = WHITE; p.font.bold = True
        elif r == rows-1:
            cell.fill.solid(); cell.fill.fore_color.rgb = LIGHT_GREEN; p.font.color.rgb = DARK_GREEN; p.font.bold = True
        else:
            cell.fill.solid(); cell.fill.fore_color.rgb = LIGHT_GRAY if r%2==0 else WHITE; p.font.color.rgb = DARK
            if c == 3: p.font.color.rgb = GREEN; p.font.bold = True
        if c > 0: p.alignment = PP_ALIGN.CENTER

stat(sl, Inches(0.8), Inches(5.0), Inches(2.8), Inches(1.5), "SOFTWARE COST", "$18,000/yr", "PayFlow AI Professional", BLUE)
stat(sl, Inches(3.9), Inches(5.0), Inches(2.8), Inches(1.5), "AP SAVINGS", "$462,000/yr", "Direct cost reduction", GREEN)
stat(sl, Inches(7.0), Inches(5.0), Inches(2.8), Inches(1.5), "ROI", "2,467%", "Return on investment", GREEN)
stat(sl, Inches(10.1), Inches(5.0), Inches(2.8), Inches(1.5), "PAYBACK", "2 Weeks", "Software pays for itself", GREEN)

# ============================================================
# SLIDE 6: REVENUE GENERATION - EARLY PAY DISCOUNTS
# ============================================================
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(sl, WHITE)
hdr(sl, "Beyond Savings: PayFlow AI Makes You Money",
    "Early-pay discount capture alone could generate $50,000 - $100,000+ per year")

# Discount opportunity chart
cd6 = CategoryChartData()
cd6.categories = ["Jan", "Feb", "Mar", "Apr", "May", "Jun", "Jul", "Aug", "Sep", "Oct", "Nov", "Dec"]
cd6.add_series("Discounts Available", (8.2, 7.8, 9.1, 8.5, 10.2, 9.8, 11.0, 10.5, 9.3, 8.8, 9.5, 10.3))
cd6.add_series("Discounts Captured (Manual)", (1.2, 0.8, 1.5, 0.9, 1.1, 1.3, 1.0, 0.7, 1.4, 0.9, 1.2, 1.1))
cd6.add_series("Discounts Captured (PayFlow AI)", (7.0, 6.6, 7.8, 7.2, 8.7, 8.3, 9.4, 9.0, 7.9, 7.5, 8.1, 8.8))

ch6 = sl.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(0.8), Inches(1.5), Inches(7.0), Inches(3.5), cd6).chart
ch6.has_legend = True; ch6.legend.position = XL_LEGEND_POSITION.BOTTOM; ch6.legend.font.size = Pt(10)
ch6.series[0].format.fill.solid(); ch6.series[0].format.fill.fore_color.rgb = MEDIUM_GRAY
ch6.series[1].format.fill.solid(); ch6.series[1].format.fill.fore_color.rgb = RED
ch6.series[2].format.fill.solid(); ch6.series[2].format.fill.fore_color.rgb = GREEN
ch6.plots[0].gap_width = 50
ch6.value_axis.tick_labels.font.size = Pt(9)
ch6.value_axis.title = None
ch6.category_axis.tick_labels.font.size = Pt(9)

# Right side explanation
txt(sl, Inches(8.3), Inches(1.5), Inches(4.5), Inches(0.4), "How It Works:", sz=20, c=NAVY, b=True)
bul(sl, Inches(8.5), Inches(2.0), Inches(4.3), Inches(3.0), [
    "Many vendors offer 2/10 Net 30 terms (2% discount if paid within 10 days)",
    "Manual AP takes 14+ days - you miss the window every time",
    "PayFlow AI processes invoices same-day and flags discount opportunities",
    "AI prioritizes discount invoices in the approval queue",
    "Monthly report shows exactly how much you captured",
], sz=13, c=GRAY, sp=12)

# Bottom comparison
stat(sl, Inches(0.8), Inches(5.3), Inches(3.7), Inches(1.5), "DISCOUNTS CAPTURED (MANUAL)", "$12,800/yr",
     "Only ~12% of available discounts", RED)
stat(sl, Inches(4.8), Inches(5.3), Inches(3.7), Inches(1.5), "DISCOUNTS CAPTURED (PAYFLOW AI)", "$96,300/yr",
     "~85% of available discounts", GREEN)
stat(sl, Inches(8.8), Inches(5.3), Inches(3.7), Inches(1.5), "ADDITIONAL REVENUE", "$83,500/yr",
     "Money you're currently leaving on the table", TEAL)

# ============================================================
# SLIDE 7: CFO FINANCIAL INTELLIGENCE
# ============================================================
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(sl, WHITE)
hdr(sl, "Financial Intelligence a CFO Actually Wants",
    "Real-time visibility into cash position, job profitability, and vendor performance")

cards_data = [
    ("Cash Flow Forecasting", [
        "Predict AP obligations 30-60 days out",
        "'You need $890K to cover payables this month'",
        "Model payment timing scenarios",
        "Plan draws and credit line usage",
    ], PURPLE, "$890K",  "payables due\nnext 30 days"),
    ("Job Profitability Alerts", [
        "Real-time cost vs budget by job",
        "'Job #2847 is 73% done but used 89% of budget'",
        "Catch bleeding projects before it's too late",
        "Alert PMs and CFO automatically",
    ], ORANGE, "89%", "of plumbing budget\nconsumed at 73%\ncomplete"),
    ("Over/Under Billing", [
        "Track billings vs work completed per job",
        "Required for WIP reporting (banks & bonding)",
        "Currently takes days in spreadsheets",
        "PayFlow does it in real-time automatically",
    ], TEAL, "$245K", "underbilled across\n4 active jobs"),
    ("Vendor Scorecard", [
        "Rate vendors on invoice accuracy",
        "Track pricing consistency over time",
        "Monitor COI and compliance status",
        "Objective data for negotiations",
    ], NAVY, "12%", "error rate on\nvendor invoices =\nleverage"),
]

for i, (title, items, color, big_num, big_sub) in enumerate(cards_data):
    col = i % 2; row = i // 2
    l = Inches(0.8) + Inches(col * 6.2)
    t = Inches(1.5) + Inches(row * 2.7)

    shape(sl, l, t, Inches(5.8), Inches(2.4), WHITE, color)
    rect(sl, l+Inches(0.05), t+Inches(0.05), Inches(5.7), Inches(0.08), color)

    # Left side - title and bullets
    txt(sl, l+Inches(0.2), t+Inches(0.2), Inches(3.5), Inches(0.35), title, sz=17, c=NAVY, b=True)
    bul(sl, l+Inches(0.2), t+Inches(0.6), Inches(3.5), Inches(1.6), items, sz=11, c=GRAY, sp=5)

    # Right side - big number
    shape(sl, l+Inches(3.9), t+Inches(0.25), Inches(1.7), Inches(1.9), LIGHT_GRAY)
    txt(sl, l+Inches(3.9), t+Inches(0.35), Inches(1.7), Inches(0.6), big_num, sz=28, c=color, b=True, a=PP_ALIGN.CENTER)
    txt(sl, l+Inches(3.9), t+Inches(0.95), Inches(1.7), Inches(1.0), big_sub, sz=9, c=GRAY, a=PP_ALIGN.CENTER)

# ============================================================
# SLIDE 8: AP AGING & TAX OPTIMIZATION
# ============================================================
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(sl, WHITE)
hdr(sl, "Visibility You Don't Have Today")

# AP Aging chart
cd8 = CategoryChartData()
cd8.categories = ["Current\n(0-30 days)", "31-60\ndays", "61-90\ndays", "90+\ndays"]
cd8.add_series("Amount ($K)", (280, 140, 85, 35))
ch8 = sl.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(0.8), Inches(1.5), Inches(5.5), Inches(2.8), cd8).chart
ch8.has_legend = False
for i, c in enumerate([GREEN, GOLD, ORANGE, RED]):
    ch8.series[0].points[i].format.fill.solid()
    ch8.series[0].points[i].format.fill.fore_color.rgb = c
p8 = ch8.plots[0]; p8.has_data_labels = True; p8.data_labels.font.size = Pt(12)
p8.data_labels.font.bold = True; p8.data_labels.number_format = '$#,##0"K"'
p8.data_labels.position = XL_LABEL_POSITION.OUTSIDE_END; p8.gap_width = 60
ch8.value_axis.visible = False; ch8.category_axis.tick_labels.font.size = Pt(10)

txt(sl, Inches(0.8), Inches(1.2), Inches(5.5), Inches(0.3), "AP Aging Dashboard", sz=18, c=NAVY, b=True, a=PP_ALIGN.CENTER)

# Right side - what you see
txt(sl, Inches(6.8), Inches(1.2), Inches(5.5), Inches(0.3), "What This Tells You:", sz=18, c=NAVY, b=True)
bul(sl, Inches(7.0), Inches(1.6), Inches(5.3), Inches(2.5), [
    "Payment health at a glance - are you falling behind?",
    "Vendor relationship risk - 90+ day payables = credit holds",
    "Credit holds stop materials from shipping to your job sites",
    "Trend analysis: 'Your 90+ day payables grew 40% in 3 months'",
    "CFO sees this in real-time, not at month-end close",
], sz=13, c=GRAY, sp=8)

# Bottom section - additional value
txt(sl, Inches(0.8), Inches(4.6), Inches(11.5), Inches(0.4), "More Built-In Financial Intelligence:", sz=20, c=NAVY, b=True)

bottom_cards = [
    ("Tax Optimization", [
        "Flag tax-exempt construction purchases",
        "Track use tax on out-of-state buys",
        "Auto-identify 1099 vendors",
        "Potential savings: $20-50K/year",
    ], TEAL),
    ("Audit-Ready Package", [
        "One-click auditor export",
        "Complete invoice register with images",
        "Full approval chain documentation",
        "Weeks of audit prep become instant",
    ], NAVY),
    ("Bonding Capacity", [
        "Real-time working capital tracking",
        "Over/under billing for WIP reports",
        "AP aging for bonding applications",
        "Know your capacity before you bid",
    ], PURPLE),
    ("Retention Tracker", [
        "Retainage held per vendor per job",
        "Release conditions and deadlines",
        "Prevents premature releases",
        "Cash flow impact of held retainage",
    ], ORANGE),
]

for i, (title, items, color) in enumerate(bottom_cards):
    l = Inches(0.6) + Inches(i * 3.15)
    card(sl, l, Inches(5.1), Inches(2.95), Inches(2.0), title, items, color, title_sz=14)

# ============================================================
# SLIDE 9: 3-YEAR TOTAL VALUE
# ============================================================
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(sl, WHITE)
hdr(sl, "3-Year Total Value: $1.7 Million+",
    "AP savings + discount capture + tax optimization + avoided losses")

# Stacked bar chart
cd9 = CategoryChartData()
cd9.categories = ["Year 1", "Year 2", "Year 3"]
cd9.add_series("AP Cost Savings", (462000, 487000, 513000))
cd9.add_series("Early-Pay Discounts Captured", (83500, 92000, 101000))
cd9.add_series("Tax Optimization", (25000, 35000, 40000))
cd9.add_series("Avoided Job Losses (early alerts)", (50000, 75000, 100000))

ch9 = sl.shapes.add_chart(XL_CHART_TYPE.COLUMN_STACKED, Inches(0.8), Inches(1.5), Inches(7.0), Inches(4.0), cd9).chart
ch9.has_legend = True; ch9.legend.position = XL_LEGEND_POSITION.BOTTOM; ch9.legend.font.size = Pt(11)
ch9.series[0].format.fill.solid(); ch9.series[0].format.fill.fore_color.rgb = GREEN
ch9.series[1].format.fill.solid(); ch9.series[1].format.fill.fore_color.rgb = TEAL
ch9.series[2].format.fill.solid(); ch9.series[2].format.fill.fore_color.rgb = PURPLE
ch9.series[3].format.fill.solid(); ch9.series[3].format.fill.fore_color.rgb = ORANGE
ch9.plots[0].has_data_labels = True; ch9.plots[0].data_labels.font.size = Pt(10)
ch9.plots[0].data_labels.font.bold = True; ch9.plots[0].data_labels.number_format = '$#,##0'
ch9.plots[0].gap_width = 80
ch9.value_axis.visible = False

# Right side summary
txt(sl, Inches(8.3), Inches(1.5), Inches(4.5), Inches(0.4), "3-Year Value Summary", sz=22, c=NAVY, b=True)

val_table = [
    ["Value Source", "Year 1", "Year 2", "Year 3", "Total"],
    ["AP Cost Savings", "$462K", "$487K", "$513K", "$1.46M"],
    ["Discount Capture", "$84K", "$92K", "$101K", "$277K"],
    ["Tax Optimization", "$25K", "$35K", "$40K", "$100K"],
    ["Avoided Job Losses", "$50K", "$75K", "$100K", "$225K"],
    ["TOTAL VALUE", "$621K", "$689K", "$754K", "$2.06M"],
    ["Software Cost", "$18K", "$18K", "$18K", "$54K"],
    ["NET VALUE", "$603K", "$671K", "$736K", "$2.01M"],
]

rows9, cols9 = len(val_table), len(val_table[0])
tbl9 = sl.shapes.add_table(rows9, cols9, Inches(8.3), Inches(2.0), Inches(4.5), Inches(3.2)).table
tbl9.columns[0].width = Inches(1.4)
for c in range(1, 5): tbl9.columns[c].width = Inches(0.78)

for r in range(rows9):
    for c in range(cols9):
        cell = tbl9.cell(r, c); cell.text = val_table[r][c]
        p = cell.text_frame.paragraphs[0]; p.font.size = Pt(10); p.font.name = "Calibri"
        if r == 0:
            cell.fill.solid(); cell.fill.fore_color.rgb = NAVY; p.font.color.rgb = WHITE; p.font.bold = True
        elif r == 5:
            cell.fill.solid(); cell.fill.fore_color.rgb = LIGHT_BLUE; p.font.color.rgb = NAVY; p.font.bold = True
        elif r == 7:
            cell.fill.solid(); cell.fill.fore_color.rgb = LIGHT_GREEN; p.font.color.rgb = DARK_GREEN; p.font.bold = True
        else:
            cell.fill.solid(); cell.fill.fore_color.rgb = LIGHT_GRAY if r%2==0 else WHITE; p.font.color.rgb = DARK
        if c > 0: p.alignment = PP_ALIGN.CENTER

shape(sl, Inches(8.3), Inches(5.5), Inches(4.5), Inches(1.3), LIGHT_GREEN, GREEN)
txt(sl, Inches(8.3), Inches(5.6), Inches(4.5), Inches(0.3), "3-Year Net Value", sz=14, c=DARK_GREEN, b=True, a=PP_ALIGN.CENTER)
txt(sl, Inches(8.3), Inches(5.9), Inches(4.5), Inches(0.6), "$2.01 Million", sz=40, c=DARK_GREEN, b=True, a=PP_ALIGN.CENTER)
txt(sl, Inches(8.3), Inches(6.4), Inches(4.5), Inches(0.3), "vs $54K software investment (3,728% ROI)", sz=11, c=DARK_GREEN, a=PP_ALIGN.CENTER)

# ============================================================
# SLIDE 10: CONSTRUCTION SPECIFIC
# ============================================================
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(sl, WHITE)
hdr(sl, "Built for Construction, Not Adapted From Generic AP",
    "Generic tools handle 3 of 13 construction requirements. We handle all 13.")

comp = [
    ["Capability", "Generic AP\n(Stampli, Vic.ai)", "PayFlow AI"],
    ["Standard invoice extraction", "Yes", "Yes"],
    ["GL code prediction", "Yes", "Yes"],
    ["Approval workflows", "Yes", "Yes"],
    ["AIA G702/G703 billing", "No", "Yes"],
    ["Lien waiver tracking", "No", "Yes"],
    ["Retention management", "No", "Yes"],
    ["5D job costing (GL+job+phase+code+type)", "No", "Yes"],
    ["Change order tracking", "No", "Yes"],
    ["COI expiration alerts", "No", "Yes"],
    ["Prevailing wage verification", "No", "Yes"],
    ["Multi-entity routing", "No", "Yes"],
    ["Sage integration", "Limited", "Native"],
    ["Procore integration", "No", "Native"],
]
rows10, cols10 = len(comp), len(comp[0])
tbl10 = sl.shapes.add_table(rows10, cols10, Inches(0.8), Inches(1.5), Inches(8.0), Inches(5.5)).table
tbl10.columns[0].width = Inches(3.5); tbl10.columns[1].width = Inches(2.2); tbl10.columns[2].width = Inches(2.3)
for r in range(rows10):
    for c in range(cols10):
        cell = tbl10.cell(r, c); cell.text = comp[r][c]
        p = cell.text_frame.paragraphs[0]; p.font.size = Pt(13); p.font.name = "Calibri"
        if r == 0:
            cell.fill.solid(); cell.fill.fore_color.rgb = NAVY; p.font.color.rgb = WHITE; p.font.bold = True
        else:
            cell.fill.solid(); cell.fill.fore_color.rgb = LIGHT_GRAY if r%2==0 else WHITE; p.font.color.rgb = DARK
            if c==1 and comp[r][c]=="No": p.font.color.rgb = RED; p.font.bold = True
            elif c==2 and comp[r][c] in ("Yes","Native"): p.font.color.rgb = GREEN; p.font.bold = True
        if c > 0: p.alignment = PP_ALIGN.CENTER

shape(sl, Inches(9.3), Inches(1.5), Inches(3.5), Inches(5.5), LIGHT_BLUE, BLUE)
txt(sl, Inches(9.5), Inches(1.7), Inches(3.1), Inches(0.4), "The Gap", sz=22, c=NAVY, b=True, a=PP_ALIGN.CENTER)
bul(sl, Inches(9.5), Inches(2.2), Inches(3.1), Inches(4.5), [
    "Generic AP tools handle 3 of 13 construction requirements.",
    "PayFlow AI handles all 13.",
    "Your AP team currently fills this gap manually.",
    "That's where the $576K goes.",
    "Plus: cash flow forecasting, discount capture, and job profitability alerts that no AP tool offers.",
], sz=14, c=NAVY, sp=16)

# ============================================================
# SLIDE 11: INTEGRATIONS
# ============================================================
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(sl, WHITE)
hdr(sl, "Plugs Into Sage & Procore On Day One", "No rip-and-replace. Works with what you already use.")

integrations = [
    ("Sage", "Your Accounting System", ["Two-way sync: chart of accounts", "Import vendors, jobs, cost codes", "Push approved invoices directly", "Sync payment status"], NAVY),
    ("Procore", "Your Project Management", ["Pull POs, budgets, commitments", "Import change orders", "Three-way match against POs", "Real-time contract tracking"], BLUE),
    ("Email", "Invoice Ingestion", ["Dedicated inbox for your company", "AI identifies invoices vs spam", "Auto-processes on arrival", "Vendors email invoices directly"], GREEN),
    ("Banking", "Payments", ["Schedule approved payments", "ACH, check, or card", "Auto-reconciliation", "Cash flow visibility"], ORANGE),
]
for i, (name, sub, items, color) in enumerate(integrations):
    l = Inches(0.5) + Inches(i * 3.2); t = Inches(1.5)
    sh = shape(sl, l, t, Inches(2.9), Inches(4.0), WHITE, color)
    rect(sl, l+Inches(0.05), t+Inches(0.05), Inches(2.8), Inches(0.65), color)
    txt(sl, l+Inches(0.1), t+Inches(0.07), Inches(2.7), Inches(0.4), name, sz=22, c=WHITE, b=True, a=PP_ALIGN.CENTER)
    txt(sl, l+Inches(0.1), t+Inches(0.38), Inches(2.7), Inches(0.3), sub, sz=10, c=RGBColor(0xd0,0xe0,0xf0), a=PP_ALIGN.CENTER)
    bul(sl, l+Inches(0.15), t+Inches(0.9), Inches(2.6), Inches(2.8), items, sz=13, c=GRAY, sp=10)

shape(sl, Inches(0.8), Inches(5.9), Inches(11.7), Inches(1.0), LIGHT_BLUE, BLUE)
txt(sl, Inches(1.0), Inches(6.0), Inches(11.3), Inches(0.3),
    "Invoice In  >  AI Extracts  >  Matches Procore PO  >  Routes for Approval  >  Posts to Sage  >  Payment Scheduled", sz=14, c=NAVY, b=True, a=PP_ALIGN.CENTER)
txt(sl, Inches(1.0), Inches(6.35), Inches(11.3), Inches(0.4),
    "Your team handles exceptions. The AI handles the other 85%.", sz=14, c=NAVY, a=PP_ALIGN.CENTER)

# ============================================================
# SLIDE 12: ACCURACY & SECURITY
# ============================================================
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(sl, WHITE)
hdr(sl, "Accuracy & Security You Can Trust")

# Accuracy chart
cd12 = CategoryChartData()
cd12.categories = ["Invoice #", "Date", "Vendor", "Amount", "Line Items", "PO #", "GL Code"]
cd12.add_series("AI Accuracy %", (99.5, 98.5, 97.5, 99.5, 96, 97, 92))
ch12 = sl.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(0.8), Inches(1.3), Inches(5.5), Inches(2.8), cd12).chart
ch12.has_legend = False
ch12.series[0].format.fill.solid(); ch12.series[0].format.fill.fore_color.rgb = BLUE
p12 = ch12.plots[0]; p12.has_data_labels = True; p12.data_labels.font.size = Pt(10)
p12.data_labels.font.bold = True; p12.data_labels.number_format = '0.0"%"'
p12.data_labels.position = XL_LABEL_POSITION.OUTSIDE_END; p12.gap_width = 60
ch12.value_axis.minimum_scale = 80; ch12.value_axis.maximum_scale = 100
ch12.value_axis.tick_labels.font.size = Pt(9); ch12.category_axis.tick_labels.font.size = Pt(9)

txt(sl, Inches(6.8), Inches(1.3), Inches(5.5), Inches(0.4), "Built-In Safety Net", sz=20, c=NAVY, b=True)
bul(sl, Inches(7.0), Inches(1.8), Inches(5.3), Inches(2.5), [
    "Confidence scoring on every field",
    "Low-confidence invoices route to humans automatically",
    "Every AI decision in full audit trail",
    "Math validation on all line items and totals",
    "Duplicate detection across all vendors",
    "Human override on any field, always",
    "Shadow mode proves accuracy before going live",
], sz=12, c=GRAY, sp=6)

# Security row
security = [
    ("Encryption", "AES-256 at rest\nTLS 1.2+ in transit"),
    ("AI Privacy", "Never trains on your data\nAPI processing only"),
    ("Audit Trail", "Every action logged\nFull compliance history"),
    ("Access Control", "Role-based: AP, PM,\nController, CFO"),
    ("SOC 2", "Enterprise controls\nAnnual audits"),
    ("Data Ownership", "Your data, always\nFull export anytime"),
]
for i, (title, desc) in enumerate(security):
    l = Inches(0.8) + Inches(i * 2.05)
    sh = shape(sl, l, Inches(4.5), Inches(1.9), Inches(2.3), LIGHT_GRAY, MEDIUM_GRAY)
    txt(sl, l+Inches(0.1), Inches(4.6), Inches(1.7), Inches(0.3), title, sz=13, c=NAVY, b=True, a=PP_ALIGN.CENTER)
    multiline(sl, l+Inches(0.1), Inches(4.95), Inches(1.7), Inches(1.5), desc.split("\n"), sz=11, c=GRAY, a=PP_ALIGN.CENTER)

# ============================================================
# SLIDE 13: IMPLEMENTATION
# ============================================================
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(sl, WHITE)
hdr(sl, "Live in 8 Weeks - Zero Disruption",
    "Your team keeps working normally while we set up in parallel")

phases = [
    ("Weeks 1-2", "Setup", "Connect Sage & Procore\nImport chart of accounts\nConfigure cost codes", BLUE),
    ("Weeks 3-4", "Calibrate", "Process real invoices\nTune AI to your patterns\nVerify accuracy", BLUE),
    ("Weeks 5-6", "Shadow Mode", "AI runs alongside team\nCompare results\nBuild confidence", GREEN),
    ("Weeks 7-8", "Go Live", "AI handles routine\nTeam handles exceptions\nStart saving day one", GREEN),
]
for i, (time, title, desc, color) in enumerate(phases):
    l = Inches(0.5) + Inches(i * 3.2); t = Inches(1.5)
    sh = shape(sl, l, t, Inches(2.95), Inches(2.8), WHITE, color)
    rect(sl, l+Inches(0.05), t+Inches(0.05), Inches(2.85), Inches(0.5), color)
    txt(sl, l+Inches(0.1), t+Inches(0.08), Inches(2.75), Inches(0.4), time, sz=14, c=WHITE, b=True, a=PP_ALIGN.CENTER)
    txt(sl, l+Inches(0.15), t+Inches(0.65), Inches(2.65), Inches(0.35), title, sz=18, c=NAVY, b=True, a=PP_ALIGN.CENTER)
    multiline(sl, l+Inches(0.15), t+Inches(1.1), Inches(2.65), Inches(1.5), desc.split("\n"), sz=12, c=GRAY, a=PP_ALIGN.CENTER, sp=4)

# Guarantees
guarantees = [
    ("Zero Disruption", "Your team works normally.\nAI runs in parallel."),
    ("Immediate Value", "ROI from week one\nof go-live."),
    ("Risk-Free", "Shadow mode proves it.\nCancel anytime."),
]
for i, (gt, gd) in enumerate(guarantees):
    l = Inches(0.8) + Inches(i * 4.1)
    shape(sl, l, Inches(4.7), Inches(3.7), Inches(1.3), LIGHT_BLUE, BLUE)
    txt(sl, l+Inches(0.2), Inches(4.8), Inches(3.3), Inches(0.3), gt, sz=16, c=NAVY, b=True, a=PP_ALIGN.CENTER)
    multiline(sl, l+Inches(0.2), Inches(5.15), Inches(3.3), Inches(0.7), gd.split("\n"), sz=12, c=NAVY, a=PP_ALIGN.CENTER)

# What shadow mode is
shape(sl, Inches(0.8), Inches(6.2), Inches(11.7), Inches(0.9), LIGHT_GREEN, GREEN)
txt(sl, Inches(1.0), Inches(6.25), Inches(11.3), Inches(0.3),
    "What is Shadow Mode?", sz=14, c=DARK_GREEN, b=True, a=PP_ALIGN.CENTER)
txt(sl, Inches(1.0), Inches(6.55), Inches(11.3), Inches(0.4),
    "AI processes every invoice alongside your team - but doesn't post anything. You compare results side by side. When you trust the numbers, flip the switch.", sz=13, c=DARK_GREEN, a=PP_ALIGN.CENTER)

# ============================================================
# SLIDE 14: PRICING
# ============================================================
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(sl, WHITE)
hdr(sl, "Investment")

tiers = [
    ("Starter", "$500", "/month", [
        "Up to 1,000 invoices/mo",
        "5 users",
        "Sage integration",
        "AI extraction & coding",
        "Email ingestion",
        "AP aging dashboard",
        "Standard support",
    ], BLUE, False),
    ("Professional", "$1,500", "/month", [
        "Up to 5,000 invoices/mo",
        "Unlimited users",
        "Sage + Procore integration",
        "AIA billing & lien waivers",
        "Retention tracking",
        "Fraud detection",
        "Cash flow forecasting",
        "Early-pay discount capture",
        "Job profitability alerts",
        "Vendor spend analytics",
        "Priority support",
    ], GREEN, True),
    ("Enterprise", "Custom", "", [
        "Unlimited invoices",
        "Unlimited users",
        "All integrations",
        "Over/under billing (WIP)",
        "Bonding capacity tracking",
        "Tax optimization",
        "Audit-ready package",
        "Vendor scorecard",
        "Dedicated account manager",
        "Custom workflows & SLA",
    ], NAVY, False),
]
for i, (name, price, period, items, color, rec) in enumerate(tiers):
    l = Inches(1.0) + Inches(i * 3.9); t = Inches(1.3); h = Inches(5.6)
    if rec:
        rect(sl, l-Inches(0.05), t-Inches(0.35), Inches(3.6), Inches(0.35), color)
        txt(sl, l-Inches(0.05), t-Inches(0.35), Inches(3.6), Inches(0.3), "RECOMMENDED", sz=11, c=WHITE, b=True, a=PP_ALIGN.CENTER)
    sh = shape(sl, l, t, Inches(3.5), h, WHITE, color)
    if rec: sh.line.width = Pt(3)
    txt(sl, l+Inches(0.1), t+Inches(0.15), Inches(3.3), Inches(0.4), name, sz=22, c=color, b=True, a=PP_ALIGN.CENTER)
    txt(sl, l+Inches(0.1), t+Inches(0.55), Inches(3.3), Inches(0.6), price, sz=42, c=DARK, b=True, a=PP_ALIGN.CENTER)
    if period: txt(sl, l+Inches(0.1), t+Inches(1.05), Inches(3.3), Inches(0.3), period, sz=14, c=GRAY, a=PP_ALIGN.CENTER)
    for j, item in enumerate(items):
        txt(sl, l+Inches(0.25), t+Inches(1.4)+Inches(j*0.35), Inches(3.0), Inches(0.3), "  "+item, sz=12, c=GRAY)

txt(sl, Inches(0.8), Inches(7.05), Inches(11.5), Inches(0.3),
    "All plans: AI extraction, GL coding, approval workflows, duplicate detection, audit trail. Annual billing available (2 months free).", sz=12, c=GRAY, a=PP_ALIGN.CENTER)

# ============================================================
# SLIDE 15: TOTAL VALUE SUMMARY
# ============================================================
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(sl, WHITE)
hdr(sl, "Total First-Year Value: $603,000+",
    "For a $18,000 investment")

# Value waterfall
cd15 = CategoryChartData()
cd15.categories = ["AP Cost\nSavings", "Early-Pay\nDiscounts", "Tax\nOptimization", "Avoided\nJob Losses", "Software\nCost", "NET VALUE"]
cd15.add_series("Value", (462000, 83500, 25000, 50000, -18000, 602500))

ch15 = sl.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(1.5), Inches(1.5), Inches(10.0), Inches(3.5), cd15).chart
ch15.has_legend = False

for i, c in enumerate([GREEN, TEAL, PURPLE, ORANGE, RED, GREEN]):
    ch15.series[0].points[i].format.fill.solid()
    ch15.series[0].points[i].format.fill.fore_color.rgb = c

p15 = ch15.plots[0]; p15.has_data_labels = True; p15.data_labels.font.size = Pt(12)
p15.data_labels.font.bold = True; p15.data_labels.number_format = '$#,##0'; p15.gap_width = 40
ch15.value_axis.visible = False; ch15.category_axis.tick_labels.font.size = Pt(11)

# Summary cards
stat(sl, Inches(0.8), Inches(5.3), Inches(2.3), Inches(1.5), "AP SAVINGS", "$462K", "81% cost reduction", GREEN)
stat(sl, Inches(3.4), Inches(5.3), Inches(2.3), Inches(1.5), "DISCOUNTS", "$83.5K", "Money found", TEAL)
stat(sl, Inches(6.0), Inches(5.3), Inches(2.3), Inches(1.5), "TAX + LOSSES", "$75K", "Optimization + prevention", PURPLE)
stat(sl, Inches(8.6), Inches(5.3), Inches(2.3), Inches(1.5), "SOFTWARE", "-$18K", "Annual cost", RED)
stat(sl, Inches(11.2), Inches(5.3), Inches(1.7), Inches(1.5), "NET", "$603K", "First year", GREEN)

# ============================================================
# SLIDE 16: CTA
# ============================================================
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(sl, NAVY)
rect(sl, Inches(0), Inches(3.0), Inches(13.333), Inches(0.06), BLUE)

txt(sl, Inches(1), Inches(0.6), Inches(11.3), Inches(0.8),
    "Ready to Save $603,000 in Year One?", sz=44, c=WHITE, b=True, a=PP_ALIGN.CENTER)
txt(sl, Inches(2), Inches(1.5), Inches(9.3), Inches(0.5),
    "Three steps to transform your AP - zero risk:", sz=20, c=RGBColor(0xa0,0xce,0xf0), a=PP_ALIGN.CENTER)

cta = [
    ("1", "Free Invoice Test", "Share 50 sample invoices. We run them through the AI\nand show you extraction results. No commitment. No cost."),
    ("2", "2-Week Shadow Pilot", "AI processes your real invoices alongside your team.\nCompare accuracy side by side. See the savings in real numbers."),
    ("3", "Go Live With Confidence", "Only when you've seen the proof. Start saving from day one.\nCancel anytime. No long-term contract required."),
]
for i, (num, title, desc) in enumerate(cta):
    t = Inches(3.4) + Inches(i * 1.15)
    circle_step(sl, Inches(2.3), t, num, BLUE)
    txt(sl, Inches(3.3), t+Inches(0.05), Inches(7), Inches(0.35), title, sz=22, c=WHITE, b=True)
    multiline(sl, Inches(3.3), t+Inches(0.4), Inches(7.5), Inches(0.6), desc.split("\n"), sz=13, c=RGBColor(0xcb,0xd5,0xe0), sp=1)

txt(sl, Inches(1), Inches(6.8), Inches(11.3), Inches(0.4),
    "PayFlow AI  |  AI-Powered AP Automation for Construction  |  payflowai.com", sz=14, c=RGBColor(0x71,0x80,0x96), a=PP_ALIGN.CENTER)

# SAVE
output = "C:/users/25badmin/projects/accounts-payable-research/PayFlow_AI_Sales_Deck_v3.pptx"
prs.save(output)
print(f"Saved: {output}")
