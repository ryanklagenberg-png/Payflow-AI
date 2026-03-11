"""
PayFlow AI - How It Works (Step-by-Step)
Comprehensive PowerPoint showing the full invoice lifecycle
with KOJO, Procore, Sage, and other integrations
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_LABEL_POSITION
from pptx.chart.data import CategoryChartData

# ── Brand Colors ──
NAVY = RGBColor(0x0B, 0x1D, 0x3A)
DARK_BLUE = RGBColor(0x12, 0x2B, 0x4F)
ACCENT_BLUE = RGBColor(0x1E, 0x90, 0xFF)
LIGHT_BLUE = RGBColor(0x4D, 0xA8, 0xFF)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xE8, 0xEC, 0xF1)
MEDIUM_GRAY = RGBColor(0x8A, 0x95, 0xA5)
DARK_GRAY = RGBColor(0x2C, 0x3E, 0x50)
GREEN = RGBColor(0x2E, 0xCC, 0x71)
DARK_GREEN = RGBColor(0x1A, 0x9B, 0x50)
ORANGE = RGBColor(0xF3, 0x9C, 0x12)
RED = RGBColor(0xE7, 0x4C, 0x3C)
PURPLE = RGBColor(0x9B, 0x59, 0xB6)
TEAL = RGBColor(0x00, 0xB4, 0xD8)
KOJO_ORANGE = RGBColor(0xFF, 0x6B, 0x35)
PROCORE_ORANGE = RGBColor(0xF4, 0x7E, 0x20)
SAGE_GREEN = RGBColor(0x00, 0xA8, 0x4F)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)


def add_bg(slide, color=NAVY):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape(slide, left, top, width, height, fill_color, border_color=None, radius=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(2)
    else:
        shape.line.fill.background()
    return shape


def add_text(slide, left, top, width, height, text, font_size=18, color=WHITE, bold=False, alignment=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = alignment
    return tf


def add_bullet_list(slide, left, top, width, height, items, font_size=16, color=WHITE, spacing=Pt(8)):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.space_after = spacing
        p.level = 0
    return tf


def add_step_number(slide, left, top, number, color=ACCENT_BLUE):
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, Inches(0.7), Inches(0.7))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.text = str(number)
    p.font.size = Pt(24)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    tf.paragraphs[0].space_before = Pt(0)
    tf.paragraphs[0].space_after = Pt(0)
    return shape


def add_icon_box(slide, left, top, width, height, title, subtitle, fill_color, border_color=None):
    shape = add_shape(slide, left, top, width, height, fill_color, border_color)
    tf = shape.text_frame
    tf.word_wrap = True
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(14)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p2 = tf.add_paragraph()
    p2.text = subtitle
    p2.font.size = Pt(11)
    p2.font.color.rgb = RGBColor(0xCC, 0xCC, 0xCC)
    p2.alignment = PP_ALIGN.CENTER
    return shape


def add_arrow(slide, left, top, width, height, color=ACCENT_BLUE):
    shape = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_down_arrow(slide, left, top, width, height, color=ACCENT_BLUE):
    shape = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


# ═══════════════════════════════════════════════
# SLIDE 1: Title Slide
# ═══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)

# Accent bar
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.06), ACCENT_BLUE)

add_text(slide, Inches(1), Inches(1.5), Inches(11), Inches(1.2),
         "PayFlow AI", 54, WHITE, True, PP_ALIGN.CENTER)
add_text(slide, Inches(1), Inches(2.7), Inches(11), Inches(0.8),
         "How It Works: Step-by-Step", 36, LIGHT_BLUE, False, PP_ALIGN.CENTER)

# Subtitle
add_text(slide, Inches(2), Inches(3.8), Inches(9), Inches(0.6),
         "The Complete Invoice Lifecycle for Construction Companies", 20, MEDIUM_GRAY, False, PP_ALIGN.CENTER)

# Integration logos bar
add_shape(slide, Inches(2.5), Inches(5.0), Inches(8.333), Inches(1.2), DARK_BLUE, ACCENT_BLUE)
add_text(slide, Inches(2.5), Inches(5.05), Inches(8.333), Inches(0.5),
         "Fully Integrated With:", 14, MEDIUM_GRAY, False, PP_ALIGN.CENTER)

# Integration boxes
integrations = [
    ("KOJO", "Procurement", KOJO_ORANGE),
    ("Procore", "Project Mgmt", PROCORE_ORANGE),
    ("Sage 300 CRE", "Accounting", SAGE_GREEN),
    ("Email / Scanner", "Invoice Intake", ACCENT_BLUE),
]
x_start = 3.2
for label, sub, color in integrations:
    add_icon_box(slide, Inches(x_start), Inches(5.45), Inches(1.6), Inches(0.65), label, sub, color)
    x_start += 1.85

add_text(slide, Inches(1), Inches(6.7), Inches(11), Inches(0.4),
         "March 2026", 14, MEDIUM_GRAY, False, PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════
# SLIDE 2: The Big Picture - System Overview
# ═══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.06), ACCENT_BLUE)

add_text(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.7),
         "The Big Picture", 36, WHITE, True)
add_text(slide, Inches(0.8), Inches(0.9), Inches(11), Inches(0.5),
         "How invoices flow from receipt to payment across all connected systems", 16, MEDIUM_GRAY)

# Flow diagram - top row: external systems
# KOJO
add_icon_box(slide, Inches(0.5), Inches(1.8), Inches(2.2), Inches(1.0),
             "KOJO", "Purchase Orders & Materials", KOJO_ORANGE)
# Procore
add_icon_box(slide, Inches(3.2), Inches(1.8), Inches(2.2), Inches(1.0),
             "Procore", "Jobs, Budgets, Contracts", PROCORE_ORANGE)
# Email/Scanner
add_icon_box(slide, Inches(5.9), Inches(1.8), Inches(2.2), Inches(1.0),
             "Email / Scanner", "Invoice Intake", ACCENT_BLUE)
# Sage
add_icon_box(slide, Inches(8.6), Inches(1.8), Inches(2.2), Inches(1.0),
             "Sage 300 CRE", "GL, Payments, Reporting", SAGE_GREEN)
# Additional
add_icon_box(slide, Inches(11.3), Inches(1.8), Inches(1.7), Inches(1.0),
             "More Apps", "Viewpoint, CMiC...", PURPLE)

# Down arrows
for x in [1.3, 4.0, 6.7, 9.4, 11.85]:
    add_down_arrow(slide, Inches(x), Inches(2.9), Inches(0.35), Inches(0.4), ACCENT_BLUE)

# Central PayFlow AI box
add_shape(slide, Inches(0.5), Inches(3.5), Inches(12.5), Inches(3.5), DARK_BLUE, ACCENT_BLUE)
add_text(slide, Inches(0.5), Inches(3.55), Inches(12.5), Inches(0.6),
         "PayFlow AI Engine", 28, ACCENT_BLUE, True, PP_ALIGN.CENTER)

# Internal steps
steps = [
    ("1. Receive", "Invoices arrive via\nemail, upload, scan,\nor API", ACCENT_BLUE),
    ("2. Extract", "Claude AI reads every\nfield, line item,\nand term", RGBColor(0x00, 0x96, 0xC7)),
    ("3. Code", "Auto-assign job,\ncost code, GL,\nphase, cost type", TEAL),
    ("4. Match", "Three-way match\nvs PO (KOJO) and\nbudget (Procore)", GREEN),
    ("5. Validate", "Flag duplicates,\nmath errors, fraud,\nover-budget", ORANGE),
    ("6. Route", "Send to right\napprover based on\njob, amount, trade", PURPLE),
    ("7. Post", "Push approved\ninvoice to Sage\nfor payment", SAGE_GREEN),
]

x = 0.7
for title, desc, color in steps:
    add_shape(slide, Inches(x), Inches(4.2), Inches(1.55), Inches(2.5), color)
    add_text(slide, Inches(x + 0.05), Inches(4.3), Inches(1.45), Inches(0.4),
             title, 13, WHITE, True, PP_ALIGN.CENTER)
    add_text(slide, Inches(x + 0.05), Inches(4.75), Inches(1.45), Inches(1.8),
             desc, 11, RGBColor(0xEE, 0xEE, 0xEE), False, PP_ALIGN.CENTER)
    if x < 11.0:
        pass  # arrows between boxes handled by proximity
    x += 1.72


# ═══════════════════════════════════════════════
# SLIDE 3: Step 1 - Invoice Intake
# ═══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.06), ACCENT_BLUE)

add_step_number(slide, Inches(0.6), Inches(0.3), "1")
add_text(slide, Inches(1.5), Inches(0.3), Inches(10), Inches(0.7),
         "Invoice Intake — How Invoices Enter the System", 32, WHITE, True)
add_text(slide, Inches(1.5), Inches(0.9), Inches(10), Inches(0.5),
         "Multiple channels, one unified inbox", 16, MEDIUM_GRAY)

# Four intake channels
channels = [
    ("Email Forwarding", "invoices@yourcompany.payflowai.com",
     ["Vendors email invoices directly", "Auto-forwarded from AP inbox",
      "Attachments extracted automatically", "PDF, PNG, JPG supported"],
     ACCENT_BLUE),
    ("Direct Upload", "Drag & drop or bulk upload",
     ["Web dashboard upload", "Bulk upload (50+ at once)",
      "Mobile photo capture", "Any file format"],
     TEAL),
    ("Scanner Integration", "Network scanner auto-feed",
     ["Scan-to-folder watched by PayFlow", "Copier/scanner direct integration",
      "Batch scanning supported", "Auto-splits multi-page scans"],
     PURPLE),
    ("KOJO / Vendor Portal", "Invoices attached to POs",
     ["KOJO sends invoices with PO reference", "Vendor portal self-service",
      "Pre-linked to purchase order", "Fastest path to three-way match"],
     KOJO_ORANGE),
]

x = 0.5
for title, subtitle, items, color in channels:
    add_shape(slide, Inches(x), Inches(1.6), Inches(2.95), Inches(5.2), DARK_BLUE, color)
    add_text(slide, Inches(x + 0.15), Inches(1.75), Inches(2.65), Inches(0.5),
             title, 18, color, True, PP_ALIGN.CENTER)
    add_text(slide, Inches(x + 0.15), Inches(2.2), Inches(2.65), Inches(0.4),
             subtitle, 11, MEDIUM_GRAY, False, PP_ALIGN.CENTER)

    y = 2.7
    for item in items:
        add_shape(slide, Inches(x + 0.15), Inches(y), Inches(2.65), Inches(0.55), RGBColor(0x18, 0x30, 0x55))
        add_text(slide, Inches(x + 0.3), Inches(y + 0.05), Inches(2.4), Inches(0.45),
                 item, 12, WHITE)
        y += 0.6
    x += 3.15

# Bottom note
add_shape(slide, Inches(0.5), Inches(7.0), Inches(12.5), Inches(0.35), DARK_GREEN)
add_text(slide, Inches(0.5), Inches(7.0), Inches(12.5), Inches(0.35),
         "All channels converge into a single queue — no invoice gets lost regardless of how it arrives",
         12, WHITE, False, PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════
# SLIDE 4: Step 2 - AI Extraction
# ═══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.06), ACCENT_BLUE)

add_step_number(slide, Inches(0.6), Inches(0.3), "2")
add_text(slide, Inches(1.5), Inches(0.3), Inches(10), Inches(0.7),
         "AI Extraction — Claude Reads Every Invoice", 32, WHITE, True)
add_text(slide, Inches(1.5), Inches(0.9), Inches(10), Inches(0.5),
         "Multimodal AI sees the invoice like a human — but faster and more consistent", 16, MEDIUM_GRAY)

# Left side: What the AI sees (invoice image area)
add_shape(slide, Inches(0.5), Inches(1.6), Inches(5.5), Inches(5.5), DARK_BLUE, ACCENT_BLUE)
add_text(slide, Inches(0.5), Inches(1.65), Inches(5.5), Inches(0.5),
         "What the AI Reads from Each Invoice", 18, ACCENT_BLUE, True, PP_ALIGN.CENTER)

fields_left = [
    "Vendor name, address, phone, email",
    "Invoice number and date",
    "Due date and payment terms (Net 30, 2/10 Net 30)",
    "PO number (links to KOJO)",
    "Job name or number (links to Procore)",
    "Line items: description, qty, unit price, amount",
    "Subtotal, tax, shipping, total",
    "Retention / retainage held",
    "AIA document references (G702/G703)",
    "Lien waiver status",
    "Change order references",
    "Remit-to address and banking details",
]

y = 2.2
for field in fields_left:
    add_shape(slide, Inches(0.7), Inches(y), Inches(5.1), Inches(0.4), RGBColor(0x18, 0x30, 0x55))
    add_text(slide, Inches(0.85), Inches(y + 0.02), Inches(4.8), Inches(0.36),
             field, 12, WHITE)
    y += 0.43

# Right side: How it works
add_shape(slide, Inches(6.3), Inches(1.6), Inches(6.5), Inches(5.5), DARK_BLUE, TEAL)
add_text(slide, Inches(6.3), Inches(1.65), Inches(6.5), Inches(0.5),
         "How the AI Engine Works", 18, TEAL, True, PP_ALIGN.CENTER)

ai_steps = [
    ("Claude Sonnet 4.6 (Primary)", "Processes 80% of standard invoices\nFast, accurate, cost-effective at ~$0.01/invoice"),
    ("Claude Opus 4.6 (Complex)", "Handles AIA pay apps, multi-page invoices\nExtended thinking for complex line items"),
    ("Claude Haiku 4.5 (Triage)", "Initial classification and routing\nInstant response, lowest cost"),
    ("Structured Output", "Guaranteed JSON schema — every field maps\ndirectly to your Sage chart of accounts"),
    ("Confidence Scoring", "Every extracted field gets a confidence %\n95%+ = auto-approve, <80% = flag for review"),
    ("Learning Loop", "Corrections feed back into the system\nAccuracy improves with every invoice processed"),
]

y = 2.2
for title, desc in ai_steps:
    add_shape(slide, Inches(6.5), Inches(y), Inches(6.1), Inches(0.82), RGBColor(0x18, 0x30, 0x55))
    add_text(slide, Inches(6.65), Inches(y + 0.02), Inches(5.8), Inches(0.3),
             title, 13, TEAL, True)
    add_text(slide, Inches(6.65), Inches(y + 0.32), Inches(5.8), Inches(0.45),
             desc, 10, RGBColor(0xCC, 0xCC, 0xCC))
    y += 0.87


# ═══════════════════════════════════════════════
# SLIDE 5: Step 3 - Intelligent Job Coding
# ═══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.06), ACCENT_BLUE)

add_step_number(slide, Inches(0.6), Inches(0.3), "3")
add_text(slide, Inches(1.5), Inches(0.3), Inches(10), Inches(0.7),
         "Intelligent Job Coding — 5-Dimensional Auto-Coding", 32, WHITE, True)
add_text(slide, Inches(1.5), Inches(0.9), Inches(10), Inches(0.5),
         "AI predicts all 5 coding dimensions using signals from KOJO, Procore, and invoice history", 16, MEDIUM_GRAY)

# 5 dimensions
dims = [
    ("GL Account", "5000 - Materials", "Mapped from Sage\nchart of accounts", ACCENT_BLUE),
    ("Job", "2401 - Downtown\nMedical Center", "Matched via Procore\njob list", PROCORE_ORANGE),
    ("Phase", "200 - Rough-In", "From Procore project\nphase structure", TEAL),
    ("Cost Code", "22-310 - Plumbing\nFixtures", "Predicted from invoice\nline items + history", GREEN),
    ("Cost Type", "M - Material", "Determined by vendor\ntype and description", PURPLE),
]

x = 0.4
for title, example, source, color in dims:
    add_shape(slide, Inches(x), Inches(1.6), Inches(2.4), Inches(2.5), DARK_BLUE, color)
    add_text(slide, Inches(x), Inches(1.7), Inches(2.4), Inches(0.45),
             title, 16, color, True, PP_ALIGN.CENTER)
    add_text(slide, Inches(x + 0.1), Inches(2.15), Inches(2.2), Inches(0.7),
             example, 14, WHITE, True, PP_ALIGN.CENTER)
    add_text(slide, Inches(x + 0.1), Inches(2.9), Inches(2.2), Inches(0.7),
             source, 11, MEDIUM_GRAY, False, PP_ALIGN.CENTER)
    x += 2.55

# How signals work
add_text(slide, Inches(0.5), Inches(4.3), Inches(12), Inches(0.5),
         "How PayFlow AI Determines the Right Codes:", 20, ACCENT_BLUE, True)

signals = [
    ("PO Number Match (KOJO)", "Invoice references PO #7823 — KOJO says that PO is for Job 2401, Cost Code 22-310. Instant match.", KOJO_ORANGE),
    ("Job Reference (Procore)", "Invoice mentions 'Downtown Medical' — Procore confirms this is Job 2401. Phase and cost codes pulled from project structure.", PROCORE_ORANGE),
    ("Vendor History", "Ferguson Supply has billed to Job 2401 / Code 22-310 for the last 47 invoices. New invoice follows the same pattern.", ACCENT_BLUE),
    ("Material Analysis", "Line items include 'copper fittings' and 'PVC pipe' — AI recognizes plumbing materials, assigns cost type M.", GREEN),
    ("Delivery Address", "Ship-to address matches the Downtown Medical job site in Procore. Confirms job assignment.", TEAL),
    ("Correction Learning", "AP clerk corrected cost code from 22-200 to 22-310 twice — AI now defaults to 22-310 for this vendor/material combo.", PURPLE),
]

y = 4.8
for title, desc, color in signals:
    add_shape(slide, Inches(0.5), Inches(y), Inches(0.12), Inches(0.38), color)
    add_text(slide, Inches(0.8), Inches(y), Inches(3.5), Inches(0.38),
             title, 12, color, True)
    add_text(slide, Inches(4.3), Inches(y), Inches(8.5), Inches(0.38),
             desc, 11, RGBColor(0xCC, 0xCC, 0xCC))
    y += 0.42


# ═══════════════════════════════════════════════
# SLIDE 6: Step 4 - Three-Way Matching (KOJO + Procore)
# ═══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.06), ACCENT_BLUE)

add_step_number(slide, Inches(0.6), Inches(0.3), "4")
add_text(slide, Inches(1.5), Inches(0.3), Inches(10), Inches(0.7),
         "Three-Way Matching — Invoice vs PO vs Budget", 32, WHITE, True)
add_text(slide, Inches(1.5), Inches(0.9), Inches(10), Inches(0.5),
         "Cross-referencing KOJO purchase orders, Procore budgets, and the invoice in real-time", 16, MEDIUM_GRAY)

# Three columns for the match
# Column 1: Invoice
add_shape(slide, Inches(0.5), Inches(1.6), Inches(3.8), Inches(4.8), DARK_BLUE, ACCENT_BLUE)
add_text(slide, Inches(0.5), Inches(1.7), Inches(3.8), Inches(0.5),
         "Invoice (Received)", 20, ACCENT_BLUE, True, PP_ALIGN.CENTER)
invoice_fields = [
    ("Vendor:", "Ferguson Supply"),
    ("Invoice #:", "INV-2026-4891"),
    ("PO Reference:", "#7823"),
    ("Amount:", "$12,400.00"),
    ("Description:", "Copper fittings, PVC pipe"),
    ("Job Ref:", "Downtown Medical"),
    ("Retention:", "10% ($1,240)"),
    ("Net Due:", "$11,160.00"),
]
y = 2.3
for label, value in invoice_fields:
    add_text(slide, Inches(0.7), Inches(y), Inches(1.5), Inches(0.35),
             label, 12, MEDIUM_GRAY, True)
    add_text(slide, Inches(2.2), Inches(y), Inches(2.0), Inches(0.35),
             value, 12, WHITE)
    y += 0.37

# Arrows
add_arrow(slide, Inches(4.5), Inches(3.5), Inches(0.6), Inches(0.35), GREEN)
add_arrow(slide, Inches(8.6), Inches(3.5), Inches(0.6), Inches(0.35), GREEN)

# Column 2: KOJO PO
add_shape(slide, Inches(5.2), Inches(1.6), Inches(3.2), Inches(4.8), DARK_BLUE, KOJO_ORANGE)
add_text(slide, Inches(5.2), Inches(1.7), Inches(3.2), Inches(0.5),
         "KOJO Purchase Order", 18, KOJO_ORANGE, True, PP_ALIGN.CENTER)
po_fields = [
    ("PO #:", "7823"),
    ("Vendor:", "Ferguson Supply"),
    ("Job:", "2401"),
    ("Cost Code:", "22-310"),
    ("PO Amount:", "$15,000.00"),
    ("Previously Billed:", "$2,600.00"),
    ("Remaining:", "$12,400.00"),
    ("Status:", "Delivered"),
]
y = 2.3
for label, value in po_fields:
    add_text(slide, Inches(5.4), Inches(y), Inches(1.4), Inches(0.35),
             label, 12, MEDIUM_GRAY, True)
    add_text(slide, Inches(6.8), Inches(y), Inches(1.5), Inches(0.35),
             value, 12, WHITE)
    y += 0.37

# Column 3: Procore Budget
add_shape(slide, Inches(9.3), Inches(1.6), Inches(3.5), Inches(4.8), DARK_BLUE, PROCORE_ORANGE)
add_text(slide, Inches(9.3), Inches(1.7), Inches(3.5), Inches(0.5),
         "Procore Budget", 18, PROCORE_ORANGE, True, PP_ALIGN.CENTER)
budget_fields = [
    ("Job:", "2401 - Downtown Med"),
    ("Code:", "22-310 Plumb Fixtures"),
    ("Budget:", "$85,000.00"),
    ("Committed:", "$72,000.00"),
    ("Billed to Date:", "$58,200.00"),
    ("This Invoice:", "$12,400.00"),
    ("Remaining:", "$14,400.00"),
    ("Status:", "Within Budget"),
]
y = 2.3
for label, value in budget_fields:
    add_text(slide, Inches(9.5), Inches(y), Inches(1.5), Inches(0.35),
             label, 12, MEDIUM_GRAY, True)
    add_text(slide, Inches(11.0), Inches(y), Inches(1.7), Inches(0.35),
             value, 12, WHITE)
    y += 0.37

# Match results
add_shape(slide, Inches(0.5), Inches(6.6), Inches(12.3), Inches(0.7), DARK_GREEN)
results = [
    ("Vendor Match", 1.5), ("Amount Match", 4.0), ("PO Match", 6.3),
    ("Budget OK", 8.5), ("Delivery Confirmed", 10.8),
]
for label, x in results:
    add_text(slide, Inches(x), Inches(6.65), Inches(2.0), Inches(0.3),
             "  " + label, 13, WHITE, True, PP_ALIGN.LEFT)
    add_shape(slide, Inches(x), Inches(6.72), Inches(0.18), Inches(0.18), GREEN)


# ═══════════════════════════════════════════════
# SLIDE 7: Step 5 - Validation & Fraud Detection
# ═══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.06), ACCENT_BLUE)

add_step_number(slide, Inches(0.6), Inches(0.3), "5")
add_text(slide, Inches(1.5), Inches(0.3), Inches(10), Inches(0.7),
         "Validation & Fraud Detection", 32, WHITE, True)
add_text(slide, Inches(1.5), Inches(0.9), Inches(10), Inches(0.5),
         "AI catches errors and anomalies that humans miss", 16, MEDIUM_GRAY)

# Left column: What we catch
checks = [
    ("Duplicate Invoice Detection", "Same vendor, same amount, same date = flagged.\nAlso catches near-duplicates (resubmitted with new number).", RED),
    ("Math Verification", "Line items x qty = line total? Lines sum to subtotal?\nTax calculated correctly? Total matches?", ORANGE),
    ("Over-Budget Alert", "This invoice would push Cost Code 22-310 over the\nProcore budget. Requires controller approval.", ORANGE),
    ("Over-Billed vs Contract", "Vendor has billed $73,000 against a $72,000 commitment.\nFlags overbilling before payment.", RED),
    ("PO Mismatch", "Invoice says PO #7823 but amount exceeds PO remaining\nbalance. Holds for review.", ORANGE),
    ("Suspicious Patterns", "Round-number invoices, new vendor with high first invoice,\nremit-to address changed recently.", RED),
    ("Retention Verification", "Invoice should hold 10% retention per contract.\nAI verifies retention is calculated correctly.", TEAL),
    ("Lien Waiver Check", "Prior-period lien waiver required before processing\ncurrent payment. Flags if missing.", PURPLE),
]

y = 1.5
for title, desc, color in checks:
    add_shape(slide, Inches(0.5), Inches(y), Inches(8.5), Inches(0.72), DARK_BLUE, color)
    add_text(slide, Inches(0.7), Inches(y + 0.03), Inches(3.0), Inches(0.3),
             title, 13, color, True)
    add_text(slide, Inches(3.8), Inches(y + 0.03), Inches(5.0), Inches(0.65),
             desc, 10, RGBColor(0xCC, 0xCC, 0xCC))
    y += 0.76

# Right side: Stats
add_shape(slide, Inches(9.3), Inches(1.5), Inches(3.5), Inches(5.6), DARK_BLUE, ACCENT_BLUE)
add_text(slide, Inches(9.3), Inches(1.6), Inches(3.5), Inches(0.5),
         "Catch Rate", 20, ACCENT_BLUE, True, PP_ALIGN.CENTER)

stats = [
    ("99.2%", "Duplicate Detection"),
    ("100%", "Math Errors Caught"),
    ("97%", "Over-Budget Flagged\nBefore Approval"),
    ("94%", "Suspicious Invoices\nIdentified"),
    ("<2%", "False Positive Rate"),
    ("$0", "Fraud Losses with\nPayFlow AI"),
]

y = 2.2
for stat, label in stats:
    add_text(slide, Inches(9.5), Inches(y), Inches(1.3), Inches(0.4),
             stat, 22, GREEN if stat not in ["<2%", "$0"] else ACCENT_BLUE, True, PP_ALIGN.RIGHT)
    add_text(slide, Inches(10.85), Inches(y + 0.02), Inches(1.8), Inches(0.5),
             label, 11, MEDIUM_GRAY)
    y += 0.55


# ═══════════════════════════════════════════════
# SLIDE 8: Step 6 - Approval Routing
# ═══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.06), ACCENT_BLUE)

add_step_number(slide, Inches(0.6), Inches(0.3), "6")
add_text(slide, Inches(1.5), Inches(0.3), Inches(10), Inches(0.7),
         "Approval Routing — Right Person, Right Time", 32, WHITE, True)
add_text(slide, Inches(1.5), Inches(0.9), Inches(10), Inches(0.5),
         "Configurable approval chains based on amount, job, trade, and vendor", 16, MEDIUM_GRAY)

# Approval rules table
add_shape(slide, Inches(0.5), Inches(1.5), Inches(12.3), Inches(0.5), ACCENT_BLUE)
headers = [("Rule", 0.6, 2.5), ("Condition", 3.2, 3.0), ("Approver", 6.3, 2.5), ("Action", 8.9, 3.5)]
for text, x, w in headers:
    add_text(slide, Inches(x), Inches(1.53), Inches(w), Inches(0.4),
             text, 14, WHITE, True)

rules = [
    ("Under $5,000", "Standard invoice, PO matched", "AP Clerk auto-approve", "Posts to Sage automatically"),
    ("$5,000 - $25,000", "Within budget, PO matched", "Project Manager", "Email + mobile notification"),
    ("$25,000 - $100,000", "Any invoice this range", "PM + Controller", "Sequential approval chain"),
    ("Over $100,000", "Large invoice or no PO", "PM + Controller + CFO", "Three-level approval required"),
    ("Over Budget", "Exceeds Procore budget", "Controller (mandatory)", "Budget override required"),
    ("No PO Match", "No KOJO PO found", "AP Manager review", "Must create or link PO"),
    ("New Vendor", "First invoice from vendor", "Controller review", "Verify W-9, COI, setup in Sage"),
    ("Change Order", "References a CO", "PM who approved the CO", "Verify CO amount in Procore"),
    ("Retention Release", "Final payment / retention", "PM + Controller + CFO", "Lien waiver required first"),
]

y = 2.1
for i, (rule, condition, approver, action) in enumerate(rules):
    bg_color = DARK_BLUE if i % 2 == 0 else RGBColor(0x15, 0x2D, 0x52)
    add_shape(slide, Inches(0.5), Inches(y), Inches(12.3), Inches(0.48), bg_color)
    add_text(slide, Inches(0.6), Inches(y + 0.05), Inches(2.5), Inches(0.38),
             rule, 12, ORANGE if "Over" in rule or "No PO" in rule or "New" in rule else WHITE, True)
    add_text(slide, Inches(3.2), Inches(y + 0.05), Inches(3.0), Inches(0.38),
             condition, 11, RGBColor(0xCC, 0xCC, 0xCC))
    add_text(slide, Inches(6.3), Inches(y + 0.05), Inches(2.5), Inches(0.38),
             approver, 11, WHITE)
    add_text(slide, Inches(8.9), Inches(y + 0.05), Inches(3.5), Inches(0.38),
             action, 11, MEDIUM_GRAY)
    y += 0.5

# Mobile approval note
add_shape(slide, Inches(0.5), Inches(6.8), Inches(12.3), Inches(0.5), DARK_GREEN)
add_text(slide, Inches(0.5), Inches(6.83), Inches(12.3), Inches(0.4),
         "Approvers can review and approve from email, mobile app, or web dashboard — with full invoice image and AI summary",
         13, WHITE, False, PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════
# SLIDE 9: Step 7 - Sage Integration (Posting)
# ═══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.06), ACCENT_BLUE)

add_step_number(slide, Inches(0.6), Inches(0.3), "7")
add_text(slide, Inches(1.5), Inches(0.3), Inches(10), Inches(0.7),
         "Sage 300 CRE Integration — Approved to Posted", 32, WHITE, True)
add_text(slide, Inches(1.5), Inches(0.9), Inches(10), Inches(0.5),
         "Approved invoices push directly into Sage — no re-keying, no errors", 16, MEDIUM_GRAY)

# What syncs FROM Sage (left)
add_shape(slide, Inches(0.5), Inches(1.5), Inches(5.8), Inches(5.5), DARK_BLUE, SAGE_GREEN)
add_text(slide, Inches(0.5), Inches(1.6), Inches(5.8), Inches(0.5),
         "What We Pull FROM Sage", 20, SAGE_GREEN, True, PP_ALIGN.CENTER)
add_text(slide, Inches(0.5), Inches(2.0), Inches(5.8), Inches(0.3),
         "Keeps PayFlow AI in sync with your accounting system", 12, MEDIUM_GRAY, False, PP_ALIGN.CENTER)

from_sage = [
    ("Chart of Accounts", "Full GL structure so AI codes to valid accounts"),
    ("Vendor Master", "All vendors, payment terms, 1099 status, addresses"),
    ("Job Cost Structure", "Jobs, phases, cost codes, cost types from Sage"),
    ("Open AP", "Existing unpaid invoices to prevent duplicates"),
    ("Payment History", "What's been paid, when, check numbers"),
    ("Tax Codes", "Sales tax, use tax rates by jurisdiction"),
]
y = 2.4
for title, desc in from_sage:
    add_shape(slide, Inches(0.7), Inches(y), Inches(5.4), Inches(0.7), RGBColor(0x18, 0x30, 0x55))
    add_text(slide, Inches(0.85), Inches(y + 0.03), Inches(2.0), Inches(0.3),
             title, 13, SAGE_GREEN, True)
    add_text(slide, Inches(0.85), Inches(y + 0.33), Inches(5.1), Inches(0.3),
             desc, 11, RGBColor(0xCC, 0xCC, 0xCC))
    y += 0.75

# What pushes TO Sage (right)
add_shape(slide, Inches(6.8), Inches(1.5), Inches(6.0), Inches(5.5), DARK_BLUE, ACCENT_BLUE)
add_text(slide, Inches(6.8), Inches(1.6), Inches(6.0), Inches(0.5),
         "What We Push TO Sage", 20, ACCENT_BLUE, True, PP_ALIGN.CENTER)
add_text(slide, Inches(6.8), Inches(2.0), Inches(6.0), Inches(0.3),
         "Approved invoices post with zero manual entry", 12, MEDIUM_GRAY, False, PP_ALIGN.CENTER)

to_sage = [
    ("AP Invoice Entry", "Full invoice header: vendor, date, amount, terms, due date"),
    ("Distribution Lines", "Each line coded to GL/Job/Phase/Cost Code/Cost Type"),
    ("Retention Entry", "Retention held posts to retainage payable automatically"),
    ("Tax Allocation", "Correct tax codes applied per line item"),
    ("Document Attachment", "Original invoice PDF attached to the Sage transaction"),
    ("Approval Trail", "Who approved, when, any notes — full audit history"),
]
y = 2.4
for title, desc in to_sage:
    add_shape(slide, Inches(7.0), Inches(y), Inches(5.6), Inches(0.7), RGBColor(0x18, 0x30, 0x55))
    add_text(slide, Inches(7.15), Inches(y + 0.03), Inches(2.2), Inches(0.3),
             title, 13, ACCENT_BLUE, True)
    add_text(slide, Inches(7.15), Inches(y + 0.33), Inches(5.3), Inches(0.3),
             desc, 11, RGBColor(0xCC, 0xCC, 0xCC))
    y += 0.75

# Bottom note
add_shape(slide, Inches(0.5), Inches(7.15), Inches(12.3), Inches(0.25), RGBColor(0x15, 0x2D, 0x52))
add_text(slide, Inches(0.5), Inches(7.15), Inches(12.3), Inches(0.25),
         "Two-way sync runs every 15 minutes — PayFlow AI always reflects current Sage data",
         11, MEDIUM_GRAY, False, PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════
# SLIDE 10: KOJO Deep Dive
# ═══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.06), KOJO_ORANGE)

add_text(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.7),
         "KOJO Integration — Procurement to Payment", 32, WHITE, True)
add_text(slide, Inches(0.8), Inches(0.9), Inches(11), Inches(0.5),
         "Closing the loop between purchasing and accounts payable", 16, MEDIUM_GRAY)

# Left: What KOJO provides
add_shape(slide, Inches(0.5), Inches(1.5), Inches(5.8), Inches(5.5), DARK_BLUE, KOJO_ORANGE)
add_text(slide, Inches(0.5), Inches(1.6), Inches(5.8), Inches(0.5),
         "What PayFlow AI Gets from KOJO", 18, KOJO_ORANGE, True, PP_ALIGN.CENTER)

kojo_data = [
    ("Purchase Orders", "PO number, vendor, line items, amounts, job assignment"),
    ("Material Orders", "What was ordered, quantities, unit prices"),
    ("Delivery Status", "Confirmed deliveries — was material actually received?"),
    ("Vendor Quotes", "Original quoted prices to compare against invoice"),
    ("Change Orders", "PO modifications, additions, price adjustments"),
    ("Order History", "Full purchasing history per vendor per job"),
    ("Buyout Tracking", "Budget vs committed vs purchased amounts"),
]
y = 2.2
for title, desc in kojo_data:
    add_shape(slide, Inches(0.7), Inches(y), Inches(5.4), Inches(0.65), RGBColor(0x18, 0x30, 0x55))
    add_text(slide, Inches(0.85), Inches(y + 0.02), Inches(2.0), Inches(0.28),
             title, 12, KOJO_ORANGE, True)
    add_text(slide, Inches(0.85), Inches(y + 0.3), Inches(5.1), Inches(0.3),
             desc, 11, RGBColor(0xCC, 0xCC, 0xCC))
    y += 0.7

# Right: What this enables
add_shape(slide, Inches(6.8), Inches(1.5), Inches(6.0), Inches(5.5), DARK_BLUE, GREEN)
add_text(slide, Inches(6.8), Inches(1.6), Inches(6.0), Inches(0.5),
         "What This Enables", 18, GREEN, True, PP_ALIGN.CENTER)

enables = [
    ("Instant PO Matching", "Invoice arrives referencing PO #7823.\nKOJO confirms: vendor, amount, job — matched in seconds."),
    ("Price Verification", "KOJO quote said $14.50/fitting.\nInvoice charges $16.00/fitting.\nFlag: vendor overcharged by 10%."),
    ("Delivery Confirmation", "Invoice says 200 fittings delivered.\nKOJO shows 180 received.\nFlag: billed for 20 undelivered items."),
    ("Remaining Balance", "PO total: $15,000. Previously billed: $2,600.\nThis invoice: $12,400. Remaining: $0.\nAll checks pass."),
    ("Auto-Coding", "PO was coded to Job 2401 / 22-310 in KOJO.\nInvoice inherits the same codes automatically.\nNo manual coding needed."),
]
y = 2.2
for title, desc in enables:
    add_shape(slide, Inches(7.0), Inches(y), Inches(5.6), Inches(0.95), RGBColor(0x18, 0x30, 0x55))
    add_text(slide, Inches(7.15), Inches(y + 0.02), Inches(5.3), Inches(0.28),
             title, 13, GREEN, True)
    add_text(slide, Inches(7.15), Inches(y + 0.3), Inches(5.3), Inches(0.6),
             desc, 10, RGBColor(0xCC, 0xCC, 0xCC))
    y += 1.0


# ═══════════════════════════════════════════════
# SLIDE 11: Procore Deep Dive
# ═══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.06), PROCORE_ORANGE)

add_text(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.7),
         "Procore Integration — Project Intelligence for AP", 32, WHITE, True)
add_text(slide, Inches(0.8), Inches(0.9), Inches(11), Inches(0.5),
         "Every invoice validated against real-time project data", 16, MEDIUM_GRAY)

# Left column
add_shape(slide, Inches(0.5), Inches(1.5), Inches(6.0), Inches(5.5), DARK_BLUE, PROCORE_ORANGE)
add_text(slide, Inches(0.5), Inches(1.6), Inches(6.0), Inches(0.5),
         "What PayFlow AI Gets from Procore", 18, PROCORE_ORANGE, True, PP_ALIGN.CENTER)

procore_data = [
    ("Active Projects", "Job list, addresses, status, project managers"),
    ("Budget Data", "Original budget, revised budget, approved changes per cost code"),
    ("Commitments", "Subcontracts, POs — contracted amounts and billed-to-date"),
    ("Change Orders", "Approved COs with amounts, tied to specific cost codes"),
    ("Cost Codes", "Full cost code structure unique to each project"),
    ("Subcontractors", "Sub list, contact info, trade, compliance status"),
    ("Daily Logs", "Field activity records — verify labor invoices against actual work"),
    ("RFIs & Submittals", "Context for material changes that affect invoicing"),
]
y = 2.1
for title, desc in procore_data:
    add_shape(slide, Inches(0.7), Inches(y), Inches(5.6), Inches(0.6), RGBColor(0x18, 0x30, 0x55))
    add_text(slide, Inches(0.85), Inches(y + 0.02), Inches(2.0), Inches(0.25),
             title, 12, PROCORE_ORANGE, True)
    add_text(slide, Inches(0.85), Inches(y + 0.28), Inches(5.3), Inches(0.28),
             desc, 10, RGBColor(0xCC, 0xCC, 0xCC))
    y += 0.63

# Right column
add_shape(slide, Inches(6.8), Inches(1.5), Inches(6.0), Inches(5.5), DARK_BLUE, GREEN)
add_text(slide, Inches(6.8), Inches(1.6), Inches(6.0), Inches(0.5),
         "What This Enables", 18, GREEN, True, PP_ALIGN.CENTER)

procore_enables = [
    ("Budget Guardian", "Every invoice checked against Procore budget.\nAlert before a cost code goes over budget — not after."),
    ("Contract Compliance", "Sub billed $73K against $72K commitment.\nFlagged: overbilling by $1,000 before payment."),
    ("Change Order Verification", "Invoice references CO #14.\nProcore confirms CO #14 approved for $8,500.\nAmount validated automatically."),
    ("Over/Under Billing", "Compare % complete in Procore vs % billed.\nCatch front-loading before it becomes a cash problem."),
    ("Job Profitability", "Real-time cost vs budget by job.\nAlert PMs when a job is trending over budget."),
    ("Compliance Check", "Sub's COI expired 3 days ago.\nHold payment until updated certificate received."),
]
y = 2.1
for title, desc in procore_enables:
    add_shape(slide, Inches(7.0), Inches(y), Inches(5.6), Inches(0.82), RGBColor(0x18, 0x30, 0x55))
    add_text(slide, Inches(7.15), Inches(y + 0.02), Inches(5.3), Inches(0.25),
             title, 13, GREEN, True)
    add_text(slide, Inches(7.15), Inches(y + 0.28), Inches(5.3), Inches(0.5),
             desc, 10, RGBColor(0xCC, 0xCC, 0xCC))
    y += 0.85


# ═══════════════════════════════════════════════
# SLIDE 12: Other Integrations You're Not Thinking Of
# ═══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.06), PURPLE)

add_text(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.7),
         "Beyond KOJO, Procore & Sage — Other Integrations", 30, WHITE, True)
add_text(slide, Inches(0.8), Inches(0.9), Inches(11), Inches(0.5),
         "Programs you may not be thinking about that add massive value", 16, MEDIUM_GRAY)

integrations_extra = [
    ("Textura / Oracle CPX", "Payment Management",
     "GC payment tracking, compliance docs,\nlien waiver exchange, pay-when-paid tracking",
     "Know when YOU'LL get paid so you know when to pay subs", ACCENT_BLUE),
    ("Viewpoint Vista / Spectrum", "Alternative ERP",
     "If a customer uses Viewpoint instead of Sage,\nsame integration model applies",
     "Expands addressable market beyond Sage shops", TEAL),
    ("CMiC", "Enterprise ERP",
     "Large ENR-ranked contractors use CMiC,\nfull job cost and AP integration",
     "Opens the door to $500M+ contractors", GREEN),
    ("DocuSign / PandaDoc", "E-Signatures",
     "Digital lien waiver signatures,\nsubcontractor agreement execution",
     "Eliminates paper chase for compliance docs", PURPLE),
    ("QuickBooks Online", "Small Contractor ERP",
     "Smaller subs and specialty contractors,\nsimpler GL and AP structure",
     "Entry-level tier for smaller companies", RGBColor(0x2C, 0xA0, 0x1C)),
    ("myCOI / PINS", "Insurance Tracking",
     "Auto-verify COI status before payment,\nalert on expiring policies",
     "Never pay a sub with expired insurance", ORANGE),
    ("Microsoft 365 / Google", "Email & Calendar",
     "Email ingestion, calendar-based approvals,\nTeams/Slack notifications",
     "Meet approvers where they already work", ACCENT_BLUE),
    ("Plangrid / Autodesk Build", "Field Management",
     "Daily reports, punch lists, inspections\nverify work completion before payment",
     "Field verification of invoiced work", PROCORE_ORANGE),
    ("LevelSet / Zlien", "Lien Waiver Mgmt",
     "Automated preliminary notices, lien tracking,\nwaiver exchange platform",
     "Complete lien risk management", RED),
    ("Banking APIs (Plaid)", "Payment Execution",
     "ACH payments, positive pay files,\nreal-time payment status",
     "End-to-end: invoice to actual payment", DARK_GREEN),
]

# Display in two columns
col1 = integrations_extra[:5]
col2 = integrations_extra[5:]

y = 1.5
for title, category, desc, value, color in col1:
    add_shape(slide, Inches(0.3), Inches(y), Inches(6.2), Inches(1.0), DARK_BLUE, color)
    add_text(slide, Inches(0.5), Inches(y + 0.03), Inches(2.3), Inches(0.25),
             title, 13, color, True)
    add_text(slide, Inches(2.8), Inches(y + 0.03), Inches(1.5), Inches(0.25),
             category, 10, MEDIUM_GRAY)
    add_text(slide, Inches(0.5), Inches(y + 0.3), Inches(2.8), Inches(0.6),
             desc, 9, RGBColor(0xBB, 0xBB, 0xBB))
    add_text(slide, Inches(3.4), Inches(y + 0.3), Inches(3.0), Inches(0.6),
             value, 10, WHITE)
    y += 1.08

y = 1.5
for title, category, desc, value, color in col2:
    add_shape(slide, Inches(6.8), Inches(y), Inches(6.2), Inches(1.0), DARK_BLUE, color)
    add_text(slide, Inches(7.0), Inches(y + 0.03), Inches(2.5), Inches(0.25),
             title, 13, color, True)
    add_text(slide, Inches(9.5), Inches(y + 0.03), Inches(1.5), Inches(0.25),
             category, 10, MEDIUM_GRAY)
    add_text(slide, Inches(7.0), Inches(y + 0.3), Inches(2.8), Inches(0.6),
             desc, 9, RGBColor(0xBB, 0xBB, 0xBB))
    add_text(slide, Inches(9.9), Inches(y + 0.3), Inches(3.0), Inches(0.6),
             value, 10, WHITE)
    y += 1.08

# Bottom note
add_shape(slide, Inches(0.3), Inches(7.0), Inches(12.7), Inches(0.35), DARK_BLUE, PURPLE)
add_text(slide, Inches(0.3), Inches(7.02), Inches(12.7), Inches(0.3),
         "Phase 1: KOJO + Procore + Sage   |   Phase 2: Textura + DocuSign + myCOI   |   Phase 3: Viewpoint + CMiC + Banking",
         12, PURPLE, False, PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════
# SLIDE 13: Construction-Specific Features
# ═══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.06), ACCENT_BLUE)

add_text(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.7),
         "Construction-Specific AP Features", 32, WHITE, True)
add_text(slide, Inches(0.8), Inches(0.9), Inches(11), Inches(0.5),
         "Features that generic AP tools like Stampli and Bill.com simply don't have", 16, MEDIUM_GRAY)

features = [
    ("AIA G702/G703 Processing",
     "AI reads progress billing applications — the most complex documents in construction AP. "
     "Extracts schedule of values, % complete per line, retention, stored materials, change orders. "
     "Compares current application against prior period. Flags overbilling by line item.",
     ACCENT_BLUE),
    ("Lien Waiver Management",
     "Tracks all 4 types: conditional/unconditional, progress/final. Blocks payment if prior-period "
     "waiver is missing. Auto-generates waiver requests. Ensures you never pay without proper waiver.",
     GREEN),
    ("Retention / Retainage Tracking",
     "Tracks retention held per vendor per job. Calculates 10% (or custom %) automatically on each draw. "
     "Separate aging for retention payable. Alerts when substantial completion triggers retention release.",
     ORANGE),
    ("Certified Payroll & Prevailing Wage",
     "Validates labor invoices against prevailing wage rates by trade and jurisdiction. "
     "Cross-references certified payroll reports. Critical for government/public works projects.",
     PURPLE),
    ("Compliance Dashboard",
     "Real-time view: which subs have current COI, W-9, business license, union status. "
     "Blocks payment to non-compliant vendors. Auto-alerts subs 30 days before expiration.",
     TEAL),
    ("Change Order Workflow",
     "Invoices referencing change orders validated against approved COs in Procore. "
     "Tracks: original contract + approved changes = current contract value. Prevents billing beyond scope.",
     RED),
]

y = 1.5
for title, desc, color in features:
    add_shape(slide, Inches(0.5), Inches(y), Inches(12.3), Inches(0.9), DARK_BLUE, color)
    add_text(slide, Inches(0.7), Inches(y + 0.05), Inches(3.0), Inches(0.3),
             title, 15, color, True)
    add_text(slide, Inches(3.8), Inches(y + 0.05), Inches(8.8), Inches(0.8),
             desc, 11, RGBColor(0xCC, 0xCC, 0xCC))
    y += 0.96


# ═══════════════════════════════════════════════
# SLIDE 14: Real Example — Full Invoice Journey
# ═══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.06), GREEN)

add_text(slide, Inches(0.8), Inches(0.15), Inches(11), Inches(0.6),
         "Real Example: One Invoice, Start to Finish", 30, WHITE, True)
add_text(slide, Inches(0.8), Inches(0.7), Inches(11), Inches(0.4),
         "Ferguson Supply invoice for copper fittings — Downtown Medical Center project", 14, MEDIUM_GRAY)

journey_steps = [
    ("9:02 AM", "Invoice Arrives",
     "Ferguson emails invoice #INV-2026-4891 for $12,400\nto invoices@markiii.payflowai.com",
     ACCENT_BLUE),
    ("9:02 AM", "AI Extraction (3 sec)",
     "Claude reads the invoice image. Extracts: vendor, amount,\nPO #7823, 14 line items, 10% retention, Net 30 terms.",
     RGBColor(0x00, 0x96, 0xC7)),
    ("9:02 AM", "KOJO PO Match",
     "PO #7823 found in KOJO: Ferguson, Job 2401, $15K total,\n$2,600 previously billed. $12,400 remaining = exact match.",
     KOJO_ORANGE),
    ("9:02 AM", "Procore Budget Check",
     "Job 2401, Code 22-310: $85K budget, $58.2K billed.\nThis invoice brings it to $70.6K. Within budget.",
     PROCORE_ORANGE),
    ("9:02 AM", "Auto-Coded (5D)",
     "GL: 5000 | Job: 2401 | Phase: 200 | Code: 22-310 | Type: M\nConfidence: 98% — PO match + vendor history.",
     TEAL),
    ("9:02 AM", "Validation Passed",
     "No duplicates. Math verified. Within budget. Within PO.\nRetention calculated correctly. No flags.",
     GREEN),
    ("9:02 AM", "Routed to PM",
     "$12,400 > $5K threshold. Sent to Project Manager (Jake)\nfor approval. Email + mobile push notification.",
     PURPLE),
    ("9:14 AM", "PM Approves (Mobile)",
     "Jake reviews on his phone at the job site. Sees invoice\nimage, AI summary, PO match. Taps 'Approve'.",
     GREEN),
    ("9:15 AM", "Posted to Sage",
     "Invoice posted to Sage AP: Vendor 4021, Job 2401-200-22310-M.\nRetention of $1,240 posted to retainage payable. PDF attached.",
     SAGE_GREEN),
    ("9:15 AM", "Done",
     "Total time: 13 minutes (12 waiting for PM).\nAI processing time: 8 seconds. Manual effort: 0 minutes.",
     GREEN),
]

y = 1.15
for time, title, desc, color in journey_steps:
    # Time column
    add_text(slide, Inches(0.3), Inches(y + 0.03), Inches(0.9), Inches(0.25),
             time, 10, MEDIUM_GRAY, False, PP_ALIGN.RIGHT)
    # Dot
    add_shape(slide, Inches(1.3), Inches(y + 0.1), Inches(0.15), Inches(0.15), color)
    # Vertical line (except last)
    if title != "Done":
        add_shape(slide, Inches(1.355), Inches(y + 0.27), Inches(0.03), Inches(0.35), RGBColor(0x30, 0x40, 0x55))
    # Content
    add_text(slide, Inches(1.6), Inches(y), Inches(2.3), Inches(0.28),
             title, 13, color, True)
    add_text(slide, Inches(3.9), Inches(y), Inches(9.0), Inches(0.5),
             desc, 10, RGBColor(0xCC, 0xCC, 0xCC))
    y += 0.57

# Bottom highlight
add_shape(slide, Inches(0.5), Inches(7.0), Inches(12.3), Inches(0.35), DARK_GREEN)
add_text(slide, Inches(0.5), Inches(7.02), Inches(12.3), Inches(0.3),
         "Old way: 25 minutes of manual data entry + chasing approvals for days  |  PayFlow AI: 13 minutes, zero manual entry",
         12, WHITE, True, PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════
# SLIDE 15: Shadow Mode — Risk-Free Onboarding
# ═══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.06), TEAL)

add_text(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.7),
         "Shadow Mode — Risk-Free Onboarding", 32, WHITE, True)
add_text(slide, Inches(0.8), Inches(0.9), Inches(11), Inches(0.5),
         "AI processes alongside your team without posting anything — prove accuracy before going live", 16, MEDIUM_GRAY)

# How shadow mode works
add_shape(slide, Inches(0.5), Inches(1.5), Inches(6.0), Inches(5.0), DARK_BLUE, TEAL)
add_text(slide, Inches(0.5), Inches(1.6), Inches(6.0), Inches(0.5),
         "How It Works", 20, TEAL, True, PP_ALIGN.CENTER)

shadow_steps = [
    ("Week 1-2: Parallel Processing",
     "Your AP team works normally. PayFlow AI processes the same\ninvoices in the background. Nothing posts to Sage."),
    ("Week 2-3: Compare Results",
     "Side-by-side comparison: what your team coded vs what AI coded.\nMeasure accuracy by field, by vendor, by job."),
    ("Week 3-4: Accuracy Report",
     "Detailed report: AI was right 94% of the time on GL coding,\n97% on job assignment, 99% on amounts."),
    ("Week 4: Fine-Tune",
     "Adjust AI prompts based on Mark III's specific patterns.\nTrain on corrections. Re-run shadow mode."),
    ("Week 5: Go Live",
     "Confidence proven. Switch to live mode.\nAI codes, team reviews, approved invoices post to Sage."),
]

y = 2.1
for title, desc in shadow_steps:
    add_shape(slide, Inches(0.7), Inches(y), Inches(5.6), Inches(0.85), RGBColor(0x18, 0x30, 0x55))
    add_text(slide, Inches(0.85), Inches(y + 0.02), Inches(5.3), Inches(0.28),
             title, 13, TEAL, True)
    add_text(slide, Inches(0.85), Inches(y + 0.32), Inches(5.3), Inches(0.5),
             desc, 10, RGBColor(0xCC, 0xCC, 0xCC))
    y += 0.9

# Right: Why this matters
add_shape(slide, Inches(6.8), Inches(1.5), Inches(6.0), Inches(5.0), DARK_BLUE, GREEN)
add_text(slide, Inches(6.8), Inches(1.6), Inches(6.0), Inches(0.5),
         "Why This Matters to a CFO", 20, GREEN, True, PP_ALIGN.CENTER)

cfo_reasons = [
    ("Zero Risk", "Nothing changes in your current workflow until\nyou've seen proof the AI works."),
    ("Measurable Accuracy", "Hard numbers, not promises. See exactly\nhow AI performs on YOUR invoices."),
    ("No Disruption", "AP team keeps working normally during shadow.\nNo training, no new processes yet."),
    ("Fast Time-to-Value", "4-5 weeks from start to live.\nNot a 6-month implementation nightmare."),
    ("Confidence", "Your team trusts the AI because they've\nseen it work on their own data."),
]

y = 2.1
for title, desc in cfo_reasons:
    add_shape(slide, Inches(7.0), Inches(y), Inches(5.6), Inches(0.85), RGBColor(0x18, 0x30, 0x55))
    add_text(slide, Inches(7.15), Inches(y + 0.02), Inches(5.3), Inches(0.28),
             title, 14, GREEN, True)
    add_text(slide, Inches(7.15), Inches(y + 0.32), Inches(5.3), Inches(0.5),
             desc, 11, RGBColor(0xCC, 0xCC, 0xCC))
    y += 0.9


# ═══════════════════════════════════════════════
# SLIDE 16: Data Flow Architecture
# ═══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.06), ACCENT_BLUE)

add_text(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.7),
         "Data Flow Architecture", 32, WHITE, True)
add_text(slide, Inches(0.8), Inches(0.9), Inches(11), Inches(0.5),
         "How data moves between all connected systems", 16, MEDIUM_GRAY)

# Center: PayFlow AI
add_shape(slide, Inches(4.5), Inches(2.8), Inches(4.3), Inches(2.2), DARK_BLUE, ACCENT_BLUE)
add_text(slide, Inches(4.5), Inches(3.0), Inches(4.3), Inches(0.5),
         "PayFlow AI", 24, ACCENT_BLUE, True, PP_ALIGN.CENTER)
add_text(slide, Inches(4.5), Inches(3.5), Inches(4.3), Inches(1.2),
         "Claude AI Engine\nInvoice Processing\nMatching & Coding\nValidation & Routing\nReporting & Analytics",
         11, RGBColor(0xCC, 0xCC, 0xCC), False, PP_ALIGN.CENTER)

# KOJO - top left
add_shape(slide, Inches(0.5), Inches(1.5), Inches(3.0), Inches(1.3), DARK_BLUE, KOJO_ORANGE)
add_text(slide, Inches(0.5), Inches(1.6), Inches(3.0), Inches(0.35),
         "KOJO", 18, KOJO_ORANGE, True, PP_ALIGN.CENTER)
add_text(slide, Inches(0.5), Inches(1.95), Inches(3.0), Inches(0.7),
         "POs, Materials, Quotes\nDelivery Confirmations",
         11, RGBColor(0xBB, 0xBB, 0xBB), False, PP_ALIGN.CENTER)
add_arrow(slide, Inches(3.6), Inches(2.0), Inches(0.8), Inches(0.3), KOJO_ORANGE)

# Procore - top right
add_shape(slide, Inches(9.8), Inches(1.5), Inches(3.0), Inches(1.3), DARK_BLUE, PROCORE_ORANGE)
add_text(slide, Inches(9.8), Inches(1.6), Inches(3.0), Inches(0.35),
         "Procore", 18, PROCORE_ORANGE, True, PP_ALIGN.CENTER)
add_text(slide, Inches(9.8), Inches(1.95), Inches(3.0), Inches(0.7),
         "Jobs, Budgets, COs\nCommitments, Daily Logs",
         11, RGBColor(0xBB, 0xBB, 0xBB), False, PP_ALIGN.CENTER)
# Left arrow (data flows in)
shape = slide.shapes.add_shape(MSO_SHAPE.LEFT_ARROW, Inches(8.9), Inches(2.0), Inches(0.8), Inches(0.3))
shape.fill.solid()
shape.fill.fore_color.rgb = PROCORE_ORANGE
shape.line.fill.background()

# Sage - bottom center
add_shape(slide, Inches(4.5), Inches(5.5), Inches(4.3), Inches(1.3), DARK_BLUE, SAGE_GREEN)
add_text(slide, Inches(4.5), Inches(5.6), Inches(4.3), Inches(0.35),
         "Sage 300 CRE", 18, SAGE_GREEN, True, PP_ALIGN.CENTER)
add_text(slide, Inches(4.5), Inches(5.95), Inches(4.3), Inches(0.7),
         "GL, AP, Vendors, Job Cost\nPayments, Tax, Reporting",
         11, RGBColor(0xBB, 0xBB, 0xBB), False, PP_ALIGN.CENTER)
add_down_arrow(slide, Inches(6.4), Inches(5.1), Inches(0.35), Inches(0.35), SAGE_GREEN)
# Up arrow too (bidirectional)
shape = slide.shapes.add_shape(MSO_SHAPE.UP_ARROW, Inches(6.9), Inches(5.1), Inches(0.35), Inches(0.35))
shape.fill.solid()
shape.fill.fore_color.rgb = SAGE_GREEN
shape.line.fill.background()

# Email/Scanner - left middle
add_shape(slide, Inches(0.5), Inches(3.3), Inches(3.0), Inches(1.0), DARK_BLUE, ACCENT_BLUE)
add_text(slide, Inches(0.5), Inches(3.35), Inches(3.0), Inches(0.35),
         "Email / Upload / Scanner", 14, ACCENT_BLUE, True, PP_ALIGN.CENTER)
add_text(slide, Inches(0.5), Inches(3.7), Inches(3.0), Inches(0.5),
         "Invoice Intake", 11, RGBColor(0xBB, 0xBB, 0xBB), False, PP_ALIGN.CENTER)
add_arrow(slide, Inches(3.6), Inches(3.65), Inches(0.8), Inches(0.3), ACCENT_BLUE)

# Approvers - right middle
add_shape(slide, Inches(9.8), Inches(3.3), Inches(3.0), Inches(1.0), DARK_BLUE, PURPLE)
add_text(slide, Inches(9.8), Inches(3.35), Inches(3.0), Inches(0.35),
         "Approvers", 14, PURPLE, True, PP_ALIGN.CENTER)
add_text(slide, Inches(9.8), Inches(3.7), Inches(3.0), Inches(0.5),
         "PM, Controller, CFO", 11, RGBColor(0xBB, 0xBB, 0xBB), False, PP_ALIGN.CENTER)
add_arrow(slide, Inches(8.9), Inches(3.65), Inches(0.8), Inches(0.3), PURPLE)

# Compliance tools - bottom left
add_shape(slide, Inches(0.5), Inches(5.5), Inches(3.0), Inches(1.3), DARK_BLUE, ORANGE)
add_text(slide, Inches(0.5), Inches(5.6), Inches(3.0), Inches(0.35),
         "Compliance Tools", 14, ORANGE, True, PP_ALIGN.CENTER)
add_text(slide, Inches(0.5), Inches(5.95), Inches(3.0), Inches(0.7),
         "myCOI, DocuSign\nLevelSet, Banking APIs",
         11, RGBColor(0xBB, 0xBB, 0xBB), False, PP_ALIGN.CENTER)

# Banking - bottom right
add_shape(slide, Inches(9.8), Inches(5.5), Inches(3.0), Inches(1.3), DARK_BLUE, DARK_GREEN)
add_text(slide, Inches(9.8), Inches(5.6), Inches(3.0), Inches(0.35),
         "Banking / Payments", 14, DARK_GREEN, True, PP_ALIGN.CENTER)
add_text(slide, Inches(9.8), Inches(5.95), Inches(3.0), Inches(0.7),
         "ACH, Checks, Positive Pay\nPayment Execution",
         11, RGBColor(0xBB, 0xBB, 0xBB), False, PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════
# SLIDE 17: Security & Compliance
# ═══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.06), GREEN)

add_text(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.7),
         "Security & Data Protection", 32, WHITE, True)
add_text(slide, Inches(0.8), Inches(0.9), Inches(11), Inches(0.5),
         "Enterprise-grade security for your financial data", 16, MEDIUM_GRAY)

security_items = [
    ("SOC 2 Type II Compliant", "Annual audit of security controls. Industry standard for financial SaaS.", GREEN),
    ("256-bit AES Encryption", "All data encrypted at rest. TLS 1.3 in transit. Zero plain-text storage.", ACCENT_BLUE),
    ("Role-Based Access Control", "Users only see what they need. AP clerk vs PM vs Controller vs CFO.", TEAL),
    ("Complete Audit Trail", "Every action logged: who viewed, edited, approved, at what time.", PURPLE),
    ("Procore Data Compliance", "We do NOT train AI on Procore data. Used only for matching and display.", PROCORE_ORANGE),
    ("Multi-Tenant Isolation", "Each customer's data is completely isolated. No cross-tenant access.", ACCENT_BLUE),
    ("Automatic Backups", "Daily encrypted backups with 90-day retention. Point-in-time recovery.", GREEN),
    ("API Security", "OAuth 2.0 for all integrations. API keys rotated regularly. Rate limiting.", ORANGE),
]

y = 1.5
for i, (title, desc, color) in enumerate(security_items):
    col = 0 if i < 4 else 1
    row = i if i < 4 else i - 4
    x = 0.5 if col == 0 else 6.8
    yy = 1.5 + row * 1.3
    add_shape(slide, Inches(x), Inches(yy), Inches(6.0), Inches(1.1), DARK_BLUE, color)
    add_text(slide, Inches(x + 0.2), Inches(yy + 0.08), Inches(5.6), Inches(0.35),
             title, 16, color, True)
    add_text(slide, Inches(x + 0.2), Inches(yy + 0.45), Inches(5.6), Inches(0.5),
             desc, 12, RGBColor(0xCC, 0xCC, 0xCC))

# Bottom
add_shape(slide, Inches(0.5), Inches(6.9), Inches(12.3), Inches(0.4), DARK_GREEN)
add_text(slide, Inches(0.5), Inches(6.92), Inches(12.3), Inches(0.35),
         "Your financial data is treated with the same security as a bank — because that's the standard construction CFOs expect",
         12, WHITE, False, PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════
# SLIDE 18: Summary — The Complete Picture
# ═══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.06), ACCENT_BLUE)

add_text(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.7),
         "The Complete Picture", 36, WHITE, True)
add_text(slide, Inches(0.8), Inches(0.9), Inches(11), Inches(0.5),
         "From invoice receipt to payment — fully automated, fully integrated, construction-specific", 16, MEDIUM_GRAY)

# Summary flow
summary_steps = [
    ("1", "RECEIVE", "Email, upload,\nscan, KOJO", ACCENT_BLUE),
    ("2", "EXTRACT", "Claude AI reads\nevery field", RGBColor(0x00, 0x96, 0xC7)),
    ("3", "CODE", "5D auto-coding\nfrom KOJO/Procore", TEAL),
    ("4", "MATCH", "3-way: Invoice\nvs PO vs Budget", GREEN),
    ("5", "VALIDATE", "Duplicates, math,\nfraud, compliance", ORANGE),
    ("6", "APPROVE", "Smart routing\nto right person", PURPLE),
    ("7", "POST", "Push to Sage\nfor payment", SAGE_GREEN),
]

x = 0.4
for num, title, desc, color in summary_steps:
    add_shape(slide, Inches(x), Inches(1.5), Inches(1.65), Inches(2.2), color)
    add_text(slide, Inches(x), Inches(1.55), Inches(1.65), Inches(0.5),
             num, 36, WHITE, True, PP_ALIGN.CENTER)
    add_text(slide, Inches(x), Inches(2.05), Inches(1.65), Inches(0.4),
             title, 14, WHITE, True, PP_ALIGN.CENTER)
    add_text(slide, Inches(x + 0.05), Inches(2.5), Inches(1.55), Inches(0.8),
             desc, 11, RGBColor(0xEE, 0xEE, 0xEE), False, PP_ALIGN.CENTER)
    if x < 11.0:
        add_arrow(slide, Inches(x + 1.7), Inches(2.45), Inches(0.2), Inches(0.2), RGBColor(0x30, 0x40, 0x55))
    x += 1.82

# Key metrics
add_text(slide, Inches(0.5), Inches(4.0), Inches(12), Inches(0.5),
         "By the Numbers", 22, ACCENT_BLUE, True)

metrics = [
    ("3,000", "Invoices/month\nprocessed", ACCENT_BLUE),
    ("8 sec", "Average AI\nprocessing time", TEAL),
    ("95%+", "Auto-coded\ncorrectly", GREEN),
    ("80%", "Require zero\nmanual touch", GREEN),
    ("$0.06", "Cost per\ninvoice", ACCENT_BLUE),
    ("13 min", "Avg receipt\nto approval", PURPLE),
    ("$462K", "Annual AP\ncost savings", GREEN),
]

x = 0.4
for value, label, color in metrics:
    add_shape(slide, Inches(x), Inches(4.5), Inches(1.65), Inches(1.5), DARK_BLUE, color)
    add_text(slide, Inches(x), Inches(4.6), Inches(1.65), Inches(0.5),
             value, 28, color, True, PP_ALIGN.CENTER)
    add_text(slide, Inches(x + 0.05), Inches(5.1), Inches(1.55), Inches(0.7),
             label, 12, RGBColor(0xCC, 0xCC, 0xCC), False, PP_ALIGN.CENTER)
    x += 1.82

# Integration bar at bottom
add_shape(slide, Inches(0.5), Inches(6.3), Inches(12.3), Inches(1.0), DARK_BLUE, ACCENT_BLUE)
add_text(slide, Inches(0.5), Inches(6.35), Inches(12.3), Inches(0.35),
         "Fully Integrated Ecosystem", 16, ACCENT_BLUE, True, PP_ALIGN.CENTER)

integrations_final = [
    ("KOJO", KOJO_ORANGE), ("Procore", PROCORE_ORANGE), ("Sage 300 CRE", SAGE_GREEN),
    ("Textura", ACCENT_BLUE), ("DocuSign", PURPLE), ("myCOI", ORANGE),
    ("Viewpoint", TEAL), ("Banking APIs", DARK_GREEN),
]
x = 1.0
for name, color in integrations_final:
    add_shape(slide, Inches(x), Inches(6.75), Inches(1.3), Inches(0.4), color)
    add_text(slide, Inches(x), Inches(6.77), Inches(1.3), Inches(0.35),
             name, 10, WHITE, True, PP_ALIGN.CENTER)
    x += 1.45


# ═══════════════════════════════════════════════
# SAVE
# ═══════════════════════════════════════════════
output_path = "C:/users/25badmin/projects/accounts-payable-research/PayFlow_AI_How_It_Works.pptx"
prs.save(output_path)
print(f"Saved: {output_path}")
