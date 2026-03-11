from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# Colors
NAVY = RGBColor(0x1a, 0x36, 0x5d)
BLUE = RGBColor(0x2b, 0x6c, 0xb0)
LIGHT_BLUE = RGBColor(0xeb, 0xf4, 0xff)
GREEN = RGBColor(0x38, 0xa1, 0x69)
LIGHT_GREEN = RGBColor(0xc6, 0xf6, 0xd5)
RED = RGBColor(0xe5, 0x3e, 0x3e)
WHITE = RGBColor(0xff, 0xff, 0xff)
DARK = RGBColor(0x1a, 0x20, 0x2c)
GRAY = RGBColor(0x4a, 0x55, 0x68)
LIGHT_GRAY = RGBColor(0xf7, 0xfa, 0xfc)
MEDIUM_GRAY = RGBColor(0xe2, 0xe8, 0xf0)
ORANGE = RGBColor(0xdd, 0x6b, 0x20)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

def add_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_shape(slide, left, top, width, height, fill_color, border_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape

def add_rect(slide, left, top, width, height, fill_color):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape

def add_text_box(slide, left, top, width, height, text, font_size=18, color=DARK, bold=False, alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_bullet_list(slide, left, top, width, height, items, font_size=16, color=DARK):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = "Calibri"
        p.space_after = Pt(8)
        p.level = 0
    return txBox

def add_stat_card(slide, left, top, width, height, label, value, sublabel="", value_color=NAVY):
    card = add_shape(slide, left, top, width, height, WHITE, MEDIUM_GRAY)
    add_text_box(slide, left + Inches(0.2), top + Inches(0.15), width - Inches(0.4), Inches(0.3),
                 label, font_size=11, color=GRAY, bold=False)
    add_text_box(slide, left + Inches(0.2), top + Inches(0.45), width - Inches(0.4), Inches(0.5),
                 value, font_size=32, color=value_color, bold=True)
    if sublabel:
        add_text_box(slide, left + Inches(0.2), top + Inches(1.0), width - Inches(0.4), Inches(0.3),
                     sublabel, font_size=10, color=GRAY)

# ============================================================
# SLIDE 1: TITLE
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
add_bg(slide, NAVY)

# Accent bar
add_rect(slide, Inches(0), Inches(3.4), Inches(13.333), Inches(0.06), BLUE)

add_text_box(slide, Inches(1), Inches(1.5), Inches(11), Inches(1),
             "PayFlow AI", font_size=54, color=WHITE, bold=True)
add_text_box(slide, Inches(1), Inches(2.3), Inches(11), Inches(0.8),
             "AI-Powered Accounts Payable Automation for Construction",
             font_size=28, color=RGBColor(0xa0, 0xce, 0xf0), bold=False)

add_text_box(slide, Inches(1), Inches(3.8), Inches(11), Inches(0.5),
             "Eliminate manual invoice processing. Save hundreds of thousands per year.",
             font_size=18, color=RGBColor(0xcb, 0xd5, 0xe0))
add_text_box(slide, Inches(1), Inches(5.5), Inches(11), Inches(0.4),
             "Confidential Presentation  |  2026",
             font_size=14, color=RGBColor(0x71, 0x80, 0x96))

# ============================================================
# SLIDE 2: THE PROBLEM
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.08), NAVY)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.6),
             "The Problem: Manual AP Is Costing You a Fortune",
             font_size=32, color=NAVY, bold=True)

# Stat cards row
add_stat_card(slide, Inches(0.8), Inches(1.3), Inches(2.7), Inches(1.3),
              "COST PER INVOICE", "$15-$25", "Manual processing (construction)", RED)
add_stat_card(slide, Inches(3.8), Inches(1.3), Inches(2.7), Inches(1.3),
              "PROCESSING TIME", "14.6 days", "Average receipt to payment", RED)
add_stat_card(slide, Inches(6.8), Inches(1.3), Inches(2.7), Inches(1.3),
              "ERROR RATE", "2-5%", "Manual data entry errors", RED)
add_stat_card(slide, Inches(9.8), Inches(1.3), Inches(2.7), Inches(1.3),
              "MISSED DISCOUNTS", "$$$", "2/10 Net 30 terms wasted", RED)

add_text_box(slide, Inches(0.8), Inches(3.0), Inches(11.5), Inches(0.5),
             "Construction AP is even harder than normal AP:",
             font_size=20, color=NAVY, bold=True)

bullets = [
    "AIA progress billing (G702/G703) takes 30+ minutes to review manually per application",
    "Lien waivers must be collected and verified before every payment — or risk a lien on the property",
    "Every invoice must be coded to the correct job, phase, cost code, and cost type (5 dimensions, not 1)",
    "Retention tracking across dozens of active projects with different retainage rates",
    "Change orders constantly modify contract values — easy to overpay without catching it",
    "Compliance documents (COIs, W-9s, sub agreements) expire and must be tracked per vendor",
]
add_bullet_list(slide, Inches(1.0), Inches(3.5), Inches(11), Inches(3.5), bullets, font_size=16, color=GRAY)

# ============================================================
# SLIDE 3: THE COST (CFO NUMBERS)
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.08), NAVY)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.6),
             "What Manual AP Actually Costs You",
             font_size=32, color=NAVY, bold=True)

add_text_box(slide, Inches(0.8), Inches(1.2), Inches(11), Inches(0.5),
             "Based on 2,000 invoices/month (typical multi-trade contractor):",
             font_size=18, color=GRAY)

# Cost breakdown table
table_data = [
    ["Cost Component", "Per Invoice", "Monthly", "Annual"],
    ["AP Staff Labor (data entry, coding, routing, filing)", "$10.00", "$20,000", "$240,000"],
    ["Error Correction & Rework", "$2.00", "$4,000", "$48,000"],
    ["Overhead (office, software, printing)", "$1.00", "$2,000", "$24,000"],
    ["Late Payments & Missed Discounts", "$2.50", "$5,000", "$60,000"],
    ["Vendor Calls & Follow-ups", "$0.50", "$1,000", "$12,000"],
    ["TOTAL COST", "$16.00", "$32,000", "$384,000"],
]

rows, cols = len(table_data), len(table_data[0])
tbl = slide.shapes.add_table(rows, cols, Inches(0.8), Inches(1.8), Inches(11.5), Inches(3.0)).table

tbl.columns[0].width = Inches(5.5)
tbl.columns[1].width = Inches(1.8)
tbl.columns[2].width = Inches(2.0)
tbl.columns[3].width = Inches(2.2)

for r in range(rows):
    for c in range(cols):
        cell = tbl.cell(r, c)
        cell.text = table_data[r][c]
        p = cell.text_frame.paragraphs[0]
        p.font.size = Pt(14)
        p.font.name = "Calibri"
        if r == 0:
            cell.fill.solid()
            cell.fill.fore_color.rgb = NAVY
            p.font.color.rgb = WHITE
            p.font.bold = True
        elif r == rows - 1:
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(0xfe, 0xd7, 0xd7)
            p.font.color.rgb = RED
            p.font.bold = True
        else:
            cell.fill.solid()
            cell.fill.fore_color.rgb = LIGHT_GRAY if r % 2 == 0 else WHITE
            p.font.color.rgb = DARK
        if c > 0:
            p.alignment = PP_ALIGN.CENTER

# Big number callout
callout = add_shape(slide, Inches(3.5), Inches(5.2), Inches(6.3), Inches(1.5), RGBColor(0xfe, 0xd7, 0xd7), RED)
add_text_box(slide, Inches(3.7), Inches(5.3), Inches(5.9), Inches(0.5),
             "You're spending up to", font_size=18, color=RED, bold=False, alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(3.7), Inches(5.7), Inches(5.9), Inches(0.7),
             "$384,000 / year", font_size=44, color=RED, bold=True, alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(3.7), Inches(6.3), Inches(5.9), Inches(0.4),
             "on a process that AI can do in seconds", font_size=16, color=RED, alignment=PP_ALIGN.CENTER)

# ============================================================
# SLIDE 4: THE SOLUTION
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.08), NAVY)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.6),
             "The Solution: PayFlow AI",
             font_size=32, color=NAVY, bold=True)
add_text_box(slide, Inches(0.8), Inches(1.0), Inches(11), Inches(0.5),
             "AI that reads, codes, matches, and routes invoices — built specifically for construction.",
             font_size=18, color=GRAY)

# Feature cards
features = [
    ("AI Invoice Reading", "Upload a PDF, photo, or email.\nAI extracts every field in seconds.\n99% accuracy, any format.", BLUE),
    ("5D Job Costing", "AI predicts GL account, job number,\nphase, cost code, and cost type.\nLearns from your patterns.", BLUE),
    ("Three-Way PO Match", "Automatically matches invoices\nagainst POs and receiving reports.\nFlags variances instantly.", BLUE),
    ("AIA Billing Support", "Reads G702/G703 progress billing.\nExtracts schedule of values.\nVerifies against contract.", GREEN),
    ("Lien Waiver Tracking", "AI reads and classifies waivers.\nBlocks payment if missing.\nTracks per vendor per job.", GREEN),
    ("Fraud Detection", "Detects duplicates, overbilling,\nprice anomalies, and suspicious\npatterns in real-time.", ORANGE),
]

for i, (title, desc, color) in enumerate(features):
    col = i % 3
    row = i // 3
    left = Inches(0.8) + Inches(col * 4.1)
    top = Inches(1.8) + Inches(row * 2.5)

    card = add_shape(slide, left, top, Inches(3.7), Inches(2.1), LIGHT_GRAY, MEDIUM_GRAY)
    # Color accent bar at top of card
    add_rect(slide, left + Inches(0.05), top + Inches(0.05), Inches(3.6), Inches(0.08), color)
    add_text_box(slide, left + Inches(0.2), top + Inches(0.25), Inches(3.3), Inches(0.4),
                 title, font_size=18, color=NAVY, bold=True)

    txBox = slide.shapes.add_textbox(left + Inches(0.2), top + Inches(0.7), Inches(3.3), Inches(1.3))
    tf = txBox.text_frame
    tf.word_wrap = True
    for line in desc.split("\n"):
        p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(13)
        p.font.color.rgb = GRAY
        p.font.name = "Calibri"
        p.space_after = Pt(2)

# ============================================================
# SLIDE 5: HOW IT WORKS
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.08), NAVY)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.6),
             "How It Works: 5 Seconds, Not 5 Days",
             font_size=32, color=NAVY, bold=True)

steps = [
    ("1", "Invoice Arrives", "Email, upload, scan,\nor vendor portal", BLUE),
    ("2", "AI Reads It", "Extracts all fields,\nline items, amounts", BLUE),
    ("3", "AI Codes It", "Predicts GL, job, phase,\ncost code, cost type", BLUE),
    ("4", "AI Validates", "PO match, fraud check,\nretention, lien waivers", GREEN),
    ("5", "Routes or Approves", "Auto-approve if confident.\nRoute to human if not.", GREEN),
    ("6", "Posts to Sage", "Approved invoices sync\ndirectly to your ERP", NAVY),
]

for i, (num, title, desc, color) in enumerate(steps):
    left = Inches(0.5) + Inches(i * 2.1)
    top = Inches(1.5)

    # Circle with number
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, left + Inches(0.65), top, Inches(0.7), Inches(0.7))
    circle.fill.solid()
    circle.fill.fore_color.rgb = color
    circle.line.fill.background()
    tf = circle.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.text = num
    p.font.size = Pt(24)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    tf.paragraphs[0].space_before = Pt(4)

    # Arrow (except last)
    if i < len(steps) - 1:
        arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, left + Inches(1.5), top + Inches(0.2), Inches(0.5), Inches(0.3))
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = MEDIUM_GRAY
        arrow.line.fill.background()

    add_text_box(slide, left + Inches(0.1), top + Inches(0.9), Inches(1.8), Inches(0.4),
                 title, font_size=14, color=NAVY, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, left + Inches(0.1), top + Inches(1.3), Inches(1.8), Inches(0.8),
                 desc, font_size=11, color=GRAY, alignment=PP_ALIGN.CENTER)

# Bottom section - before/after
add_rect(slide, Inches(0.8), Inches(3.5), Inches(5.6), Inches(3.3), RGBColor(0xfe, 0xd7, 0xd7))
add_text_box(slide, Inches(1.0), Inches(3.6), Inches(5.2), Inches(0.4),
             "TODAY (Manual Process)", font_size=18, color=RED, bold=True, alignment=PP_ALIGN.CENTER)
today_items = [
    "AP clerk opens email, prints or downloads invoice",
    "Manually types vendor, amount, date into Sage",
    "Looks up job number and cost code in spreadsheet",
    "Walks invoice to PM for approval (or emails and waits)",
    "PM reviews, asks questions, eventually approves",
    "AP clerk posts to Sage, files paper copy",
    "Total: 20-30 minutes per invoice, 14+ days to pay",
]
add_bullet_list(slide, Inches(1.2), Inches(4.1), Inches(5), Inches(2.5), today_items, font_size=12, color=RGBColor(0xc5, 0x30, 0x30))

add_rect(slide, Inches(6.9), Inches(3.5), Inches(5.6), Inches(3.3), LIGHT_GREEN)
add_text_box(slide, Inches(7.1), Inches(3.6), Inches(5.2), Inches(0.4),
             "WITH PAYFLOW AI", font_size=18, color=GREEN, bold=True, alignment=PP_ALIGN.CENTER)
ai_items = [
    "Invoice arrives via email or upload",
    "AI reads and extracts all fields automatically",
    "AI codes to correct GL, job, phase, cost code",
    "AI matches against PO, checks for anomalies",
    "Auto-approves if confident; routes to PM if not",
    "Approved invoices post directly to Sage",
    "Total: 5 seconds per invoice, same-day processing",
]
add_bullet_list(slide, Inches(7.3), Inches(4.1), Inches(5), Inches(2.5), ai_items, font_size=12, color=RGBColor(0x27, 0x67, 0x49))

# ============================================================
# SLIDE 6: ROI / SAVINGS
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.08), NAVY)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.6),
             "The ROI: Payback in Weeks, Not Years",
             font_size=32, color=NAVY, bold=True)

# Savings comparison
table_data2 = [
    ["", "Manual (Current)", "With PayFlow AI", "Savings"],
    ["Cost per invoice", "$16-$25", "$2-$4", "$12-$21"],
    ["Processing time", "20-30 minutes", "5 seconds", "99.5% faster"],
    ["Error rate", "2-5%", "<0.5%", "90% fewer errors"],
    ["Invoices per FTE per hour", "2.9", "50+", "17x more productive"],
    ["Late payment rate", "15-20%", "<2%", "Early-pay discounts captured"],
]

rows2, cols2 = len(table_data2), len(table_data2[0])
tbl2 = slide.shapes.add_table(rows2, cols2, Inches(0.8), Inches(1.3), Inches(11.5), Inches(2.5)).table
tbl2.columns[0].width = Inches(3.0)
tbl2.columns[1].width = Inches(2.8)
tbl2.columns[2].width = Inches(2.8)
tbl2.columns[3].width = Inches(2.9)

for r in range(rows2):
    for c in range(cols2):
        cell = tbl2.cell(r, c)
        cell.text = table_data2[r][c]
        p = cell.text_frame.paragraphs[0]
        p.font.size = Pt(14)
        p.font.name = "Calibri"
        if r == 0:
            cell.fill.solid()
            cell.fill.fore_color.rgb = NAVY
            p.font.color.rgb = WHITE
            p.font.bold = True
        else:
            cell.fill.solid()
            cell.fill.fore_color.rgb = LIGHT_GRAY if r % 2 == 0 else WHITE
            p.font.color.rgb = DARK
            if c == 3:
                p.font.color.rgb = GREEN
                p.font.bold = True
        if c > 0:
            p.alignment = PP_ALIGN.CENTER

# Big savings callout
add_text_box(slide, Inches(0.8), Inches(4.1), Inches(11.5), Inches(0.4),
             "Annual Impact (at 2,000 invoices/month):", font_size=18, color=NAVY, bold=True)

add_stat_card(slide, Inches(0.8), Inches(4.6), Inches(3.5), Inches(1.5),
              "CURRENT ANNUAL AP COST", "$384,000", "Labor + errors + overhead + hidden costs", RED)
add_stat_card(slide, Inches(4.8), Inches(4.6), Inches(3.5), Inches(1.5),
              "WITH PAYFLOW AI", "$48,000", "Reduced staff + software subscription", BLUE)
add_stat_card(slide, Inches(8.8), Inches(4.6), Inches(3.7), Inches(1.5),
              "ANNUAL SAVINGS", "$336,000", "87% reduction in AP processing costs", GREEN)

add_text_box(slide, Inches(0.8), Inches(6.4), Inches(11.5), Inches(0.5),
             "ROI: The software pays for itself in the first month.",
             font_size=20, color=GREEN, bold=True, alignment=PP_ALIGN.CENTER)

# ============================================================
# SLIDE 7: CONSTRUCTION-SPECIFIC
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.08), NAVY)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.6),
             "Built for Construction, Not Adapted From Generic AP",
             font_size=32, color=NAVY, bold=True)
add_text_box(slide, Inches(0.8), Inches(1.0), Inches(11), Inches(0.5),
             "Generic AP tools don't understand construction. We do.",
             font_size=18, color=GRAY)

# Two column comparison
add_rect(slide, Inches(0.8), Inches(1.7), Inches(5.6), Inches(5.0), RGBColor(0xfe, 0xf2, 0xf2))
add_text_box(slide, Inches(1.0), Inches(1.8), Inches(5.2), Inches(0.4),
             "What Generic AP Tools Can't Do", font_size=18, color=RED, bold=True)
generic_items = [
    "Read AIA G702/G703 progress billing forms",
    "Track retention/retainage per vendor per job",
    "Require lien waivers before releasing payment",
    "Code to 5 dimensions (GL + job + phase + cost code + type)",
    "Verify invoices against evolving contract values",
    "Detect overbilling beyond contract + change orders",
    "Track COI expiration and subcontractor compliance",
    "Handle multi-entity routing (Mechanical LLC, Electrical LLC)",
    "Understand prevailing wage and certified payroll",
]
add_bullet_list(slide, Inches(1.2), Inches(2.3), Inches(5), Inches(4.2), generic_items, font_size=13, color=RGBColor(0xc5, 0x30, 0x30))

add_rect(slide, Inches(6.9), Inches(1.7), Inches(5.6), Inches(5.0), RGBColor(0xf0, 0xff, 0xf4))
add_text_box(slide, Inches(7.1), Inches(1.8), Inches(5.2), Inches(0.4),
             "What PayFlow AI Does Natively", font_size=18, color=GREEN, bold=True)
payflow_items = [
    "AI reads and extracts G702/G703 fields automatically",
    "Retention dashboard tracks held amounts across all jobs",
    "Blocks payment if conditional lien waiver is missing",
    "AI predicts all 5 cost dimensions per line item",
    "Tracks original contract + change orders vs. billed-to-date",
    "Flags invoices that exceed approved contract amounts",
    "Alerts when insurance certificates expire",
    "Routes to correct entity based on job and trade",
    "Verifies labor rates against prevailing wage tables",
]
add_bullet_list(slide, Inches(7.3), Inches(2.3), Inches(5), Inches(4.2), payflow_items, font_size=13, color=RGBColor(0x27, 0x67, 0x49))

# ============================================================
# SLIDE 8: INTEGRATIONS
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.08), NAVY)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.6),
             "Integrates With Your Existing Systems",
             font_size=32, color=NAVY, bold=True)
add_text_box(slide, Inches(0.8), Inches(1.0), Inches(11), Inches(0.5),
             "No rip-and-replace. PayFlow AI plugs into what you already use.",
             font_size=18, color=GRAY)

integrations = [
    ("Sage", "Your Accounting System", "Two-way sync: chart of accounts,\nvendors, jobs, cost codes.\nApproved invoices post directly.", NAVY),
    ("Procore", "Your Project Management", "Pull POs, change orders, budgets,\ncost codes. Three-way matching.\nReal-time contract tracking.", BLUE),
    ("Email", "Invoice Ingestion", "Dedicated inbox for your company.\nVendors email invoices directly.\nAI processes automatically.", GREEN),
    ("Bank / Payment", "Payment Execution", "Schedule payments for approved\ninvoices. ACH, check, or card.\nReconcile automatically.", ORANGE),
]

for i, (name, subtitle, desc, color) in enumerate(integrations):
    left = Inches(0.6) + Inches(i * 3.15)
    top = Inches(1.8)
    card = add_shape(slide, left, top, Inches(2.9), Inches(3.2), LIGHT_GRAY, MEDIUM_GRAY)
    add_rect(slide, left + Inches(0.05), top + Inches(0.05), Inches(2.8), Inches(0.08), color)

    add_text_box(slide, left + Inches(0.2), top + Inches(0.3), Inches(2.5), Inches(0.4),
                 name, font_size=24, color=color, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, left + Inches(0.2), top + Inches(0.75), Inches(2.5), Inches(0.3),
                 subtitle, font_size=12, color=GRAY, alignment=PP_ALIGN.CENTER)

    txBox = slide.shapes.add_textbox(left + Inches(0.2), top + Inches(1.2), Inches(2.5), Inches(1.8))
    tf = txBox.text_frame
    tf.word_wrap = True
    for line in desc.split("\n"):
        p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(13)
        p.font.color.rgb = GRAY
        p.font.name = "Calibri"
        p.alignment = PP_ALIGN.CENTER
        p.space_after = Pt(2)

# Data flow
add_text_box(slide, Inches(0.8), Inches(5.4), Inches(11.5), Inches(0.4),
             "Data Flow:  Invoices In  >  AI Processes  >  Human Reviews (if needed)  >  Posts to Sage  >  Payment Scheduled",
             font_size=16, color=NAVY, bold=True, alignment=PP_ALIGN.CENTER)

add_shape(slide, Inches(1.5), Inches(6.0), Inches(10.3), Inches(0.9), LIGHT_BLUE, BLUE)
add_text_box(slide, Inches(1.7), Inches(6.1), Inches(9.9), Inches(0.7),
             "Your team focuses on exceptions and approvals. The AI handles the other 80%+ automatically.",
             font_size=16, color=NAVY, alignment=PP_ALIGN.CENTER)

# ============================================================
# SLIDE 9: SECURITY & COMPLIANCE
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.08), NAVY)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.6),
             "Security & Compliance",
             font_size=32, color=NAVY, bold=True)
add_text_box(slide, Inches(0.8), Inches(1.0), Inches(11), Inches(0.5),
             "Your financial data is protected at every step.",
             font_size=18, color=GRAY)

security_items = [
    ("Data Encryption", "AES-256 encryption at rest, TLS 1.2+ in transit.\nYour data is encrypted end-to-end, always."),
    ("AI Privacy", "Invoice data is processed via API — never used to train AI models.\nYour data stays your data. Period."),
    ("Audit Trail", "Every action logged: who uploaded, what AI extracted,\nwho approved, what changed, and when."),
    ("Role-Based Access", "AP clerks, project managers, controllers, and CFOs\neach see only what they need."),
    ("SOC 2 Compliance", "Enterprise-grade security controls.\nAnnual third-party audits."),
    ("Data Ownership", "You own your data. Always. Full export available.\nIf you leave, your data leaves with you."),
]

for i, (title, desc) in enumerate(security_items):
    col = i % 2
    row = i // 2
    left = Inches(0.8) + Inches(col * 6.2)
    top = Inches(1.7) + Inches(row * 1.7)

    card = add_shape(slide, left, top, Inches(5.8), Inches(1.4), LIGHT_GRAY, MEDIUM_GRAY)
    add_text_box(slide, left + Inches(0.3), top + Inches(0.15), Inches(5.2), Inches(0.35),
                 title, font_size=16, color=NAVY, bold=True)
    txBox = slide.shapes.add_textbox(left + Inches(0.3), top + Inches(0.5), Inches(5.2), Inches(0.8))
    tf = txBox.text_frame
    tf.word_wrap = True
    for line in desc.split("\n"):
        p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(12)
        p.font.color.rgb = GRAY
        p.font.name = "Calibri"
        p.space_after = Pt(1)

# ============================================================
# SLIDE 10: IMPLEMENTATION
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.08), NAVY)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.6),
             "Implementation: Fast, Low-Risk, High-Reward",
             font_size=32, color=NAVY, bold=True)

phases = [
    ("Week 1-2", "Setup & Connect", "Connect Sage and Procore.\nImport chart of accounts,\njobs, vendors, cost codes.", BLUE),
    ("Week 3-4", "AI Training", "Process your real invoices.\nTune AI to your patterns.\nVerify accuracy.", BLUE),
    ("Week 5-6", "Shadow Mode", "AI processes alongside your\nteam. Compare results.\nBuild confidence.", GREEN),
    ("Week 7-8", "Go Live", "AI handles routine invoices.\nTeam focuses on exceptions.\nStart saving immediately.", GREEN),
]

for i, (timeline, title, desc, color) in enumerate(phases):
    left = Inches(0.6) + Inches(i * 3.15)
    top = Inches(1.5)

    card = add_shape(slide, left, top, Inches(2.9), Inches(3.0), WHITE, color)
    add_rect(slide, left + Inches(0.05), top + Inches(0.05), Inches(2.8), Inches(0.5), color)
    add_text_box(slide, left + Inches(0.1), top + Inches(0.08), Inches(2.7), Inches(0.45),
                 timeline, font_size=16, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, left + Inches(0.2), top + Inches(0.65), Inches(2.5), Inches(0.4),
                 title, font_size=18, color=NAVY, bold=True, alignment=PP_ALIGN.CENTER)

    txBox = slide.shapes.add_textbox(left + Inches(0.2), top + Inches(1.1), Inches(2.5), Inches(1.7))
    tf = txBox.text_frame
    tf.word_wrap = True
    for line in desc.split("\n"):
        p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(13)
        p.font.color.rgb = GRAY
        p.font.name = "Calibri"
        p.alignment = PP_ALIGN.CENTER
        p.space_after = Pt(3)

# Bottom callouts
add_shape(slide, Inches(0.8), Inches(5.0), Inches(3.7), Inches(1.2), LIGHT_BLUE, BLUE)
add_text_box(slide, Inches(1.0), Inches(5.1), Inches(3.3), Inches(0.35),
             "Zero Disruption", font_size=16, color=NAVY, bold=True, alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(1.0), Inches(5.45), Inches(3.3), Inches(0.6),
             "Your team keeps working normally.\nAI runs in parallel until you're ready.",
             font_size=12, color=NAVY, alignment=PP_ALIGN.CENTER)

add_shape(slide, Inches(4.8), Inches(5.0), Inches(3.7), Inches(1.2), LIGHT_GREEN, GREEN)
add_text_box(slide, Inches(5.0), Inches(5.1), Inches(3.3), Inches(0.35),
             "Immediate Value", font_size=16, color=RGBColor(0x27, 0x67, 0x49), bold=True, alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(5.0), Inches(5.45), Inches(3.3), Inches(0.6),
             "See ROI from the first week of\ngo-live. No 6-month wait.",
             font_size=12, color=RGBColor(0x27, 0x67, 0x49), alignment=PP_ALIGN.CENTER)

add_shape(slide, Inches(8.8), Inches(5.0), Inches(3.7), Inches(1.2), RGBColor(0xff, 0xf5, 0xeb), ORANGE)
add_text_box(slide, Inches(9.0), Inches(5.1), Inches(3.3), Inches(0.35),
             "Low Risk", font_size=16, color=ORANGE, bold=True, alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(9.0), Inches(5.45), Inches(3.3), Inches(0.6),
             "Shadow mode proves accuracy\nbefore you commit. No surprises.",
             font_size=12, color=ORANGE, alignment=PP_ALIGN.CENTER)

# ============================================================
# SLIDE 11: PRICING
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.08), NAVY)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.6),
             "Investment",
             font_size=32, color=NAVY, bold=True)

# Pricing tiers
tiers = [
    ("Starter", "$500", "/month", "Up to 1,000 invoices\n5 users\nSage integration\nEmail ingestion\nStandard support", BLUE),
    ("Professional", "$1,500", "/month", "Up to 5,000 invoices\nUnlimited users\nSage + Procore integration\nAIA billing support\nLien waiver tracking\nPriority support", GREEN),
    ("Enterprise", "Custom", "", "Unlimited invoices\nUnlimited users\nAll integrations\nDedicated account manager\nCustom workflows\nSLA guarantee", NAVY),
]

for i, (name, price, period, features_text, color) in enumerate(tiers):
    left = Inches(1.0) + Inches(i * 3.9)
    top = Inches(1.3)

    card = add_shape(slide, left, top, Inches(3.5), Inches(4.8), WHITE, color)
    add_rect(slide, left + Inches(0.05), top + Inches(0.05), Inches(3.4), Inches(0.08), color)

    add_text_box(slide, left + Inches(0.1), top + Inches(0.3), Inches(3.3), Inches(0.4),
                 name, font_size=20, color=color, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, left + Inches(0.1), top + Inches(0.8), Inches(3.3), Inches(0.6),
                 price, font_size=40, color=DARK, bold=True, alignment=PP_ALIGN.CENTER)
    if period:
        add_text_box(slide, left + Inches(0.1), top + Inches(1.3), Inches(3.3), Inches(0.3),
                     period, font_size=14, color=GRAY, alignment=PP_ALIGN.CENTER)

    txBox = slide.shapes.add_textbox(left + Inches(0.3), top + Inches(1.7), Inches(2.9), Inches(2.8))
    tf = txBox.text_frame
    tf.word_wrap = True
    for line in features_text.split("\n"):
        p = tf.add_paragraph()
        p.text = "  " + line
        p.font.size = Pt(13)
        p.font.color.rgb = GRAY
        p.font.name = "Calibri"
        p.space_after = Pt(6)

# Bottom note
add_text_box(slide, Inches(0.8), Inches(6.4), Inches(11.5), Inches(0.5),
             "All plans include: AI invoice extraction, GL coding, fraud detection, approval workflows, and audit trails.",
             font_size=14, color=GRAY, alignment=PP_ALIGN.CENTER)

# ============================================================
# SLIDE 12: NEXT STEPS / CTA
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, NAVY)
add_rect(slide, Inches(0), Inches(3.2), Inches(13.333), Inches(0.06), BLUE)

add_text_box(slide, Inches(1), Inches(1.2), Inches(11), Inches(0.8),
             "Let's Get Started", font_size=44, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(2), Inches(2.2), Inches(9), Inches(0.6),
             "Three simple steps to transform your AP process:",
             font_size=20, color=RGBColor(0xa0, 0xce, 0xf0), alignment=PP_ALIGN.CENTER)

steps_text = [
    ("1.", "Share 50 sample invoices", "We'll run them through the AI and show you the results — no commitment."),
    ("2.", "2-week shadow pilot", "AI processes your invoices alongside your team. Compare accuracy side by side."),
    ("3.", "Go live with confidence", "Only when you've seen the results. Start saving from day one."),
]

for i, (num, title, desc) in enumerate(steps_text):
    top = Inches(3.6) + Inches(i * 1.0)
    add_text_box(slide, Inches(2.5), top, Inches(0.5), Inches(0.4),
                 num, font_size=24, color=RGBColor(0xa0, 0xce, 0xf0), bold=True)
    add_text_box(slide, Inches(3.1), top, Inches(3), Inches(0.4),
                 title, font_size=20, color=WHITE, bold=True)
    add_text_box(slide, Inches(3.1), top + Inches(0.35), Inches(7), Inches(0.4),
                 desc, font_size=14, color=RGBColor(0xcb, 0xd5, 0xe0))

add_text_box(slide, Inches(1), Inches(6.5), Inches(11), Inches(0.5),
             "PayFlow AI  |  AI-Powered AP Automation for Construction",
             font_size=16, color=RGBColor(0x71, 0x80, 0x96), alignment=PP_ALIGN.CENTER)

# SAVE
output_path = "C:/users/25badmin/projects/accounts-payable-research/PayFlow_AI_Sales_Deck.pptx"
prs.save(output_path)
print(f"Presentation saved: {output_path}")
