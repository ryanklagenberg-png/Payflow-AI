from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_LABEL_POSITION
from pptx.chart.data import CategoryChartData

# Colors
NAVY = RGBColor(0x1a, 0x36, 0x5d)
BLUE = RGBColor(0x2b, 0x6c, 0xb0)
LIGHT_BLUE = RGBColor(0xeb, 0xf4, 0xff)
GREEN = RGBColor(0x38, 0xa1, 0x69)
LIGHT_GREEN = RGBColor(0xc6, 0xf6, 0xd5)
DARK_GREEN = RGBColor(0x27, 0x67, 0x49)
RED = RGBColor(0xe5, 0x3e, 0x3e)
LIGHT_RED = RGBColor(0xfe, 0xd7, 0xd7)
WHITE = RGBColor(0xff, 0xff, 0xff)
DARK = RGBColor(0x1a, 0x20, 0x2c)
GRAY = RGBColor(0x4a, 0x55, 0x68)
LIGHT_GRAY = RGBColor(0xf7, 0xfa, 0xfc)
MEDIUM_GRAY = RGBColor(0xe2, 0xe8, 0xf0)
ORANGE = RGBColor(0xdd, 0x6b, 0x20)
GOLD = RGBColor(0xd6, 0x9e, 0x2e)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)


def add_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape(slide, left, top, width, height, fill_color, border_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape


def add_rect(slide, left, top, width, height, fill_color):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape


def text(slide, left, top, width, height, txt, size=18, color=DARK, bold=False, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = txt
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = "Calibri"
    p.alignment = align
    return box


def bullets(slide, left, top, width, height, items, size=15, color=DARK, spacing=8):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.name = "Calibri"
        p.space_after = Pt(spacing)


def stat_card(slide, left, top, w, h, label, value, sub="", vcolor=NAVY):
    add_shape(slide, left, top, w, h, WHITE, MEDIUM_GRAY)
    text(slide, left + Inches(0.15), top + Inches(0.1), w - Inches(0.3), Inches(0.25),
         label, size=10, color=GRAY, bold=True)
    text(slide, left + Inches(0.15), top + Inches(0.35), w - Inches(0.3), Inches(0.5),
         value, size=30, color=vcolor, bold=True, align=PP_ALIGN.CENTER)
    if sub:
        text(slide, left + Inches(0.15), top + Inches(0.85), w - Inches(0.3), Inches(0.25),
             sub, size=9, color=GRAY, align=PP_ALIGN.CENTER)


def slide_header(slide, title, subtitle=""):
    add_rect(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.08), NAVY)
    text(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.6), title, size=32, color=NAVY, bold=True)
    if subtitle:
        text(slide, Inches(0.8), Inches(0.9), Inches(11), Inches(0.4), subtitle, size=17, color=GRAY)


def add_table(slide, data, left, top, col_widths, row_height=Inches(0.4)):
    rows, cols = len(data), len(data[0])
    total_w = sum(col_widths)
    tbl = slide.shapes.add_table(rows, cols, left, top, total_w, row_height * rows).table
    for c, w in enumerate(col_widths):
        tbl.columns[c].width = w
    for r in range(rows):
        for c in range(cols):
            cell = tbl.cell(r, c)
            cell.text = data[r][c]
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(13)
            p.font.name = "Calibri"
            if r == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = NAVY
                p.font.color.rgb = WHITE
                p.font.bold = True
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = LIGHT_GRAY if r % 2 == 0 else WHITE
                p.font.color.rgb = DARK
            if c > 0:
                p.alignment = PP_ALIGN.CENTER
    return tbl


# ============================================================
# SLIDE 1: TITLE
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, NAVY)
add_rect(slide, Inches(0), Inches(3.4), Inches(13.333), Inches(0.06), BLUE)

text(slide, Inches(1), Inches(1.3), Inches(11), Inches(1),
     "PayFlow AI", size=58, color=WHITE, bold=True)
text(slide, Inches(1), Inches(2.3), Inches(11), Inches(0.8),
     "AI-Powered Accounts Payable Automation for Construction",
     size=26, color=RGBColor(0xa0, 0xce, 0xf0))
text(slide, Inches(1), Inches(3.8), Inches(11), Inches(0.5),
     "Stop paying $576,000/year to process invoices manually.",
     size=20, color=RGBColor(0xcb, 0xd5, 0xe0))
text(slide, Inches(1), Inches(4.3), Inches(11), Inches(0.5),
     "Start processing them in 5 seconds for pennies.",
     size=20, color=RGBColor(0xcb, 0xd5, 0xe0))
text(slide, Inches(1), Inches(5.8), Inches(11), Inches(0.4),
     "Confidential  |  2026", size=14, color=RGBColor(0x71, 0x80, 0x96))

# ============================================================
# SLIDE 2: THE PROBLEM - BY THE NUMBERS
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
slide_header(slide, "Your AP Department By The Numbers",
             "Based on 3,000 invoices per month")

stat_card(slide, Inches(0.8), Inches(1.5), Inches(2.5), Inches(1.2),
          "INVOICES / MONTH", "3,000", "", NAVY)
stat_card(slide, Inches(3.6), Inches(1.5), Inches(2.5), Inches(1.2),
          "COST PER INVOICE", "$16", "Industry average (construction)", RED)
stat_card(slide, Inches(6.4), Inches(1.5), Inches(2.5), Inches(1.2),
          "MONTHLY AP COST", "$48,000", "Labor + errors + overhead", RED)
stat_card(slide, Inches(9.2), Inches(1.5), Inches(2.5), Inches(1.2),
          "ANNUAL AP COST", "$576,000", "Going out the door every year", RED)

# Pie chart - where the money goes
chart_data = CategoryChartData()
chart_data.categories = ["AP Staff Labor\n$360K", "Error Correction\n$72K",
                         "Late Payments &\nMissed Discounts\n$90K",
                         "Overhead\n$36K", "Vendor\nFollow-ups\n$18K"]
chart_data.add_series("Cost", (360, 72, 90, 36, 18))

chart = slide.shapes.add_chart(
    XL_CHART_TYPE.PIE, Inches(0.8), Inches(3.0), Inches(5.5), Inches(4.0), chart_data
).chart
chart.has_legend = True
chart.legend.position = XL_LEGEND_POSITION.BOTTOM
chart.legend.include_in_layout = False
chart.legend.font.size = Pt(10)
chart.legend.font.color.rgb = GRAY

plot = chart.plots[0]
plot.has_data_labels = True
data_labels = plot.data_labels
data_labels.font.size = Pt(11)
data_labels.font.bold = True
data_labels.font.color.rgb = WHITE
data_labels.number_format = '0"%"'

series = chart.series[0]
colors = [NAVY, RED, ORANGE, BLUE, GRAY]
for i, color in enumerate(colors):
    pt = series.points[i]
    pt.format.fill.solid()
    pt.format.fill.fore_color.rgb = color

# Right side - what's eating your money
text(slide, Inches(6.8), Inches(3.0), Inches(5.5), Inches(0.4),
     "Where Your $576,000 Goes:", size=20, color=NAVY, bold=True)

money_items = [
    "62% ($360K) — Staff manually typing invoice data into Sage",
    "16% ($90K) — Late payment penalties and missed early-pay discounts",
    "13% ($72K) — Finding and fixing data entry errors",
    "6% ($36K) — Office space, software, printing, filing",
    "3% ($18K) — Calling vendors about payment status",
]
bullets(slide, Inches(7.0), Inches(3.5), Inches(5.5), Inches(3.0), money_items, size=14, color=GRAY)

text(slide, Inches(6.8), Inches(6.0), Inches(5.5), Inches(0.5),
     "The #1 cost driver is manual data entry — exactly what AI eliminates.",
     size=14, color=RED, bold=True)

# ============================================================
# SLIDE 3: THE SOLUTION
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
slide_header(slide, "PayFlow AI: What It Does",
             "AI reads, codes, matches, and routes invoices — in seconds, not days.")

# 6 feature cards in 2 rows
features = [
    ("AI Invoice Extraction", "Reads any invoice format — PDF, scan,\nphoto, email. Extracts every field\nwith 99% accuracy. No templates.", BLUE),
    ("5-Dimensional Job Costing", "Predicts GL account, job number,\nphase, cost code, and cost type\nfor every line item automatically.", BLUE),
    ("Three-Way PO Matching", "Matches invoices against POs and\nreceiving reports. Flags price and\nquantity variances instantly.", BLUE),
    ("AIA Progress Billing", "Reads G702/G703 forms natively.\nExtracts schedule of values, tracks\npercent complete, verifies totals.", GREEN),
    ("Lien Waiver Management", "AI reads and classifies waivers.\nBlocks payment if missing.\n4 waiver types, state-aware.", GREEN),
    ("Fraud & Anomaly Detection", "Catches duplicates, overbilling,\nprice spikes, and suspicious patterns.\nReal-time risk scoring.", ORANGE),
]

for i, (title, desc, color) in enumerate(features):
    col = i % 3
    row = i // 3
    left = Inches(0.8) + Inches(col * 4.05)
    top = Inches(1.6) + Inches(row * 2.7)
    card = add_shape(slide, left, top, Inches(3.7), Inches(2.3), WHITE, MEDIUM_GRAY)
    add_rect(slide, left + Inches(0.05), top + Inches(0.05), Inches(3.6), Inches(0.08), color)
    text(slide, left + Inches(0.2), top + Inches(0.25), Inches(3.3), Inches(0.35),
         title, size=17, color=NAVY, bold=True)
    box = slide.shapes.add_textbox(left + Inches(0.2), top + Inches(0.7), Inches(3.3), Inches(1.4))
    tf = box.text_frame
    tf.word_wrap = True
    for line in desc.split("\n"):
        p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(13)
        p.font.color.rgb = GRAY
        p.font.name = "Calibri"
        p.space_after = Pt(2)

# ============================================================
# SLIDE 4: HOW IT WORKS - VISUAL FLOW
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
slide_header(slide, "How It Works")

steps = [
    ("1", "Invoice\nArrives", "Email, upload,\nscan, or portal", BLUE),
    ("2", "AI Reads\nDocument", "Extracts all fields\nin 3-5 seconds", BLUE),
    ("3", "AI Codes\nto Job", "GL, job, phase,\ncost code, type", BLUE),
    ("4", "AI Validates\n& Matches", "PO match, fraud\ncheck, compliance", GREEN),
    ("5", "Route or\nAuto-Approve", "85%+ go straight\nthrough", GREEN),
    ("6", "Posts to\nSage", "Syncs directly\nto your ERP", NAVY),
]

for i, (num, title, desc, color) in enumerate(steps):
    left = Inches(0.4) + Inches(i * 2.15)
    top = Inches(1.4)
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, left + Inches(0.55), top, Inches(0.9), Inches(0.9))
    circle.fill.solid()
    circle.fill.fore_color.rgb = color
    circle.line.fill.background()
    tf = circle.text_frame
    p = tf.paragraphs[0]
    p.text = num
    p.font.size = Pt(28)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    if i < len(steps) - 1:
        arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, left + Inches(1.6), top + Inches(0.25), Inches(0.45), Inches(0.4))
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = MEDIUM_GRAY
        arrow.line.fill.background()

    text(slide, left, top + Inches(1.1), Inches(2.0), Inches(0.6),
         title, size=14, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
    text(slide, left, top + Inches(1.7), Inches(2.0), Inches(0.5),
         desc, size=11, color=GRAY, align=PP_ALIGN.CENTER)

# Before/After bar chart
chart_data2 = CategoryChartData()
chart_data2.categories = ["Time per\nInvoice", "Cost per\nInvoice", "Error\nRate", "Days to\nPayment"]
chart_data2.add_series("Manual (Today)", (25, 16, 4, 14.6))
chart_data2.add_series("PayFlow AI", (0.08, 3, 0.5, 1))

chart_frame = slide.shapes.add_chart(
    XL_CHART_TYPE.BAR_CLUSTERED, Inches(0.8), Inches(3.2), Inches(6.0), Inches(3.8), chart_data2
)
chart2 = chart_frame.chart
chart2.has_legend = True
chart2.legend.position = XL_LEGEND_POSITION.BOTTOM
chart2.legend.font.size = Pt(11)

chart2.series[0].format.fill.solid()
chart2.series[0].format.fill.fore_color.rgb = RED
chart2.series[1].format.fill.solid()
chart2.series[1].format.fill.fore_color.rgb = GREEN

plot2 = chart2.plots[0]
plot2.has_data_labels = True
plot2.data_labels.font.size = Pt(10)
plot2.data_labels.font.bold = True
plot2.gap_width = 80

cat_axis = chart2.category_axis
cat_axis.tick_labels.font.size = Pt(10)
val_axis = chart2.value_axis
val_axis.visible = False

# Right side stats
text(slide, Inches(7.2), Inches(3.2), Inches(5.5), Inches(0.4),
     "The Transformation:", size=20, color=NAVY, bold=True)

transform = [
    ("Processing Time:", "25 min  -->  5 seconds", "(99.7% faster)"),
    ("Cost Per Invoice:", "$16.00  -->  $3.00", "(81% cheaper)"),
    ("Error Rate:", "4%  -->  0.5%", "(87% fewer errors)"),
    ("Days to Payment:", "14.6 days  -->  1 day", "(93% faster)"),
    ("Straight-Through:", "<20%  -->  85%+", "(no human touch needed)"),
]
for i, (label, change, pct) in enumerate(transform):
    top = Inches(3.7) + Inches(i * 0.6)
    text(slide, Inches(7.2), top, Inches(2.2), Inches(0.3), label, size=13, color=NAVY, bold=True)
    text(slide, Inches(9.4), top, Inches(2.5), Inches(0.3), change, size=13, color=DARK)
    text(slide, Inches(11.7), top, Inches(1.2), Inches(0.3), pct, size=12, color=GREEN, bold=True)

# ============================================================
# SLIDE 5: ROI - THE MONEY SLIDE
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
slide_header(slide, "The ROI: $462,000 Saved Per Year",
             "Based on your volume of 3,000 invoices per month")

# Annual cost comparison bar chart
chart_data3 = CategoryChartData()
chart_data3.categories = ["Annual AP Cost"]
chart_data3.add_series("Current (Manual)", (576000,))
chart_data3.add_series("With PayFlow AI", (114000,))
chart_data3.add_series("Your Savings", (462000,))

chart_frame3 = slide.shapes.add_chart(
    XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(0.8), Inches(1.5), Inches(5.5), Inches(3.5), chart_data3
)
chart3 = chart_frame3.chart
chart3.has_legend = True
chart3.legend.position = XL_LEGEND_POSITION.BOTTOM
chart3.legend.font.size = Pt(11)

chart3.series[0].format.fill.solid()
chart3.series[0].format.fill.fore_color.rgb = RED
chart3.series[1].format.fill.solid()
chart3.series[1].format.fill.fore_color.rgb = BLUE
chart3.series[2].format.fill.solid()
chart3.series[2].format.fill.fore_color.rgb = GREEN

plot3 = chart3.plots[0]
plot3.has_data_labels = True
plot3.data_labels.font.size = Pt(14)
plot3.data_labels.font.bold = True
plot3.data_labels.number_format = '$#,##0'
plot3.gap_width = 60

val_axis3 = chart3.value_axis
val_axis3.visible = False

# Right side - detailed savings
text(slide, Inches(6.8), Inches(1.5), Inches(5.5), Inches(0.4),
     "Savings Breakdown", size=22, color=NAVY, bold=True)

savings_table = [
    ["Category", "Current", "With AI", "Savings"],
    ["AP Staff Labor", "$360,000", "$84,000", "$276,000"],
    ["Error Correction", "$72,000", "$7,200", "$64,800"],
    ["Late Payments / Discounts", "$90,000", "$9,000", "$81,000"],
    ["Overhead", "$36,000", "$12,000", "$24,000"],
    ["Vendor Follow-ups", "$18,000", "$1,800", "$16,200"],
    ["TOTAL", "$576,000", "$114,000", "$462,000"],
]

rows, cols = len(savings_table), len(savings_table[0])
tbl = slide.shapes.add_table(rows, cols, Inches(6.8), Inches(2.0), Inches(5.7), Inches(2.8)).table
tbl.columns[0].width = Inches(2.0)
tbl.columns[1].width = Inches(1.1)
tbl.columns[2].width = Inches(1.1)
tbl.columns[3].width = Inches(1.5)

for r in range(rows):
    for c in range(cols):
        cell = tbl.cell(r, c)
        cell.text = savings_table[r][c]
        p = cell.text_frame.paragraphs[0]
        p.font.size = Pt(12)
        p.font.name = "Calibri"
        if r == 0:
            cell.fill.solid()
            cell.fill.fore_color.rgb = NAVY
            p.font.color.rgb = WHITE
            p.font.bold = True
        elif r == rows - 1:
            cell.fill.solid()
            cell.fill.fore_color.rgb = LIGHT_GREEN
            p.font.color.rgb = DARK_GREEN
            p.font.bold = True
        else:
            cell.fill.solid()
            cell.fill.fore_color.rgb = LIGHT_GRAY if r % 2 == 0 else WHITE
            p.font.color.rgb = DARK
            if c == 3:
                p.font.color.rgb = GREEN
                p.font.bold = True
        if c > 0:
            p.alignment = PP_ALIGN.CENTER

# Bottom callouts
stat_card(slide, Inches(0.8), Inches(5.3), Inches(2.8), Inches(1.5),
          "SOFTWARE COST", "$18,000/yr", "PayFlow AI subscription", BLUE)
stat_card(slide, Inches(3.9), Inches(5.3), Inches(2.8), Inches(1.5),
          "ANNUAL SAVINGS", "$462,000", "Net savings after software cost", GREEN)
stat_card(slide, Inches(7.0), Inches(5.3), Inches(2.8), Inches(1.5),
          "ROI", "2,467%", "Return on investment", GREEN)
stat_card(slide, Inches(10.1), Inches(5.3), Inches(2.8), Inches(1.5),
          "PAYBACK PERIOD", "2 Weeks", "Software pays for itself", GREEN)

# ============================================================
# SLIDE 6: 3-YEAR PROJECTION
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
slide_header(slide, "3-Year Financial Impact",
             "Cumulative savings assuming 5% annual invoice growth")

# 3-year line/bar chart
chart_data4 = CategoryChartData()
chart_data4.categories = ["Year 1", "Year 2", "Year 3"]
chart_data4.add_series("Manual AP Cost", (576000, 604800, 635040))
chart_data4.add_series("With PayFlow AI", (114000, 118200, 122400))
chart_data4.add_series("Cumulative Savings", (462000, 948600, 1461240))

chart_frame4 = slide.shapes.add_chart(
    XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(0.8), Inches(1.5), Inches(7.5), Inches(4.5), chart_data4
)
chart4 = chart_frame4.chart
chart4.has_legend = True
chart4.legend.position = XL_LEGEND_POSITION.BOTTOM
chart4.legend.font.size = Pt(12)

chart4.series[0].format.fill.solid()
chart4.series[0].format.fill.fore_color.rgb = RED
chart4.series[1].format.fill.solid()
chart4.series[1].format.fill.fore_color.rgb = BLUE
chart4.series[2].format.fill.solid()
chart4.series[2].format.fill.fore_color.rgb = GREEN

plot4 = chart4.plots[0]
plot4.has_data_labels = True
plot4.data_labels.font.size = Pt(12)
plot4.data_labels.font.bold = True
plot4.data_labels.number_format = '$#,##0'
plot4.gap_width = 60

val_axis4 = chart4.value_axis
val_axis4.visible = False

# Right side - 3 year summary
text(slide, Inches(8.8), Inches(1.5), Inches(4), Inches(0.4),
     "3-Year Summary", size=22, color=NAVY, bold=True)

yr_table = [
    ["", "Year 1", "Year 2", "Year 3"],
    ["Invoices/Month", "3,000", "3,150", "3,308"],
    ["Manual AP Cost", "$576K", "$605K", "$635K"],
    ["PayFlow AI Cost", "$114K", "$118K", "$122K"],
    ["Annual Savings", "$462K", "$487K", "$513K"],
    ["Cumulative Savings", "$462K", "$949K", "$1.46M"],
]

rows5, cols5 = len(yr_table), len(yr_table[0])
tbl5 = slide.shapes.add_table(rows5, cols5, Inches(8.8), Inches(2.0), Inches(4.0), Inches(2.5)).table
tbl5.columns[0].width = Inches(1.4)
for c in range(1, 4):
    tbl5.columns[c].width = Inches(0.87)

for r in range(rows5):
    for c in range(cols5):
        cell = tbl5.cell(r, c)
        cell.text = yr_table[r][c]
        p = cell.text_frame.paragraphs[0]
        p.font.size = Pt(11)
        p.font.name = "Calibri"
        if r == 0:
            cell.fill.solid()
            cell.fill.fore_color.rgb = NAVY
            p.font.color.rgb = WHITE
            p.font.bold = True
        elif r == rows5 - 1:
            cell.fill.solid()
            cell.fill.fore_color.rgb = LIGHT_GREEN
            p.font.bold = True
            p.font.color.rgb = DARK_GREEN
        else:
            cell.fill.solid()
            cell.fill.fore_color.rgb = LIGHT_GRAY if r % 2 == 0 else WHITE
            p.font.color.rgb = DARK
        if c > 0:
            p.alignment = PP_ALIGN.CENTER

# Big number
add_shape(slide, Inches(8.8), Inches(4.8), Inches(4.0), Inches(1.5), LIGHT_GREEN, GREEN)
text(slide, Inches(8.8), Inches(4.9), Inches(4.0), Inches(0.3),
     "3-Year Total Savings", size=14, color=DARK_GREEN, bold=True, align=PP_ALIGN.CENTER)
text(slide, Inches(8.8), Inches(5.2), Inches(4.0), Inches(0.7),
     "$1.46 Million", size=40, color=DARK_GREEN, bold=True, align=PP_ALIGN.CENTER)
text(slide, Inches(8.8), Inches(5.9), Inches(4.0), Inches(0.3),
     "Redirected to projects, equipment, and growth", size=11, color=DARK_GREEN, align=PP_ALIGN.CENTER)

# ============================================================
# SLIDE 7: CONSTRUCTION SPECIFIC - WHY GENERIC FAILS
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
slide_header(slide, "Why Generic AP Software Fails Construction",
             "Your AP process has requirements that Stampli, Vic.ai, and Bill.com can't handle")

# Spider/radar chart substitute - comparison table with checkmarks
comp_data = [
    ["Capability", "Generic AP\n(Stampli, Vic.ai)", "PayFlow AI"],
    ["Read standard invoices", "Yes", "Yes"],
    ["GL code prediction", "Yes", "Yes"],
    ["Approval workflows", "Yes", "Yes"],
    ["AIA G702/G703 billing", "No", "Yes"],
    ["Lien waiver tracking", "No", "Yes"],
    ["Retention management", "No", "Yes"],
    ["5D job costing", "No", "Yes"],
    ["Change order tracking", "No", "Yes"],
    ["COI expiration alerts", "No", "Yes"],
    ["Prevailing wage verification", "No", "Yes"],
    ["Multi-entity routing", "No", "Yes"],
    ["Sage integration", "Limited", "Native"],
    ["Procore integration", "No", "Native"],
]

rows6, cols6 = len(comp_data), len(comp_data[0])
tbl6 = slide.shapes.add_table(rows6, cols6, Inches(0.8), Inches(1.5), Inches(8.0), Inches(5.5)).table
tbl6.columns[0].width = Inches(3.5)
tbl6.columns[1].width = Inches(2.2)
tbl6.columns[2].width = Inches(2.3)

for r in range(rows6):
    for c in range(cols6):
        cell = tbl6.cell(r, c)
        cell.text = comp_data[r][c]
        p = cell.text_frame.paragraphs[0]
        p.font.size = Pt(13)
        p.font.name = "Calibri"
        if r == 0:
            cell.fill.solid()
            cell.fill.fore_color.rgb = NAVY
            p.font.color.rgb = WHITE
            p.font.bold = True
        else:
            cell.fill.solid()
            cell.fill.fore_color.rgb = LIGHT_GRAY if r % 2 == 0 else WHITE
            p.font.color.rgb = DARK
            if c == 1 and comp_data[r][c] == "No":
                p.font.color.rgb = RED
                p.font.bold = True
            elif c == 2 and comp_data[r][c] == "Yes":
                p.font.color.rgb = GREEN
                p.font.bold = True
            elif c == 2 and comp_data[r][c] == "Native":
                p.font.color.rgb = GREEN
                p.font.bold = True
        if c > 0:
            p.alignment = PP_ALIGN.CENTER

# Right side callout
add_shape(slide, Inches(9.3), Inches(1.5), Inches(3.5), Inches(5.5), LIGHT_BLUE, BLUE)
text(slide, Inches(9.5), Inches(1.7), Inches(3.1), Inches(0.4),
     "The Gap", size=22, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
gap_items = [
    "Generic AP tools handle 3 out of 13 construction requirements.",
    "PayFlow AI handles all 13.",
    "That's not a feature gap — it's a different product for a different industry.",
    "Your AP team currently fills this gap manually. That's where the $576K goes.",
]
bullets(slide, Inches(9.5), Inches(2.2), Inches(3.1), Inches(4.5), gap_items, size=14, color=NAVY, spacing=16)

# ============================================================
# SLIDE 8: INTEGRATIONS
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
slide_header(slide, "Plugs Into Your Existing Systems",
             "No rip-and-replace. Works with Sage and Procore on day one.")

integrations = [
    ("Sage", "Accounting", [
        "Two-way sync of chart of accounts",
        "Import vendors, jobs, cost codes",
        "Push approved invoices directly",
        "Sync payment status back",
    ], NAVY),
    ("Procore", "Project Management", [
        "Pull POs, budgets, commitments",
        "Import change orders",
        "Three-way match against POs",
        "Real-time contract tracking",
    ], BLUE),
    ("Email", "Invoice Ingestion", [
        "Dedicated inbox per company",
        "AI identifies invoices vs. spam",
        "Auto-processes on arrival",
        "Vendors email directly",
    ], GREEN),
    ("Banking", "Payments", [
        "Schedule approved payments",
        "ACH, check, or card",
        "Auto-reconciliation",
        "Cash flow visibility",
    ], ORANGE),
]

for i, (name, subtitle, items, color) in enumerate(integrations):
    left = Inches(0.5) + Inches(i * 3.2)
    top = Inches(1.6)
    card = add_shape(slide, left, top, Inches(2.9), Inches(4.0), WHITE, color)
    add_rect(slide, left + Inches(0.05), top + Inches(0.05), Inches(2.8), Inches(0.7), color)
    text(slide, left + Inches(0.1), top + Inches(0.08), Inches(2.7), Inches(0.4),
         name, size=22, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    text(slide, left + Inches(0.1), top + Inches(0.4), Inches(2.7), Inches(0.3),
         subtitle, size=11, color=RGBColor(0xd0, 0xe0, 0xf0), align=PP_ALIGN.CENTER)
    bullets(slide, left + Inches(0.2), top + Inches(0.95), Inches(2.5), Inches(2.8),
            items, size=13, color=GRAY, spacing=10)

# Bottom flow
add_shape(slide, Inches(0.8), Inches(6.0), Inches(11.7), Inches(0.9), LIGHT_BLUE, BLUE)
text(slide, Inches(1.0), Inches(6.1), Inches(11.3), Inches(0.7),
     "Invoice arrives via email  -->  AI extracts & codes  -->  Matches against Procore PO  -->  Routes for approval  -->  Posts to Sage  -->  Payment scheduled",
     size=14, color=NAVY, bold=True, align=PP_ALIGN.CENTER)

# ============================================================
# SLIDE 9: ACCURACY & TRUST
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
slide_header(slide, "AI Accuracy You Can Trust",
             "Built-in safeguards ensure nothing gets posted without verification")

# Accuracy chart
chart_data5 = CategoryChartData()
chart_data5.categories = ["Invoice\nNumber", "Date", "Vendor\nName", "Total\nAmount",
                          "Line Items", "PO Number", "GL Code\nPrediction"]
chart_data5.add_series("AI Accuracy %", (99.5, 98.5, 97.5, 99.5, 96, 97, 92))

chart_frame5 = slide.shapes.add_chart(
    XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(0.8), Inches(1.5), Inches(6.5), Inches(3.5), chart_data5
)
chart5 = chart_frame5.chart
chart5.has_legend = False

series5 = chart5.series[0]
series5.format.fill.solid()
series5.format.fill.fore_color.rgb = BLUE

plot5 = chart5.plots[0]
plot5.has_data_labels = True
plot5.data_labels.font.size = Pt(11)
plot5.data_labels.font.bold = True
plot5.data_labels.number_format = '0.0"%"'
plot5.data_labels.position = XL_LABEL_POSITION.OUTSIDE_END
plot5.gap_width = 60

val_axis5 = chart5.value_axis
val_axis5.minimum_scale = 80
val_axis5.maximum_scale = 100
val_axis5.tick_labels.font.size = Pt(10)

cat_axis5 = chart5.category_axis
cat_axis5.tick_labels.font.size = Pt(10)

# Right side - trust features
text(slide, Inches(7.8), Inches(1.5), Inches(5), Inches(0.4),
     "Built-In Safety Net", size=22, color=NAVY, bold=True)

trust_items = [
    "Confidence scoring on every field — you see exactly how sure the AI is",
    "Low-confidence invoices route to humans automatically — never auto-posted blindly",
    "Every AI decision logged in full audit trail — who, what, when, why",
    "Math validation — line items must sum to totals or it's flagged",
    "Duplicate detection — same invoice number, same amount, same vendor",
    "Human override on any field — your team has final say, always",
    "Shadow mode — run AI alongside your team before going live",
]
bullets(slide, Inches(8.0), Inches(2.0), Inches(4.8), Inches(3.5), trust_items, size=13, color=GRAY, spacing=10)

# Bottom bar
add_shape(slide, Inches(0.8), Inches(5.3), Inches(11.7), Inches(1.5), LIGHT_GREEN, GREEN)
text(slide, Inches(1.0), Inches(5.4), Inches(11.3), Inches(0.4),
     "The AI handles the 85% that's routine. Your team focuses on the 15% that needs judgment.",
     size=18, color=DARK_GREEN, bold=True, align=PP_ALIGN.CENTER)
text(slide, Inches(1.0), Inches(5.9), Inches(11.3), Inches(0.6),
     "That's not replacing your team — it's giving them superpowers. One AP clerk with PayFlow AI\noutperforms a team of five doing it manually.",
     size=14, color=DARK_GREEN, align=PP_ALIGN.CENTER)

# ============================================================
# SLIDE 10: SECURITY
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
slide_header(slide, "Enterprise-Grade Security",
             "Your financial data is protected at every layer")

security = [
    ("Encryption", "AES-256 at rest\nTLS 1.2+ in transit\nEnd-to-end protection"),
    ("AI Data Privacy", "Invoices processed via API only\nNever used to train AI models\nYour data stays yours"),
    ("Access Control", "Role-based permissions\nAP clerk, PM, Controller, CFO\nEach sees only what they need"),
    ("Audit Trail", "Every action logged\nWho uploaded, extracted, approved\nFull compliance history"),
    ("SOC 2", "Enterprise security controls\nAnnual third-party audits\nIndustry standard compliance"),
    ("Data Ownership", "You own your data, always\nFull export available anytime\nPortable if you leave"),
]

for i, (title, desc) in enumerate(security):
    col = i % 3
    row = i // 3
    left = Inches(0.8) + Inches(col * 4.05)
    top = Inches(1.5) + Inches(row * 2.5)
    card = add_shape(slide, left, top, Inches(3.7), Inches(2.1), LIGHT_GRAY, MEDIUM_GRAY)
    text(slide, left + Inches(0.2), top + Inches(0.2), Inches(3.3), Inches(0.35),
         title, size=18, color=NAVY, bold=True)
    box = slide.shapes.add_textbox(left + Inches(0.2), top + Inches(0.6), Inches(3.3), Inches(1.3))
    tf = box.text_frame
    tf.word_wrap = True
    for line in desc.split("\n"):
        p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(13)
        p.font.color.rgb = GRAY
        p.font.name = "Calibri"
        p.space_after = Pt(4)

# ============================================================
# SLIDE 11: IMPLEMENTATION TIMELINE
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
slide_header(slide, "Live in 8 Weeks — Zero Disruption",
             "Your team keeps working normally while we set up in parallel")

# Gantt-like timeline
phases = [
    ("Weeks 1-2", "Setup & Connect", "Connect Sage & Procore. Import\nchart of accounts, vendors,\njobs, cost codes.", BLUE, 0, 2),
    ("Weeks 3-4", "AI Calibration", "Process your real invoices.\nTune AI to your patterns.\nVerify accuracy.", BLUE, 2, 2),
    ("Weeks 5-6", "Shadow Mode", "AI runs alongside your team.\nCompare results side by side.\nBuild confidence.", GREEN, 4, 2),
    ("Weeks 7-8", "Go Live", "AI handles routine invoices.\nTeam focuses on exceptions.\nStart saving immediately.", GREEN, 6, 2),
]

# Timeline bar background
add_rect(slide, Inches(0.8), Inches(1.5), Inches(11.7), Inches(0.5), MEDIUM_GRAY)
for i in range(9):
    x = Inches(0.8) + Inches(i * (11.7 / 8))
    text(slide, x - Inches(0.2), Inches(1.52), Inches(0.5), Inches(0.4),
         f"W{i}" if i > 0 else "Start", size=9, color=GRAY, align=PP_ALIGN.CENTER)

for i, (timeline, title, desc, color, start, length) in enumerate(phases):
    # Gantt bar
    bar_left = Inches(0.8) + Inches(start * (11.7 / 8))
    bar_width = Inches(length * (11.7 / 8))
    bar = add_rect(slide, bar_left, Inches(1.5), bar_width, Inches(0.5), color)
    text(slide, bar_left, Inches(1.52), bar_width, Inches(0.45),
         f"{timeline}: {title}", size=11, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

    # Detail cards below
    left = Inches(0.5) + Inches(i * 3.15)
    top = Inches(2.3)
    card = add_shape(slide, left, top, Inches(2.95), Inches(2.5), WHITE, color)
    add_rect(slide, left + Inches(0.05), top + Inches(0.05), Inches(2.85), Inches(0.45), color)
    text(slide, left + Inches(0.1), top + Inches(0.08), Inches(2.75), Inches(0.4),
         timeline, size=14, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    text(slide, left + Inches(0.15), top + Inches(0.6), Inches(2.65), Inches(0.35),
         title, size=16, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
    box = slide.shapes.add_textbox(left + Inches(0.15), top + Inches(1.0), Inches(2.65), Inches(1.3))
    tf = box.text_frame
    tf.word_wrap = True
    for line in desc.split("\n"):
        p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(12)
        p.font.color.rgb = GRAY
        p.font.name = "Calibri"
        p.alignment = PP_ALIGN.CENTER
        p.space_after = Pt(3)

# Bottom guarantees
guarantees = [
    ("Zero Disruption", "Your team keeps working normally.\nAI runs in parallel until ready."),
    ("Immediate Value", "See ROI from the first week\nof go-live. No 6-month wait."),
    ("Risk-Free", "Shadow mode proves accuracy\nbefore you commit. Cancel anytime."),
]
for i, (gtitle, gdesc) in enumerate(guarantees):
    left = Inches(0.8) + Inches(i * 4.1)
    card = add_shape(slide, left, Inches(5.2), Inches(3.7), Inches(1.6), LIGHT_BLUE, BLUE)
    text(slide, left + Inches(0.2), Inches(5.3), Inches(3.3), Inches(0.35),
         gtitle, size=16, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
    box = slide.shapes.add_textbox(left + Inches(0.2), Inches(5.7), Inches(3.3), Inches(0.9))
    tf = box.text_frame
    tf.word_wrap = True
    for line in gdesc.split("\n"):
        p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(12)
        p.font.color.rgb = NAVY
        p.font.name = "Calibri"
        p.alignment = PP_ALIGN.CENTER

# ============================================================
# SLIDE 12: PRICING
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
slide_header(slide, "Investment")

tiers = [
    ("Starter", "$500", "/month", [
        "Up to 1,000 invoices/mo",
        "5 users",
        "Sage integration",
        "AI extraction & coding",
        "Email ingestion",
        "Standard support",
    ], BLUE, False),
    ("Professional", "$1,500", "/month", [
        "Up to 5,000 invoices/mo",
        "Unlimited users",
        "Sage + Procore integration",
        "AIA billing & lien waivers",
        "Retention tracking",
        "Fraud detection",
        "Priority support",
    ], GREEN, True),
    ("Enterprise", "Custom", "", [
        "Unlimited invoices",
        "Unlimited users",
        "All integrations",
        "Custom workflows",
        "Dedicated account manager",
        "SLA guarantee",
        "On-site training",
    ], NAVY, False),
]

for i, (name, price, period, items, color, recommended) in enumerate(tiers):
    left = Inches(1.0) + Inches(i * 3.9)
    top = Inches(1.3)
    h = Inches(5.3)

    if recommended:
        add_rect(slide, left - Inches(0.05), top - Inches(0.35), Inches(3.6), Inches(0.35), color)
        text(slide, left - Inches(0.05), top - Inches(0.35), Inches(3.6), Inches(0.3),
             "RECOMMENDED", size=11, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

    card = add_shape(slide, left, top, Inches(3.5), h, WHITE, color)
    if recommended:
        card.line.width = Pt(3)

    text(slide, left + Inches(0.1), top + Inches(0.2), Inches(3.3), Inches(0.4),
         name, size=22, color=color, bold=True, align=PP_ALIGN.CENTER)
    text(slide, left + Inches(0.1), top + Inches(0.65), Inches(3.3), Inches(0.6),
         price, size=42, color=DARK, bold=True, align=PP_ALIGN.CENTER)
    if period:
        text(slide, left + Inches(0.1), top + Inches(1.15), Inches(3.3), Inches(0.3),
             period, size=14, color=GRAY, align=PP_ALIGN.CENTER)

    start_y = top + Inches(1.55)
    for j, item in enumerate(items):
        text(slide, left + Inches(0.3), start_y + Inches(j * 0.4), Inches(2.9), Inches(0.35),
             "  " + item, size=13, color=GRAY)

# Bottom note
text(slide, Inches(0.8), Inches(6.8), Inches(11.5), Inches(0.4),
     "All plans include: AI invoice extraction, GL coding, approval workflows, duplicate detection, and full audit trail.",
     size=13, color=GRAY, align=PP_ALIGN.CENTER)

# ============================================================
# SLIDE 13: NEXT STEPS / CTA
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, NAVY)
add_rect(slide, Inches(0), Inches(3.0), Inches(13.333), Inches(0.06), BLUE)

text(slide, Inches(1), Inches(0.8), Inches(11.3), Inches(0.8),
     "Ready to Save $462,000/Year?", size=44, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
text(slide, Inches(2), Inches(1.7), Inches(9.3), Inches(0.5),
     "Three steps to transform your AP process — zero risk:",
     size=20, color=RGBColor(0xa0, 0xce, 0xf0), align=PP_ALIGN.CENTER)

cta_steps = [
    ("1", "Share 50 Sample Invoices", "We run them through the AI and show you extraction results.\nNo commitment. No cost. Just proof."),
    ("2", "2-Week Shadow Pilot", "AI processes your real invoices alongside your team.\nYou compare accuracy. We prove the value."),
    ("3", "Go Live With Confidence", "Only when you've seen the numbers.\nStart saving from day one. Cancel anytime."),
]

for i, (num, title, desc) in enumerate(cta_steps):
    top = Inches(3.4) + Inches(i * 1.2)
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(2.5), top, Inches(0.6), Inches(0.6))
    circle.fill.solid()
    circle.fill.fore_color.rgb = BLUE
    circle.line.fill.background()
    tf = circle.text_frame
    p = tf.paragraphs[0]
    p.text = num
    p.font.size = Pt(20)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    text(slide, Inches(3.3), top + Inches(0.02), Inches(7), Inches(0.35),
         title, size=22, color=WHITE, bold=True)
    box = slide.shapes.add_textbox(Inches(3.3), top + Inches(0.4), Inches(7), Inches(0.6))
    tf2 = box.text_frame
    tf2.word_wrap = True
    for line in desc.split("\n"):
        p2 = tf2.add_paragraph()
        p2.text = line
        p2.font.size = Pt(13)
        p2.font.color.rgb = RGBColor(0xcb, 0xd5, 0xe0)
        p2.font.name = "Calibri"

text(slide, Inches(1), Inches(6.8), Inches(11.3), Inches(0.4),
     "PayFlow AI  |  AI-Powered AP Automation for Construction  |  payflowai.com",
     size=14, color=RGBColor(0x71, 0x80, 0x96), align=PP_ALIGN.CENTER)

# SAVE
output = "C:/users/25badmin/projects/accounts-payable-research/PayFlow_AI_Sales_Deck_v2.pptx"
prs.save(output)
print(f"Presentation saved: {output}")
