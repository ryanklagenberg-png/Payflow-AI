"""
AI SaaS Opportunities — CEO Presentation
PowerPoint showing high-revenue AI product opportunities
with market analysis, revenue projections, and strategic recommendations
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_LABEL_POSITION
from pptx.chart.data import CategoryChartData

# ── Brand Colors ──
NAVY = RGBColor(0x0B, 0x1D, 0x3A)
DARK_BLUE = RGBColor(0x12, 0x2B, 0x4F)
ACCENT_BLUE = RGBColor(0x1E, 0x90, 0xFF)
LIGHT_BLUE = RGBColor(0x4D, 0xA8, 0xFF)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xE8, 0xEC, 0xF1)
MEDIUM_GRAY = RGBColor(0x8A, 0x95, 0xA5)
DARK_GRAY = RGBColor(0x2C, 0x3E, 0x50)
GREEN = RGBColor(0x2E, 0xCC, 0x71)
DARK_GREEN = RGBColor(0x1A, 0x9B, 0x50)
ORANGE = RGBColor(0xF3, 0x9C, 0x12)
RED = RGBColor(0xE7, 0x4C, 0x3C)
PURPLE = RGBColor(0x9B, 0x59, 0xB6)
TEAL = RGBColor(0x00, 0xB4, 0xD8)
GOLD = RGBColor(0xF1, 0xC4, 0x0F)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)


def add_bg(slide, color=NAVY):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape(slide, left, top, width, height, fill_color, border_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(2)
    else:
        shape.line.fill.background()
    return shape


def add_text(slide, left, top, width, height, text, font_size=18, color=WHITE, bold=False, alignment=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = alignment
    return tf


def add_arrow(slide, left, top, width, height, color=ACCENT_BLUE):
    shape = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_down_arrow(slide, left, top, width, height, color=ACCENT_BLUE):
    shape = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


# ═══════════════════════════════════════════════
# SLIDE 1: Title
# ═══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.06), GOLD)

add_text(slide, Inches(1), Inches(1.2), Inches(11), Inches(1.2),
         "AI-Powered SaaS Opportunities", 48, WHITE, True, PP_ALIGN.CENTER)
add_text(slide, Inches(1), Inches(2.4), Inches(11), Inches(0.8),
         "High-Margin Products That Solve Expensive Problems", 28, GOLD, False, PP_ALIGN.CENTER)
add_text(slide, Inches(2), Inches(3.8), Inches(9), Inches(0.6),
         "Strategic Product Portfolio for Maximum Revenue", 18, MEDIUM_GRAY, False, PP_ALIGN.CENTER)

# Key stat boxes
stats = [
    ("$50B+", "AI SaaS Market\nby 2028", ACCENT_BLUE),
    ("80%", "Cost Reduction\nPotential", GREEN),
    ("10-50x", "ROI for\nCustomers", GOLD),
    ("$0.02-0.08", "AI Cost Per\nTransaction", TEAL),
]
x = 1.8
for value, label, color in stats:
    add_shape(slide, Inches(x), Inches(4.8), Inches(2.2), Inches(1.5), DARK_BLUE, color)
    add_text(slide, Inches(x), Inches(4.95), Inches(2.2), Inches(0.6),
             value, 30, color, True, PP_ALIGN.CENTER)
    add_text(slide, Inches(x + 0.1), Inches(5.5), Inches(2.0), Inches(0.7),
             label, 13, RGBColor(0xCC, 0xCC, 0xCC), False, PP_ALIGN.CENTER)
    x += 2.55

add_text(slide, Inches(1), Inches(6.8), Inches(11), Inches(0.4),
         "CEO Strategic Briefing  |  March 2026", 14, MEDIUM_GRAY, False, PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════
# SLIDE 2: The Thesis
# ═══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.06), GOLD)

add_text(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.7),
         "The Thesis: Why These Products Print Money", 34, WHITE, True)

# Three pillars
pillars = [
    ("Find Expensive Manual Processes",
     "Every industry has workers doing repetitive document processing — reading, extracting, matching, coding, routing. These tasks cost companies $15-50 per transaction in human labor.",
     "The Pain", ACCENT_BLUE),
    ("Replace with AI at 99% Lower Cost",
     "Claude AI can read documents, extract structured data, and make decisions at $0.02-0.08 per transaction. What costs a company $25 in labor costs us $0.06 in API calls.",
     "The Solution", GREEN),
    ("Charge 10% of What They're Spending",
     "If a company spends $25/transaction manually and we charge $2.50/transaction, they save 90% and we make 95%+ gross margins. Everyone wins.",
     "The Business", GOLD),
]

y = 1.3
for title, desc, label, color in pillars:
    add_shape(slide, Inches(0.5), Inches(y), Inches(12.3), Inches(1.6), DARK_BLUE, color)
    add_shape(slide, Inches(0.6), Inches(y + 0.15), Inches(1.8), Inches(0.5), color)
    add_text(slide, Inches(0.6), Inches(y + 0.17), Inches(1.8), Inches(0.45),
             label, 14, WHITE, True, PP_ALIGN.CENTER)
    add_text(slide, Inches(2.6), Inches(y + 0.1), Inches(9.8), Inches(0.4),
             title, 20, WHITE, True)
    add_text(slide, Inches(2.6), Inches(y + 0.55), Inches(9.8), Inches(0.9),
             desc, 14, RGBColor(0xCC, 0xCC, 0xCC))
    y += 1.75

# Bottom formula
add_shape(slide, Inches(1.5), Inches(6.6), Inches(10.3), Inches(0.7), DARK_BLUE, GOLD)
add_text(slide, Inches(1.5), Inches(6.65), Inches(10.3), Inches(0.6),
         "The Formula:  Expensive Manual Process  +  AI Automation  +  Vertical Expertise  =  High-Margin SaaS",
         16, GOLD, True, PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════
# SLIDE 3: Product Portfolio Overview
# ═══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.06), GOLD)

add_text(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.7),
         "Product Portfolio — 8 High-Revenue Opportunities", 32, WHITE, True)
add_text(slide, Inches(0.8), Inches(0.85), Inches(11), Inches(0.4),
         "Ranked by revenue potential, speed to market, and strategic synergy", 15, MEDIUM_GRAY)

products = [
    ("1", "PayFlow AI", "Construction AP Automation", "$1.5-10K/mo", "700K firms", "READY", GREEN),
    ("2", "AI Contract Review", "Construction Contract Analysis", "$2-10K/mo", "700K firms", "HIGH", ACCENT_BLUE),
    ("3", "AI Permit Expediting", "Automated Permit Applications", "$500-2K/permit", "1M+ permits/yr", "HIGH", ACCENT_BLUE),
    ("4", "AI Insurance Claims", "Claims Processing & Fraud Detection", "$5-50K/mo", "6K carriers", "HIGH", ACCENT_BLUE),
    ("5", "AI Medical Coding", "ICD-10/CPT Code Prediction", "$5-30K/mo", "1M+ providers", "MEDIUM", ORANGE),
    ("6", "AI Lien Management", "Notice Tracking & Waiver Exchange", "$500-2K/mo", "700K firms", "HIGH", ACCENT_BLUE),
    ("7", "AI Estimating", "Automated Bid & Proposal Generation", "$1-5K/mo", "700K firms", "MEDIUM", ORANGE),
    ("8", "AI Property Mgmt AP", "Multi-Property Invoice Processing", "$1-5K/mo", "300K firms", "HIGH", ACCENT_BLUE),
]

# Header
add_shape(slide, Inches(0.3), Inches(1.35), Inches(12.7), Inches(0.45), ACCENT_BLUE)
headers = [("#", 0.35, 0.35), ("Product", 0.75, 2.2), ("What It Does", 3.0, 2.8),
           ("Revenue/Customer", 5.9, 1.5), ("Market Size", 7.5, 1.3),
           ("Speed to Build", 8.9, 1.2), ("Tier", 10.2, 1.0)]
for text, x, w in headers:
    add_text(slide, Inches(x), Inches(1.38), Inches(w), Inches(0.4),
             text, 12, WHITE, True)

y = 1.85
for rank, name, desc, rev, market, speed, color in products:
    bg = DARK_BLUE if int(rank) % 2 == 1 else RGBColor(0x15, 0x2D, 0x52)
    add_shape(slide, Inches(0.3), Inches(y), Inches(12.7), Inches(0.6), bg)

    tier_color = GREEN if speed == "READY" else (ACCENT_BLUE if speed == "HIGH" else ORANGE)

    add_text(slide, Inches(0.4), Inches(y + 0.1), Inches(0.3), Inches(0.4),
             rank, 14, GOLD, True, PP_ALIGN.CENTER)
    add_text(slide, Inches(0.75), Inches(y + 0.1), Inches(2.2), Inches(0.4),
             name, 14, WHITE, True)
    add_text(slide, Inches(3.0), Inches(y + 0.1), Inches(2.8), Inches(0.4),
             desc, 11, RGBColor(0xCC, 0xCC, 0xCC))
    add_text(slide, Inches(5.9), Inches(y + 0.1), Inches(1.5), Inches(0.4),
             rev, 12, GREEN, True)
    add_text(slide, Inches(7.5), Inches(y + 0.1), Inches(1.3), Inches(0.4),
             market, 11, MEDIUM_GRAY)
    add_shape(slide, Inches(9.0), Inches(y + 0.15), Inches(1.0), Inches(0.3), tier_color)
    add_text(slide, Inches(9.0), Inches(y + 0.15), Inches(1.0), Inches(0.3),
             speed, 10, WHITE, True, PP_ALIGN.CENTER)

    y += 0.63

# Legend
add_shape(slide, Inches(0.3), Inches(6.95), Inches(12.7), Inches(0.4), DARK_BLUE)
add_shape(slide, Inches(0.5), Inches(7.02), Inches(0.7), Inches(0.25), GREEN)
add_text(slide, Inches(1.25), Inches(7.0), Inches(1.5), Inches(0.3),
         "Ready to build now", 11, MEDIUM_GRAY)
add_shape(slide, Inches(3.0), Inches(7.02), Inches(0.7), Inches(0.25), ACCENT_BLUE)
add_text(slide, Inches(3.75), Inches(7.0), Inches(2.5), Inches(0.3),
         "Fast to build, reuses AI engine", 11, MEDIUM_GRAY)
add_shape(slide, Inches(6.5), Inches(7.02), Inches(0.7), Inches(0.25), ORANGE)
add_text(slide, Inches(7.25), Inches(7.0), Inches(3.0), Inches(0.3),
         "Larger build, higher regulatory bar", 11, MEDIUM_GRAY)


# ═══════════════════════════════════════════════
# SLIDE 4: Product 1 — PayFlow AI (Flagship)
# ═══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.06), GREEN)

add_shape(slide, Inches(0.5), Inches(0.2), Inches(0.5), Inches(0.5), GREEN)
add_text(slide, Inches(0.5), Inches(0.22), Inches(0.5), Inches(0.45),
         "1", 24, WHITE, True, PP_ALIGN.CENTER)
add_text(slide, Inches(1.2), Inches(0.25), Inches(6), Inches(0.6),
         "PayFlow AI — Construction AP Automation", 30, WHITE, True)
add_shape(slide, Inches(9.5), Inches(0.25), Inches(1.5), Inches(0.4), GREEN)
add_text(slide, Inches(9.5), Inches(0.27), Inches(1.5), Inches(0.35),
         "FLAGSHIP", 14, WHITE, True, PP_ALIGN.CENTER)
add_shape(slide, Inches(11.2), Inches(0.25), Inches(1.8), Inches(0.4), DARK_GREEN)
add_text(slide, Inches(11.2), Inches(0.27), Inches(1.8), Inches(0.35),
         "READY TO BUILD", 12, WHITE, True, PP_ALIGN.CENTER)

# Left: The Problem
add_shape(slide, Inches(0.3), Inches(1.0), Inches(6.2), Inches(2.8), DARK_BLUE, RED)
add_text(slide, Inches(0.3), Inches(1.05), Inches(6.2), Inches(0.4),
         "The Problem", 18, RED, True, PP_ALIGN.CENTER)
problems = [
    "Construction companies process 1,000-10,000 invoices per month",
    "Manual AP processing costs $15-25 per invoice in labor",
    "5-dimensional job coding unique to construction (GL + Job + Phase + Cost Code + Cost Type)",
    "AIA progress billing, lien waivers, retention — generic AP tools can't handle these",
    "No existing product combines AI + construction-specific workflows",
    "Result: AP departments are the biggest bottleneck in construction finance",
]
y = 1.5
for item in problems:
    add_text(slide, Inches(0.5), Inches(y), Inches(5.8), Inches(0.35),
             "  " + item, 11, RGBColor(0xCC, 0xCC, 0xCC))
    add_shape(slide, Inches(0.5), Inches(y + 0.05), Inches(0.12), Inches(0.12), RED)
    y += 0.37

# Right: The Opportunity
add_shape(slide, Inches(6.8), Inches(1.0), Inches(6.2), Inches(2.8), DARK_BLUE, GREEN)
add_text(slide, Inches(6.8), Inches(1.05), Inches(6.2), Inches(0.4),
         "The Opportunity", 18, GREEN, True, PP_ALIGN.CENTER)
opps = [
    "700,000+ construction firms in the US",
    "$1.5-10K/month per customer ($18-120K ARR)",
    "Our cost: $0.06/invoice. Customer saves: $15-25/invoice.",
    "95%+ gross margins after AI processing costs",
    "Integrates with KOJO, Procore, Sage — tools they already use",
    "Procore App Marketplace = free distribution to 16,000 companies",
]
y = 1.5
for item in opps:
    add_text(slide, Inches(7.0), Inches(y), Inches(5.8), Inches(0.35),
             "  " + item, 11, RGBColor(0xCC, 0xCC, 0xCC))
    add_shape(slide, Inches(7.0), Inches(y + 0.05), Inches(0.12), Inches(0.12), GREEN)
    y += 0.37

# Revenue projection
add_shape(slide, Inches(0.3), Inches(4.0), Inches(12.7), Inches(0.4), ACCENT_BLUE)
add_text(slide, Inches(0.3), Inches(4.02), Inches(12.7), Inches(0.35),
         "Revenue Projection — PayFlow AI", 16, WHITE, True, PP_ALIGN.CENTER)

rev_data = [
    ("Year 1", "5 customers", "$7,500/mo", "$90K ARR", "Pilot + first sales"),
    ("Year 2", "25 customers", "$50,000/mo", "$600K ARR", "Procore marketplace traction"),
    ("Year 3", "100 customers", "$250,000/mo", "$3M ARR", "Sales team, enterprise tier"),
    ("Year 5", "500 customers", "$1.5M/mo", "$18M ARR", "Market leader in construction AP"),
]

add_shape(slide, Inches(0.3), Inches(4.5), Inches(12.7), Inches(0.4), RGBColor(0x15, 0x2D, 0x52))
cols = [("", 0.4, 1.0), ("Customers", 1.5, 1.3), ("MRR", 3.0, 1.5),
        ("ARR", 4.7, 1.3), ("Stage", 6.2, 3.0)]
for text, x, w in cols:
    add_text(slide, Inches(x), Inches(4.52), Inches(w), Inches(0.35),
             text, 11, ACCENT_BLUE, True)

y = 4.95
for year, cust, mrr, arr, stage in rev_data:
    bg = DARK_BLUE if rev_data.index((year, cust, mrr, arr, stage)) % 2 == 0 else RGBColor(0x15, 0x2D, 0x52)
    add_shape(slide, Inches(0.3), Inches(y), Inches(12.7), Inches(0.42), bg)
    add_text(slide, Inches(0.4), Inches(y + 0.05), Inches(1.0), Inches(0.3),
             year, 12, GOLD, True)
    add_text(slide, Inches(1.5), Inches(y + 0.05), Inches(1.3), Inches(0.3),
             cust, 12, WHITE)
    add_text(slide, Inches(3.0), Inches(y + 0.05), Inches(1.5), Inches(0.3),
             mrr, 12, GREEN, True)
    add_text(slide, Inches(4.7), Inches(y + 0.05), Inches(1.3), Inches(0.3),
             arr, 13, GREEN, True)
    add_text(slide, Inches(6.2), Inches(y + 0.05), Inches(3.0), Inches(0.3),
             stage, 11, MEDIUM_GRAY)
    y += 0.44

# Status
add_shape(slide, Inches(0.3), Inches(6.9), Inches(12.7), Inches(0.45), DARK_GREEN)
add_text(slide, Inches(0.3), Inches(6.93), Inches(12.7), Inches(0.4),
         "STATUS: Architecture complete. 8 AI prompts written. Pilot customer identified (Mark III Construction). Ready for Phase 1 development.",
         13, WHITE, True, PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════
# SLIDE 5: Product 2 — AI Contract Review
# ═══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.06), ACCENT_BLUE)

add_shape(slide, Inches(0.5), Inches(0.2), Inches(0.5), Inches(0.5), ACCENT_BLUE)
add_text(slide, Inches(0.5), Inches(0.22), Inches(0.5), Inches(0.45),
         "2", 24, WHITE, True, PP_ALIGN.CENTER)
add_text(slide, Inches(1.2), Inches(0.25), Inches(8), Inches(0.6),
         "AI Contract Review — Construction Contract Analysis", 30, WHITE, True)

# Two columns
add_shape(slide, Inches(0.3), Inches(1.0), Inches(6.2), Inches(3.2), DARK_BLUE, RED)
add_text(slide, Inches(0.3), Inches(1.05), Inches(6.2), Inches(0.4),
         "The Problem", 18, RED, True, PP_ALIGN.CENTER)
contract_problems = [
    "Construction companies sign hundreds of contracts per year",
    "Legal review costs $300-500/hour — most companies skip it",
    "Subcontractors sign contracts without fully reading them",
    "Unfavorable terms discovered only when disputes arise",
    "Pay-when-paid clauses, indemnification traps, insurance gaps",
    "One bad contract clause can cost $50K-500K in a dispute",
    "No AI tool exists specifically for construction contracts",
]
y = 1.5
for item in contract_problems:
    add_text(slide, Inches(0.5), Inches(y), Inches(5.8), Inches(0.35),
             "  " + item, 11, RGBColor(0xCC, 0xCC, 0xCC))
    add_shape(slide, Inches(0.5), Inches(y + 0.05), Inches(0.12), Inches(0.12), RED)
    y += 0.37

add_shape(slide, Inches(6.8), Inches(1.0), Inches(6.2), Inches(3.2), DARK_BLUE, GREEN)
add_text(slide, Inches(6.8), Inches(1.05), Inches(6.2), Inches(0.4),
         "What AI Contract Review Does", 18, GREEN, True, PP_ALIGN.CENTER)
contract_features = [
    "Upload any subcontract, vendor agreement, or change order",
    "AI reads the entire document in seconds",
    "Flags unfavorable clauses: pay-when-paid, broad indemnification, no-damage-for-delay",
    "Compares payment terms against your standards",
    "Checks insurance requirements match your minimums",
    "Identifies missing clauses (retainage, dispute resolution, lien rights)",
    "Generates a risk score and plain-English summary for the PM or CFO",
]
y = 1.5
for item in contract_features:
    add_text(slide, Inches(7.0), Inches(y), Inches(5.8), Inches(0.35),
             "  " + item, 11, RGBColor(0xCC, 0xCC, 0xCC))
    add_shape(slide, Inches(7.0), Inches(y + 0.05), Inches(0.12), Inches(0.12), GREEN)
    y += 0.37

# Revenue model
add_shape(slide, Inches(0.3), Inches(4.4), Inches(6.2), Inches(2.8), DARK_BLUE, GOLD)
add_text(slide, Inches(0.3), Inches(4.45), Inches(6.2), Inches(0.4),
         "Revenue Model", 18, GOLD, True, PP_ALIGN.CENTER)
rev_items = [
    ("Per contract review:", "$200-500 per contract"),
    ("Monthly subscription:", "$2,000-10,000/month"),
    ("Enterprise (unlimited):", "$50,000-100,000/year"),
    ("Our AI cost per contract:", "~$0.15-0.30 (Opus for complex legal)"),
    ("Gross margin:", "95%+"),
    ("Target market:", "700K construction firms + 50K law firms"),
]
y = 4.95
for label, value in rev_items:
    add_text(slide, Inches(0.5), Inches(y), Inches(2.8), Inches(0.3),
             label, 12, MEDIUM_GRAY, True)
    add_text(slide, Inches(3.3), Inches(y), Inches(3.0), Inches(0.3),
             value, 12, GREEN if "95" in value or "$" in value else WHITE)
    y += 0.35

# Strategic value
add_shape(slide, Inches(6.8), Inches(4.4), Inches(6.2), Inches(2.8), DARK_BLUE, ACCENT_BLUE)
add_text(slide, Inches(6.8), Inches(4.45), Inches(6.2), Inches(0.4),
         "Strategic Value", 18, ACCENT_BLUE, True, PP_ALIGN.CENTER)
strat_items = [
    "Same customers as PayFlow AI — construction companies",
    "Bundle with PayFlow: AP automation + contract review = higher ARPU",
    "Reuses the same Claude AI engine and document processing pipeline",
    "Cross-sell: every PayFlow customer is a contract review prospect",
    "Defensive moat: competitors can't easily replicate construction expertise",
    "Data flywheel: every contract reviewed improves the AI's clause library",
]
y = 4.95
for item in strat_items:
    add_text(slide, Inches(7.0), Inches(y), Inches(5.8), Inches(0.32),
             "  " + item, 11, RGBColor(0xCC, 0xCC, 0xCC))
    add_shape(slide, Inches(7.0), Inches(y + 0.05), Inches(0.12), Inches(0.12), ACCENT_BLUE)
    y += 0.35


# ═══════════════════════════════════════════════
# SLIDE 6: Product 3 — AI Permit Expediting
# ═══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.06), TEAL)

add_shape(slide, Inches(0.5), Inches(0.2), Inches(0.5), Inches(0.5), TEAL)
add_text(slide, Inches(0.5), Inches(0.22), Inches(0.5), Inches(0.45),
         "3", 24, WHITE, True, PP_ALIGN.CENTER)
add_text(slide, Inches(1.2), Inches(0.25), Inches(8), Inches(0.6),
         "AI Permit Expediting — Automated Permit Applications", 30, WHITE, True)

# Problem + Solution
add_shape(slide, Inches(0.3), Inches(1.0), Inches(6.2), Inches(2.5), DARK_BLUE, RED)
add_text(slide, Inches(0.3), Inches(1.05), Inches(6.2), Inches(0.4),
         "The Problem", 18, RED, True, PP_ALIGN.CENTER)
permit_problems = [
    "Permit applications take 2-4 weeks to prepare manually",
    "Permit expediters charge $5,000-15,000 per project",
    "Applications get rejected 30-40% of the time for errors or missing info",
    "Each rejection delays the project 2-6 weeks",
    "1M+ building permits issued annually in the US",
    "No one is automating this with AI yet — massive first-mover opportunity",
]
y = 1.5
for item in permit_problems:
    add_text(slide, Inches(0.5), Inches(y), Inches(5.8), Inches(0.32),
             "  " + item, 11, RGBColor(0xCC, 0xCC, 0xCC))
    add_shape(slide, Inches(0.5), Inches(y + 0.05), Inches(0.12), Inches(0.12), RED)
    y += 0.35

add_shape(slide, Inches(6.8), Inches(1.0), Inches(6.2), Inches(2.5), DARK_BLUE, GREEN)
add_text(slide, Inches(6.8), Inches(1.05), Inches(6.2), Inches(0.4),
         "What AI Permit Expediting Does", 18, GREEN, True, PP_ALIGN.CENTER)
permit_features = [
    "Upload construction drawings (blueprints, plans) — Claude reads them visually",
    "AI auto-fills permit application forms for the specific jurisdiction",
    "Cross-checks plans against local building codes (IBC, IRC, local amendments)",
    "Flags code violations BEFORE submission — reduces rejections by 60-80%",
    "Generates required supporting documents and calculations",
    "Tracks submission status and response deadlines",
]
y = 1.5
for item in permit_features:
    add_text(slide, Inches(7.0), Inches(y), Inches(5.8), Inches(0.32),
             "  " + item, 11, RGBColor(0xCC, 0xCC, 0xCC))
    add_shape(slide, Inches(7.0), Inches(y + 0.05), Inches(0.12), Inches(0.12), GREEN)
    y += 0.35

# Revenue model
add_shape(slide, Inches(0.3), Inches(3.7), Inches(6.2), Inches(2.0), DARK_BLUE, GOLD)
add_text(slide, Inches(0.3), Inches(3.75), Inches(6.2), Inches(0.4),
         "Revenue Model", 18, GOLD, True, PP_ALIGN.CENTER)
rev = [
    ("Per permit:", "$500-2,000 per submission"),
    ("Subscription:", "$3,000-15,000/month for high-volume firms"),
    ("Market size:", "1M+ permits/year x $500 avg = $500M+ market"),
    ("Our AI cost:", "~$0.50-2.00 per permit (blueprint reading is token-heavy)"),
    ("Gross margin:", "90%+"),
]
y = 4.2
for label, value in rev:
    add_text(slide, Inches(0.5), Inches(y), Inches(2.5), Inches(0.3),
             label, 12, MEDIUM_GRAY, True)
    add_text(slide, Inches(3.0), Inches(y), Inches(3.3), Inches(0.3),
             value, 12, GREEN if "90" in value or "$500M" in value else WHITE)
    y += 0.32

# Why now
add_shape(slide, Inches(6.8), Inches(3.7), Inches(6.2), Inches(2.0), DARK_BLUE, TEAL)
add_text(slide, Inches(6.8), Inches(3.75), Inches(6.2), Inches(0.4),
         "Why This Is a Massive Opportunity", 18, TEAL, True, PP_ALIGN.CENTER)
why = [
    "Claude is one of the first AI models that can read blueprints accurately",
    "No competitor exists — pure greenfield market",
    "Every project needs permits — residential, commercial, industrial",
    "Municipalities are moving to digital submissions (ePermitting)",
    "Same construction customers as PayFlow AI — sell as a suite",
]
y = 4.2
for item in why:
    add_text(slide, Inches(7.0), Inches(y), Inches(5.8), Inches(0.32),
             "  " + item, 11, RGBColor(0xCC, 0xCC, 0xCC))
    add_shape(slide, Inches(7.0), Inches(y + 0.05), Inches(0.12), Inches(0.12), TEAL)
    y += 0.32

# Bottom
add_shape(slide, Inches(0.3), Inches(5.9), Inches(12.7), Inches(0.4), DARK_BLUE, GOLD)
add_text(slide, Inches(0.3), Inches(5.92), Inches(12.7), Inches(0.35),
         "First-Mover Advantage: No AI-powered permit expediting tool exists in the market today",
         14, GOLD, True, PP_ALIGN.CENTER)

# How it works flow
add_text(slide, Inches(0.5), Inches(6.5), Inches(12), Inches(0.4),
         "How It Works:", 16, WHITE, True)

flow_boxes = [
    ("Upload Plans", "Blueprints, specs,\nsite plans", ACCENT_BLUE),
    ("AI Reads Plans", "Claude extracts\nscope, dimensions,\nsystems", TEAL),
    ("Code Check", "Cross-reference\nlocal building codes\nand zoning", ORANGE),
    ("Auto-Fill App", "Generate permit\napplication with\nall required fields", GREEN),
    ("Flag Issues", "Identify violations\nbefore submission\nsaves weeks", RED),
    ("Submit", "Digital submission\nto municipality\ntrack status", DARK_GREEN),
]
x = 0.4
for title, desc, color in flow_boxes:
    add_shape(slide, Inches(x), Inches(6.85), Inches(1.9), Inches(0.55), color)
    add_text(slide, Inches(x + 0.05), Inches(6.87), Inches(1.8), Inches(0.25),
             title, 11, WHITE, True, PP_ALIGN.CENTER)
    x += 2.1


# ═══════════════════════════════════════════════
# SLIDE 7: Product 4 — AI Insurance Claims
# ═══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.06), PURPLE)

add_shape(slide, Inches(0.5), Inches(0.2), Inches(0.5), Inches(0.5), PURPLE)
add_text(slide, Inches(0.5), Inches(0.22), Inches(0.5), Inches(0.45),
         "4", 24, WHITE, True, PP_ALIGN.CENTER)
add_text(slide, Inches(1.2), Inches(0.25), Inches(8), Inches(0.6),
         "AI Insurance Claims — Processing & Fraud Detection", 30, WHITE, True)

# Left
add_shape(slide, Inches(0.3), Inches(1.0), Inches(6.2), Inches(3.0), DARK_BLUE, RED)
add_text(slide, Inches(0.3), Inches(1.05), Inches(6.2), Inches(0.4),
         "The $80 Billion Problem", 18, RED, True, PP_ALIGN.CENTER)
ins_problems = [
    "300M+ insurance claims processed annually in the US",
    "Each claim costs $30-50 in manual labor to process",
    "Average claim takes 10-15 days to adjudicate",
    "Insurance fraud costs the industry $80B/year",
    "Claims adjusters manually read documents, cross-reference policies, calculate coverage",
    "Same document processing pattern as AP automation — but bigger market",
]
y = 1.55
for item in ins_problems:
    add_text(slide, Inches(0.5), Inches(y), Inches(5.8), Inches(0.35),
             "  " + item, 11, RGBColor(0xCC, 0xCC, 0xCC))
    add_shape(slide, Inches(0.5), Inches(y + 0.05), Inches(0.12), Inches(0.12), RED)
    y += 0.38

# Right
add_shape(slide, Inches(6.8), Inches(1.0), Inches(6.2), Inches(3.0), DARK_BLUE, GREEN)
add_text(slide, Inches(6.8), Inches(1.05), Inches(6.2), Inches(0.4),
         "What Our AI Does", 18, GREEN, True, PP_ALIGN.CENTER)
ins_features = [
    "Reads claim documents (photos, PDFs, medical records, police reports)",
    "Extracts all claim data into structured format",
    "Cross-references against policy terms and coverage limits",
    "Calculates covered amount vs deductible vs exclusions",
    "Fraud detection: duplicate claims, inflated amounts, suspicious patterns",
    "Routes to the right adjuster with AI recommendation and risk score",
]
y = 1.55
for item in ins_features:
    add_text(slide, Inches(7.0), Inches(y), Inches(5.8), Inches(0.35),
             "  " + item, 11, RGBColor(0xCC, 0xCC, 0xCC))
    add_shape(slide, Inches(7.0), Inches(y + 0.05), Inches(0.12), Inches(0.12), GREEN)
    y += 0.38

# Bottom: Revenue
add_shape(slide, Inches(0.3), Inches(4.2), Inches(4.0), Inches(3.0), DARK_BLUE, GOLD)
add_text(slide, Inches(0.3), Inches(4.25), Inches(4.0), Inches(0.4),
         "Revenue Potential", 18, GOLD, True, PP_ALIGN.CENTER)
ins_rev = [
    ("Per claim:", "$3-10"),
    ("Monthly (mid-size carrier):", "$5,000-50,000"),
    ("Monthly (top-20 carrier):", "$100,000-500,000"),
    ("Market (US claims x $5):", "$1.5B+"),
    ("Our AI cost per claim:", "~$0.05-0.15"),
    ("Gross margin:", "95%+"),
]
y = 4.7
for label, value in ins_rev:
    add_text(slide, Inches(0.5), Inches(y), Inches(2.2), Inches(0.3),
             label, 11, MEDIUM_GRAY, True)
    add_text(slide, Inches(2.7), Inches(y), Inches(1.5), Inches(0.3),
             value, 11, GREEN, True)
    y += 0.33

# Considerations
add_shape(slide, Inches(4.5), Inches(4.2), Inches(4.1), Inches(3.0), DARK_BLUE, ORANGE)
add_text(slide, Inches(4.5), Inches(4.25), Inches(4.1), Inches(0.4),
         "Considerations", 18, ORANGE, True, PP_ALIGN.CENTER)
cons = [
    "Regulated industry — need compliance expertise",
    "Long sales cycles with insurance carriers",
    "Existing players: Shift Technology, FRISS",
    "Need domain expertise in insurance",
    "Higher barrier but much larger market",
    "Could start with construction insurance niche",
]
y = 4.7
for item in cons:
    add_text(slide, Inches(4.7), Inches(y), Inches(3.7), Inches(0.3),
             "  " + item, 11, RGBColor(0xCC, 0xCC, 0xCC))
    add_shape(slide, Inches(4.7), Inches(y + 0.05), Inches(0.12), Inches(0.12), ORANGE)
    y += 0.33

# Opportunity
add_shape(slide, Inches(8.8), Inches(4.2), Inches(4.2), Inches(3.0), DARK_BLUE, ACCENT_BLUE)
add_text(slide, Inches(8.8), Inches(4.25), Inches(4.2), Inches(0.4),
         "Construction Insurance Angle", 16, ACCENT_BLUE, True, PP_ALIGN.CENTER)
angle = [
    "Start with workers comp claims for construction",
    "Same customers as PayFlow AI",
    "Cross-reference with project data from Procore",
    "Verify claims against daily logs and job records",
    "Unique positioning: construction + insurance + AI",
    "Niche entry → expand to broader insurance market",
]
y = 4.7
for item in angle:
    add_text(slide, Inches(9.0), Inches(y), Inches(3.8), Inches(0.3),
             "  " + item, 11, RGBColor(0xCC, 0xCC, 0xCC))
    add_shape(slide, Inches(9.0), Inches(y + 0.05), Inches(0.12), Inches(0.12), ACCENT_BLUE)
    y += 0.33


# ═══════════════════════════════════════════════
# SLIDE 8: Products 5-8 (Quick View)
# ═══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.06), GOLD)

add_text(slide, Inches(0.8), Inches(0.25), Inches(11), Inches(0.6),
         "Additional Opportunities (Products 5-8)", 30, WHITE, True)

products_5_8 = [
    ("5", "AI Medical Coding", "ICD-10 / CPT Code Prediction", PURPLE,
     ["Reads clinical documentation, predicts correct diagnosis and procedure codes",
      "Average practice loses 5-10% of revenue to coding errors — AI fixes this",
      "1M+ healthcare providers, $5-30K/month per practice",
      "HIPAA compliance required — higher barrier but massive market ($4.3T healthcare)",
      "Same pattern: document in → AI reads → structured output → human reviews"]),
    ("6", "AI Lien & Notice Management", "Deadline Tracking & Auto-Filing", GREEN,
     ["Tracks preliminary notice deadlines by state — miss one and you lose lien rights",
      "Auto-generates and sends notices before deadlines",
      "Manages waiver exchange between GCs, subs, and suppliers",
      "$500-2K/month per customer — direct add-on to PayFlow AI",
      "Same customers, same data — highest synergy of any product here"]),
    ("7", "AI Estimating & Proposals", "Automated Bid Generation", TEAL,
     ["Takes project specs or RFPs, generates detailed cost estimates",
      "Uses historical bid data to improve accuracy over time",
      "Win rates are 10-20% — AI helps bid smarter, not just faster",
      "$1-5K/month per contractor, 700K+ construction firms",
      "Medium complexity — needs material pricing databases and labor rate data"]),
    ("8", "AI Property Management AP", "Multi-Property Invoice Processing", ACCENT_BLUE,
     ["Same AP automation engine but for property management",
      "Process maintenance invoices, utility bills, vendor payments across hundreds of properties",
      "Integrates with Yardi, AppFolio, Buildium instead of Sage/Procore",
      "300K+ property management companies in the US, $1-5K/month",
      "Fastest to build — 70% of PayFlow code reused with different integrations"]),
]

y = 0.95
for num, name, subtitle, color, items in products_5_8:
    add_shape(slide, Inches(0.3), Inches(y), Inches(12.7), Inches(1.45), DARK_BLUE, color)
    add_shape(slide, Inches(0.4), Inches(y + 0.1), Inches(0.4), Inches(0.4), color)
    add_text(slide, Inches(0.4), Inches(y + 0.12), Inches(0.4), Inches(0.35),
             num, 18, WHITE, True, PP_ALIGN.CENTER)
    add_text(slide, Inches(0.9), Inches(y + 0.08), Inches(3.0), Inches(0.35),
             name, 18, color, True)
    add_text(slide, Inches(3.9), Inches(y + 0.12), Inches(3.5), Inches(0.3),
             subtitle, 12, MEDIUM_GRAY)
    item_y = y + 0.5
    for item in items:
        add_text(slide, Inches(0.7), Inches(item_y), Inches(12.0), Inches(0.25),
                 "  " + item, 10, RGBColor(0xBB, 0xBB, 0xBB))
        add_shape(slide, Inches(0.7), Inches(item_y + 0.03), Inches(0.1), Inches(0.1), color)
        item_y += 0.2
    y += 1.55


# ═══════════════════════════════════════════════
# SLIDE 9: Revenue Projection Chart
# ═══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.06), GOLD)

add_text(slide, Inches(0.8), Inches(0.25), Inches(11), Inches(0.6),
         "Combined Revenue Projection — Product Portfolio", 30, WHITE, True)
add_text(slide, Inches(0.8), Inches(0.8), Inches(11), Inches(0.4),
         "Staggered launch: each product builds on the same AI engine and customer base", 15, MEDIUM_GRAY)

# Chart
chart_data = CategoryChartData()
chart_data.categories = ["Year 1", "Year 2", "Year 3", "Year 4", "Year 5"]
chart_data.add_series("PayFlow AI", (90, 600, 3000, 8000, 18000))
chart_data.add_series("Contract Review", (0, 200, 1200, 3500, 8000))
chart_data.add_series("Permit Expediting", (0, 0, 500, 2000, 6000))
chart_data.add_series("Other Products", (0, 0, 300, 1500, 5000))

chart = slide.shapes.add_chart(
    XL_CHART_TYPE.COLUMN_STACKED,
    Inches(0.5), Inches(1.3), Inches(8.0), Inches(5.5),
    chart_data
).chart

chart.has_legend = True
chart.legend.position = XL_LEGEND_POSITION.BOTTOM
chart.legend.font.size = Pt(11)
chart.legend.font.color.rgb = WHITE
chart.legend.include_in_layout = False

plot = chart.plots[0]
plot.gap_width = 80

colors = [ACCENT_BLUE, GREEN, TEAL, PURPLE]
for i, color in enumerate(colors):
    series = plot.series[i]
    series.format.fill.solid()
    series.format.fill.fore_color.rgb = color

chart.category_axis.tick_labels.font.size = Pt(12)
chart.category_axis.tick_labels.font.color.rgb = WHITE
chart.value_axis.tick_labels.font.size = Pt(10)
chart.value_axis.tick_labels.font.color.rgb = MEDIUM_GRAY
chart.value_axis.format.line.color.rgb = RGBColor(0x30, 0x40, 0x55)
chart.category_axis.format.line.color.rgb = RGBColor(0x30, 0x40, 0x55)
chart.value_axis.major_gridlines.format.line.color.rgb = RGBColor(0x20, 0x30, 0x45)

# Value axis title
chart.value_axis.has_title = True
chart.value_axis.axis_title.text_frame.paragraphs[0].text = "Annual Revenue ($K)"
chart.value_axis.axis_title.text_frame.paragraphs[0].font.size = Pt(10)
chart.value_axis.axis_title.text_frame.paragraphs[0].font.color.rgb = MEDIUM_GRAY

# Right side: summary table
add_shape(slide, Inches(8.8), Inches(1.3), Inches(4.2), Inches(5.5), DARK_BLUE, GOLD)
add_text(slide, Inches(8.8), Inches(1.35), Inches(4.2), Inches(0.4),
         "Portfolio Revenue Summary", 16, GOLD, True, PP_ALIGN.CENTER)

summary = [
    ("Year 1", "$90K", "PayFlow only"),
    ("Year 2", "$800K", "+ Contract Review"),
    ("Year 3", "$5.0M", "+ Permit Expediting"),
    ("Year 4", "$15.0M", "+ Other products"),
    ("Year 5", "$37.0M", "Full portfolio"),
]

add_shape(slide, Inches(9.0), Inches(1.85), Inches(3.8), Inches(0.35), RGBColor(0x15, 0x2D, 0x52))
add_text(slide, Inches(9.1), Inches(1.87), Inches(1.0), Inches(0.3), "Year", 11, GOLD, True)
add_text(slide, Inches(10.1), Inches(1.87), Inches(1.0), Inches(0.3), "ARR", 11, GOLD, True)
add_text(slide, Inches(11.1), Inches(1.87), Inches(1.5), Inches(0.3), "Products Live", 11, GOLD, True)

y = 2.25
for year, arr, products_live in summary:
    add_text(slide, Inches(9.1), Inches(y), Inches(1.0), Inches(0.35),
             year, 13, WHITE, True)
    add_text(slide, Inches(10.1), Inches(y), Inches(1.0), Inches(0.35),
             arr, 14, GREEN, True)
    add_text(slide, Inches(11.1), Inches(y), Inches(1.5), Inches(0.35),
             products_live, 11, MEDIUM_GRAY)
    y += 0.45

# Assumptions
add_text(slide, Inches(9.0), Inches(4.6), Inches(3.8), Inches(0.3),
         "Key Assumptions:", 12, ACCENT_BLUE, True)
assumptions = [
    "Conservative customer acquisition",
    "Average $3K/mo per customer",
    "Each product launched 12 months apart",
    "Same sales team sells all products",
    "95%+ gross margins on all products",
    "No external funding required",
    "Bootstrapped from PayFlow revenue",
]
y = 4.9
for item in assumptions:
    add_text(slide, Inches(9.2), Inches(y), Inches(3.5), Inches(0.25),
             "  " + item, 9, RGBColor(0xBB, 0xBB, 0xBB))
    add_shape(slide, Inches(9.2), Inches(y + 0.03), Inches(0.08), Inches(0.08), ACCENT_BLUE)
    y += 0.22


# ═══════════════════════════════════════════════
# SLIDE 10: Shared AI Engine
# ═══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.06), ACCENT_BLUE)

add_text(slide, Inches(0.8), Inches(0.25), Inches(11), Inches(0.6),
         "The Unfair Advantage: One AI Engine, Many Products", 30, WHITE, True)
add_text(slide, Inches(0.8), Inches(0.8), Inches(11), Inches(0.4),
         "Every product reuses the same core technology — reducing development time and cost by 60-70%", 15, MEDIUM_GRAY)

# Center: AI Engine
add_shape(slide, Inches(4.2), Inches(2.5), Inches(5.0), Inches(2.5), DARK_BLUE, ACCENT_BLUE)
add_text(slide, Inches(4.2), Inches(2.6), Inches(5.0), Inches(0.5),
         "Core AI Engine", 24, ACCENT_BLUE, True, PP_ALIGN.CENTER)
engine_items = [
    "Document ingestion (email, upload, scan, API)",
    "Claude AI extraction (multimodal vision + structured output)",
    "Confidence scoring & validation",
    "Approval routing & workflow",
    "User dashboard & reporting",
    "Multi-tenant SaaS infrastructure (AWS)",
    "Authentication, security, audit trails",
]
y = 3.15
for item in engine_items:
    add_text(slide, Inches(4.5), Inches(y), Inches(4.5), Inches(0.25),
             "  " + item, 10, RGBColor(0xCC, 0xCC, 0xCC))
    add_shape(slide, Inches(4.5), Inches(y + 0.03), Inches(0.08), Inches(0.08), ACCENT_BLUE)
    y += 0.22

# Products around the engine
products_around = [
    ("PayFlow AI", "AP invoices", ACCENT_BLUE, 0.5, 1.5),
    ("Contract Review", "Legal docs", GREEN, 0.5, 3.5),
    ("Permit Expediting", "Blueprints", TEAL, 0.5, 5.5),
    ("Insurance Claims", "Claim docs", PURPLE, 9.8, 1.5),
    ("Medical Coding", "Clinical notes", ORANGE, 9.8, 3.5),
    ("Property Mgmt AP", "Maint invoices", RGBColor(0x2C, 0xA0, 0x1C), 9.8, 5.5),
]

for name, docs, color, x, y_pos in products_around:
    add_shape(slide, Inches(x), Inches(y_pos), Inches(3.2), Inches(1.2), DARK_BLUE, color)
    add_text(slide, Inches(x), Inches(y_pos + 0.1), Inches(3.2), Inches(0.35),
             name, 14, color, True, PP_ALIGN.CENTER)
    add_text(slide, Inches(x), Inches(y_pos + 0.5), Inches(3.2), Inches(0.3),
             "Different prompts", 11, MEDIUM_GRAY, False, PP_ALIGN.CENTER)
    add_text(slide, Inches(x), Inches(y_pos + 0.75), Inches(3.2), Inches(0.3),
             "Different integrations", 11, MEDIUM_GRAY, False, PP_ALIGN.CENTER)

# Arrows pointing to center
for y_pos in [1.9, 3.9, 5.9]:
    add_arrow(slide, Inches(3.7), Inches(y_pos), Inches(0.4), Inches(0.25), RGBColor(0x30, 0x50, 0x70))
    shape = slide.shapes.add_shape(MSO_SHAPE.LEFT_ARROW, Inches(9.3), Inches(y_pos), Inches(0.4), Inches(0.25))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0x30, 0x50, 0x70)
    shape.line.fill.background()

# Bottom
add_shape(slide, Inches(0.3), Inches(7.0), Inches(12.7), Inches(0.35), DARK_BLUE, ACCENT_BLUE)
add_text(slide, Inches(0.3), Inches(7.02), Inches(12.7), Inches(0.3),
         "Build the engine once. Change the prompts and integrations. Launch a new product in weeks, not months.",
         13, ACCENT_BLUE, True, PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════
# SLIDE 11: Competitive Moats
# ═══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.06), GOLD)

add_text(slide, Inches(0.8), Inches(0.25), Inches(11), Inches(0.6),
         "Competitive Moats — Why Competitors Can't Catch Up", 30, WHITE, True)

moats = [
    ("Vertical Expertise", "Deep domain knowledge in construction AP, contracts, permitting, and compliance. Generic AI tools can't match this specificity. We know the difference between AIA G702 and a standard invoice — Stampli doesn't.",
     ACCENT_BLUE),
    ("Data Flywheel", "Every invoice processed, every correction made, every contract reviewed makes the AI smarter. By the time a competitor starts, we have millions of data points they don't. Accuracy compounds over time.",
     GREEN),
    ("Integration Lock-In", "Connected to KOJO + Procore + Sage + 10 other tools. Ripping out PayFlow means disconnecting everything. Switching cost is months of re-integration. Customers don't switch once embedded.",
     TEAL),
    ("Multi-Product Platform", "One login, one vendor, multiple products. A CFO buys PayFlow for AP, then adds Contract Review, then Permit Expediting. Each product makes the others stickier. Competitors selling one tool can't compete with a suite.",
     GOLD),
    ("Cost Advantage", "Claude AI processes at $0.06/invoice while incumbents built their own ML models at 10x the development cost. We're API-native — no $50M investment in training infrastructure. Faster to adapt when AI models improve.",
     PURPLE),
    ("First-Mover in Construction AI", "No one has built a comprehensive AI platform for construction back-office. There are point solutions (Stampli for AP, LevelSet for liens) but no unified construction AI platform. We're building the Salesforce of construction finance.",
     RED),
]

y = 1.0
for title, desc, color in moats:
    add_shape(slide, Inches(0.3), Inches(y), Inches(12.7), Inches(0.95), DARK_BLUE, color)
    add_text(slide, Inches(0.5), Inches(y + 0.05), Inches(2.5), Inches(0.3),
             title, 15, color, True)
    add_text(slide, Inches(3.1), Inches(y + 0.05), Inches(9.6), Inches(0.85),
             desc, 11, RGBColor(0xCC, 0xCC, 0xCC))
    y += 1.02


# ═══════════════════════════════════════════════
# SLIDE 12: Execution Roadmap
# ═══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.06), GOLD)

add_text(slide, Inches(0.8), Inches(0.25), Inches(11), Inches(0.6),
         "Execution Roadmap — From First Product to Platform", 30, WHITE, True)

phases = [
    ("NOW", "Q1-Q2 2026", "Build PayFlow AI", [
        "Complete development (6-month timeline)",
        "Pilot with Mark III Construction",
        "Shadow mode → live mode → paying customer",
        "List on Procore App Marketplace",
        "Revenue: $1.5K/month (1 customer)",
    ], GREEN, "Month 1-6"),
    ("NEXT", "Q3-Q4 2026", "Scale PayFlow + Start Contract Review", [
        "Land 5 paying customers from Procore marketplace",
        "Begin building AI Contract Review (reuse 70% of engine)",
        "Bundle pricing: PayFlow + Contract Review",
        "Revenue: $7.5K-15K/month",
    ], ACCENT_BLUE, "Month 7-12"),
    ("THEN", "2027", "Launch Contract Review + Start Permits", [
        "Contract Review live, sell to existing PayFlow customers",
        "Begin Permit Expediting development",
        "Hire first salesperson",
        "Revenue: $50K+/month, 25+ customers",
    ], TEAL, "Month 13-24"),
    ("SCALE", "2028+", "Full Platform + New Verticals", [
        "Launch Permit Expediting",
        "Expand to Lien Management, Property Mgmt AP",
        "Consider insurance claims or medical coding",
        "Revenue: $250K+/month, 100+ customers",
    ], GOLD, "Month 25-36"),
]

y = 1.1
for stage, timeframe, title, items, color, months in phases:
    add_shape(slide, Inches(0.3), Inches(y), Inches(12.7), Inches(1.4), DARK_BLUE, color)
    add_shape(slide, Inches(0.4), Inches(y + 0.08), Inches(1.0), Inches(0.35), color)
    add_text(slide, Inches(0.4), Inches(y + 0.1), Inches(1.0), Inches(0.3),
             stage, 14, WHITE, True, PP_ALIGN.CENTER)
    add_text(slide, Inches(1.6), Inches(y + 0.05), Inches(1.5), Inches(0.3),
             timeframe, 12, MEDIUM_GRAY, True)
    add_text(slide, Inches(3.2), Inches(y + 0.05), Inches(4.0), Inches(0.3),
             title, 16, WHITE, True)
    add_text(slide, Inches(10.5), Inches(y + 0.05), Inches(2.3), Inches(0.3),
             months, 12, color, True, PP_ALIGN.RIGHT)
    item_y = y + 0.42
    for item in items:
        add_text(slide, Inches(1.6), Inches(item_y), Inches(10.0), Inches(0.22),
                 "  " + item, 10, RGBColor(0xBB, 0xBB, 0xBB))
        add_shape(slide, Inches(1.6), Inches(item_y + 0.03), Inches(0.08), Inches(0.08), color)
        item_y += 0.2
    y += 1.48

# Bottom: Key principle
add_shape(slide, Inches(0.3), Inches(7.05), Inches(12.7), Inches(0.35), DARK_BLUE, GOLD)
add_text(slide, Inches(0.3), Inches(7.07), Inches(12.7), Inches(0.3),
         "Principle: Each product funds the development of the next. No external funding required. Bootstrap from revenue.",
         13, GOLD, True, PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════
# SLIDE 13: Investment Required
# ═══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.06), GOLD)

add_text(slide, Inches(0.8), Inches(0.25), Inches(11), Inches(0.6),
         "Investment Required vs. Return", 30, WHITE, True)

# Left: Investment
add_shape(slide, Inches(0.3), Inches(1.0), Inches(6.2), Inches(5.8), DARK_BLUE, ORANGE)
add_text(slide, Inches(0.3), Inches(1.05), Inches(6.2), Inches(0.5),
         "What It Takes", 22, ORANGE, True, PP_ALIGN.CENTER)

invest_items = [
    ("PayFlow AI Development", "6 months, 1 developer + Claude Code", "$0 (sweat equity)"),
    ("AWS Infrastructure (Year 1)", "Servers, database, storage, security", "$5,000-8,000"),
    ("Claude API (Year 1)", "Processing invoices during dev + pilot", "$1,000-2,000"),
    ("Legal Setup", "LLC, trademark, terms of service", "$2,000-5,000"),
    ("Domain & Branding", "payflowai.com, logo, marketing site", "$500-1,000"),
    ("Business Insurance", "E&O + cyber liability", "$1,200-2,400"),
    ("Software Tools", "GitHub, monitoring, email, analytics", "$500-1,000"),
    ("SOC 2 Prep (end of Year 1)", "Compliance readiness", "$5,000-10,000"),
]

y = 1.65
total_low = 0
total_high = 0
for item, desc, cost in invest_items:
    add_text(slide, Inches(0.5), Inches(y), Inches(2.5), Inches(0.28),
             item, 11, WHITE, True)
    add_text(slide, Inches(0.5), Inches(y + 0.28), Inches(2.5), Inches(0.22),
             desc, 9, MEDIUM_GRAY)
    add_text(slide, Inches(4.5), Inches(y + 0.05), Inches(1.8), Inches(0.3),
             cost, 11, ORANGE if "sweat" in cost else GREEN, True, PP_ALIGN.RIGHT)
    y += 0.55

# Total
add_shape(slide, Inches(0.5), Inches(y + 0.1), Inches(5.8), Inches(0.45), ORANGE)
add_text(slide, Inches(0.6), Inches(y + 0.13), Inches(3.5), Inches(0.35),
         "Total Year 1 Cash Investment:", 13, WHITE, True)
add_text(slide, Inches(4.3), Inches(y + 0.13), Inches(1.8), Inches(0.35),
         "$15,000-30,000", 14, WHITE, True, PP_ALIGN.RIGHT)

# Right: Return
add_shape(slide, Inches(6.8), Inches(1.0), Inches(6.2), Inches(5.8), DARK_BLUE, GREEN)
add_text(slide, Inches(6.8), Inches(1.05), Inches(6.2), Inches(0.5),
         "What You Get Back", 22, GREEN, True, PP_ALIGN.CENTER)

return_items = [
    ("Year 1 Revenue (5 customers)", "$90,000", GREEN),
    ("Year 2 Revenue (25 customers)", "$600,000", GREEN),
    ("Year 3 Revenue (100 customers)", "$3,000,000", GREEN),
    ("Year 5 Revenue (full portfolio)", "$37,000,000", GOLD),
]

y = 1.8
for item, value, color in return_items:
    add_shape(slide, Inches(7.0), Inches(y), Inches(5.8), Inches(0.7), RGBColor(0x18, 0x30, 0x55))
    add_text(slide, Inches(7.2), Inches(y + 0.08), Inches(3.3), Inches(0.3),
             item, 13, WHITE)
    add_text(slide, Inches(10.5), Inches(y + 0.05), Inches(2.1), Inches(0.35),
             value, 20, color, True, PP_ALIGN.RIGHT)
    y += 0.8

add_text(slide, Inches(7.0), Inches(5.1), Inches(5.8), Inches(0.4),
         "Additional Value:", 14, ACCENT_BLUE, True)
extra = [
    "IP portfolio worth multiples of revenue at exit",
    "Platform play valued at 10-20x ARR by acquirers",
    "At $37M ARR: company valued at $370M-740M",
    "Potential acquirers: Procore, Sage, Oracle, Autodesk",
    "Or: continue growing independently with 95% margins",
]
y = 5.5
for item in extra:
    add_text(slide, Inches(7.2), Inches(y), Inches(5.5), Inches(0.25),
             "  " + item, 11, RGBColor(0xCC, 0xCC, 0xCC))
    add_shape(slide, Inches(7.2), Inches(y + 0.03), Inches(0.08), Inches(0.08), ACCENT_BLUE)
    y += 0.25


# ═══════════════════════════════════════════════
# SLIDE 14: The Ask / Closing
# ═══════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.06), GOLD)

add_text(slide, Inches(1), Inches(1.0), Inches(11), Inches(1.0),
         "The Opportunity Is Clear", 44, WHITE, True, PP_ALIGN.CENTER)
add_text(slide, Inches(1.5), Inches(2.2), Inches(10), Inches(0.8),
         "AI is creating a once-in-a-generation opportunity to build vertical SaaS products\nthat replace expensive manual processes at 95%+ margins.", 20, MEDIUM_GRAY, False, PP_ALIGN.CENTER)

# Three key stats
key_stats = [
    ("$15-30K", "Year 1\nInvestment", ORANGE),
    ("$37M ARR", "Year 5\nRevenue Target", GREEN),
    ("1,200x", "Potential\nReturn", GOLD),
]
x = 2.2
for value, label, color in key_stats:
    add_shape(slide, Inches(x), Inches(3.5), Inches(2.8), Inches(1.8), DARK_BLUE, color)
    add_text(slide, Inches(x), Inches(3.65), Inches(2.8), Inches(0.7),
             value, 36, color, True, PP_ALIGN.CENTER)
    add_text(slide, Inches(x + 0.2), Inches(4.3), Inches(2.4), Inches(0.8),
             label, 16, RGBColor(0xCC, 0xCC, 0xCC), False, PP_ALIGN.CENTER)
    x += 3.15

# Next steps
add_shape(slide, Inches(2.5), Inches(5.6), Inches(8.3), Inches(1.5), DARK_BLUE, GOLD)
add_text(slide, Inches(2.5), Inches(5.65), Inches(8.3), Inches(0.4),
         "Recommended Next Step", 18, GOLD, True, PP_ALIGN.CENTER)
add_text(slide, Inches(2.8), Inches(6.05), Inches(7.7), Inches(0.9),
         "Build PayFlow AI first. Prove the model with one paying customer.\nThen expand into Contract Review and Permit Expediting.\nEach product funds the next. No external capital required.",
         15, WHITE, False, PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════
# SAVE
# ═══════════════════════════════════════════════
output = "C:/users/25badmin/projects/accounts-payable-research/AI_SaaS_Opportunities_CEO.pptx"
prs.save(output)
print(f"Saved: {output}")
